from __future__ import annotations

import io
import re
from dataclasses import dataclass
from typing import Any

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


BEAMER_SYSTEM_PROMPT = """你是一个专业的 LaTeX Beamer 演示文稿生成专家。你的任务是将用户提供的文案内容转化为高质量的 LaTeX Beamer 代码，生成风格统一、结构清晰的学术演示文稿。

## 输出格式要求
1. 必须输出完整的、可直接用 xelatex 编译的 .tex 文件
2. 从 \\documentclass 开始，到 \\end{document} 结束
3. 不要输出 markdown 代码块标记
4. 不要输出任何解释文字、注释说明或无关提示
5. 只输出纯 LaTeX 代码

## 结构组织
1. 自动提取 3-6 个核心主题作为 \\section
2. 第一页为标题页，第二页为总览/回顾页
3. 每个 section 前插入目录页或章节概览页
4. 最后一个 section 为总结

## 内容规范
- 每页 3-6 个要点，控制在 1-2 行
- 数学公式使用 \\[ ... \\] 或 align 环境
- 表格使用 booktabs 三线表
- 标题页要简洁大气，内容页要有清晰的层次
- 不要生成与用户文案内容无关的填充内容
"""


def build_beamer_prompt(content: str, *, style: str = "academic", slide_count: int = 0) -> str:
    count_rule = ""
    if slide_count > 0:
        count_rule = f"\n- 页数目标：大约 {slide_count} 页（含标题页）"
    return f"""风格：{style}{count_rule}

请根据以下文案内容，生成完整的 LaTeX Beamer 演示文稿代码。
要求直接输出可编译的 .tex 文件，必须从 \\documentclass 开始，到 \\end{{document}} 结束。
如果内容较长，请自动拆分为多个 section 和 frame。

文案内容：
{content}
"""


def clean_latex_response(text: str) -> str:
    clean = (text or "").strip()
    if clean.startswith("```"):
        clean = re.sub(r"^```(?:latex|tex)?\s*", "", clean, flags=re.I)
        clean = re.sub(r"\s*```$", "", clean)
    start = clean.find("\\documentclass")
    end = clean.rfind("\\end{document}")
    if start >= 0 and end >= start:
        clean = clean[start : end + len("\\end{document}")]
    return clean.strip()


def build_local_beamer_latex(content: str, *, style: str = "academic", slide_count: int = 0) -> str:
    title, subtitle = _infer_meta(content)
    sections = _split_into_sections(content, max_sections=_target_section_count(content, slide_count))
    if not sections:
        sections = [{"title": "核心内容", "points": _points_from_text(content)}]

    section_titles = [section["title"] for section in sections]
    if not title:
        title = section_titles[0] if section_titles else "Presentation"

    body: list[str] = []
    body.append("\\begin{frame}[plain]")
    body.append("  \\titlepage")
    body.append("\\end{frame}")

    body.append("\\begin{frame}{内容总览}")
    body.append("  \\vfill")
    body.append("  \\begin{center}")
    body.append("    \\begin{minipage}{0.82\\textwidth}")
    body.append("      \\begin{itemize}")
    body.append("        \\setlength{\\itemsep}{0.3\\baselineskip}")
    for index, section_title in enumerate(section_titles, start=1):
        body.append(
            f"        \\item[\\textcolor{{black}}{{\\textbf{{{index}.}}}}] "
            f"\\textcolor{{black}}{{{_escape_tex(section_title)}}}"
        )
    body.append("      \\end{itemize}")
    body.append("    \\end{minipage}")
    body.append("  \\end{center}")
    body.append("  \\vfill")
    body.append("\\end{frame}")

    for index, section in enumerate(sections, start=1):
        section_title = section["title"]
        body.append(f"\\section{{{_escape_tex(section_title)}}}")
        body.append(f"\\begin{{frame}}{{{_escape_tex(section_title)}}}")
        body.append("  \\vfill")
        body.append("  \\begin{center}")
        body.append("    \\begin{minipage}{0.8\\textwidth}")
        body.append("      \\begin{itemize}")
        body.append("        \\setlength{\\itemsep}{0.28\\baselineskip}")
        body.append(f"        \\item[\\textcolor{{black}}{{\\textbf{{{index}.}}}}] { _escape_tex(section_title) }")
        for point_index, point in enumerate(section["points"][:5], start=1):
            overlay = f"<{point_index}->" if len(section["points"]) > 3 else ""
            body.append(f"        \\item{overlay} {_format_point(point)}")
        body.append("      \\end{itemize}")
        body.append("    \\end{minipage}")
        body.append("  \\end{center}")
        body.append("  \\vfill")
        body.append("\\end{frame}")

        if len(section["points"]) > 5:
            body.extend(_build_extra_frames(section_title, section["points"][5:]))

    preamble = [
        "\\documentclass[10pt, aspectratio=169]{ctexbeamer}",
        "\\usetheme{Madrid}",
        "\\usepackage{amsmath, amssymb, amsthm}",
        "\\usepackage{graphicx}",
        "\\usepackage{booktabs}",
        "\\usepackage{multirow}",
        "\\usepackage{caption}",
        "\\usepackage{hyperref}",
        "\\usepackage{tikz}",
        "\\usetikzlibrary{shapes.callouts, tikzmark}",
        "\\usetikzlibrary{shapes, positioning}",
        "\\definecolor{myline}{RGB}{0,116,112}",
        "\\definecolor{myblue}{RGB}{40,100,180}",
        "\\definecolor{whitebg}{RGB}{255,255,255}",
        "\\setbeamertemplate{navigation symbols}{}",
        "\\setbeamertemplate{footline}{}",
        "\\setbeamertemplate{frametitle}{%",
        "  \\vspace*{0.2cm}%",
        "  \\begin{beamercolorbox}[wd=\\paperwidth,leftskip=0.5cm,rightskip=0.5cm,ht=0.3cm,dp=0pt]{whitebg}%",
        "    \\usebeamerfont{frametitle}\\textcolor{black}{\\insertframetitle}%",
        "  \\end{beamercolorbox}%",
        "  \\vspace{0pt}%",
        "  \\begin{tikzpicture}[remember picture, overlay]",
        "    \\draw[myline, line width=1.5pt] ([yshift=-1.3cm] current page.north west) -- ([yshift=-1.3cm] current page.north east);",
        "  \\end{tikzpicture}%",
        "  \\vspace{0.1cm}%",
        "}",
        "\\setbeamertemplate{title page}{%",
        "  \\begin{tikzpicture}[remember picture, overlay]",
        "    \\draw[line width=1.5pt, color=myline] ([yshift=-40pt] current page.north west) -- ([yshift=-40pt] current page.north east);",
        "    \\node[anchor=north west, inner sep=0, minimum width=0.25\\paperwidth, minimum height=39pt, fill=gray!30, text=black, align=center] at (current page.north west) {KGTS};",
        "  \\end{tikzpicture}%",
        "  \\vspace*{36pt}",
        "  \\begin{center}",
        "    \\begin{tikzpicture}",
        "      \\node[draw=none, inner sep=8pt, fill=white, text=black, align=center, font=\\Huge\\bfseries] (titlebox) {\\inserttitle};",
        "      \\node[draw=none, rounded corners=2pt, inner sep=8pt, fill=white, text=black, align=center, font=\\large, below=5pt of titlebox] (subtitlebox) {\\insertsubtitle};",
        "      \\node[below=5pt of subtitlebox.south east, anchor=north east, align=center, text=black] {\\insertauthor \\\\[3pt]\\insertdate};",
        "    \\end{tikzpicture}",
        "  \\end{center}",
        "}",
        f"\\title{{{_escape_tex(title)}}}",
        f"\\subtitle{{{_escape_tex(subtitle)}}}",
        "\\author{}",
        "\\date{\\today}",
        "\\begin{document}",
    ]

    return "\n".join(preamble + body + ["\\end{document}"])


def parse_latex_to_slides(latex: str) -> dict[str, Any]:
    result: dict[str, Any] = {
        "title": _extract_command(latex, "title") or "Presentation",
        "subtitle": _extract_command(latex, "subtitle") or "",
        "author": _extract_command(latex, "author") or "",
        "date": _extract_command(latex, "date") or "",
        "slides": [],
    }

    body_match = re.search(r"\\begin\{document\}(.+?)\\end\{document\}", latex, re.S)
    body = body_match.group(1) if body_match else latex
    frame_pattern = re.compile(r"\\begin\{frame\}(?:\[[^\]]*\])?(?:\{([^}]*)\})?(.*?)\\end\{frame\}", re.S)
    frames = list(frame_pattern.finditer(body))

    for index, match in enumerate(frames):
        frame_head = (match.group(1) or "").strip()
        frame_body = match.group(2) or ""
        slide_type = "title" if "\\titlepage" in frame_body else ("toc" if "\\tableofcontents" in frame_body else "content")
        title = frame_head or _extract_frame_title(frame_body)
        if not title:
            title = result["title"] if slide_type == "title" else f"Slide {index + 1}"
        items = _extract_items(frame_body)
        equations = _extract_equations(frame_body)
        plain = _clean_latex_text(_remove_latex_blocks(frame_body)).strip()
        if plain and not items and slide_type != "title":
            items = [line for line in re.split(r"[\n。；;]+", plain) if line.strip()][:6]
        result["slides"].append(
            {
                "id": index,
                "type": slide_type,
                "title": title,
                "subtitle": "",
                "items": items,
                "equations": equations,
                "notes": "",
            }
        )

    if not result["slides"]:
        result["slides"].append(
            {
                "id": 0,
                "type": "title",
                "title": result["title"],
                "subtitle": result["subtitle"],
                "items": [],
                "equations": [],
                "notes": "",
            }
        )
    return result


def generate_pptx(slides_data: dict[str, Any]) -> bytes:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    for slide_data in slides_data.get("slides") or []:
        slide_type = slide_data.get("type") or "content"
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        _add_top_line(slide)
        if slide_type == "title":
            _add_title_slide(slide, slides_data, slide_data)
        elif slide_type == "toc":
            _add_toc_slide(slide, slide_data)
        else:
            _add_content_slide(slide, slide_data)

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def _extract_command(latex: str, cmd: str) -> str | None:
    match = re.search(rf"\\{cmd}(?:\[[^\]]*\])?\{{", latex)
    if not match:
        return None
    content = _match_braces(latex, match.end() - 1)
    return _clean_latex_text(content) if content else None


def _match_braces(text: str, open_pos: int) -> str | None:
    depth = 0
    for index in range(open_pos, len(text)):
        if text[index] == "{":
            depth += 1
        elif text[index] == "}":
            depth -= 1
            if depth == 0:
                return text[open_pos + 1 : index]
    return None


def _extract_frame_title(frame: str) -> str:
    match = re.search(r"\\frametitle\{", frame)
    if match:
        content = _match_braces(frame, match.end() - 1)
        if content:
            return _clean_latex_text(content)
    header = re.match(r"\s*\{", frame)
    if header:
        content = _match_braces(frame, header.end() - 1)
        if content:
            return _clean_latex_text(content)
    return ""


def _extract_items(frame: str) -> list[str]:
    items: list[str] = []
    pattern = re.compile(r"\\item(?:<[^>]+>)?(?:\[[^\]]*\])?\s*(.*?)(?=\\item|\\end\{itemize\}|\\end\{enumerate\}|$)", re.S)
    for match in pattern.finditer(frame):
        item = _clean_latex_text(match.group(1)).strip()
        if item and not item.startswith("\\"):
            items.append(item)
    return items


def _extract_equations(frame: str) -> list[str]:
    equations: list[str] = []
    patterns = [
        r"\\\[(.+?)\\\]",
        r"\$\$(.+?)\$\$",
        r"\\begin\{equation\*?\}(.+?)\\end\{equation\*?\}",
        r"\\begin\{align\*?\}(.+?)\\end\{align\*?\}",
    ]
    for pattern in patterns:
        for match in re.finditer(pattern, frame, re.S):
            equations.append(match.group(1).strip())
    return equations


def _remove_latex_blocks(text: str) -> str:
    text = re.sub(r"\\begin\{itemize\}.+?\\end\{itemize\}", "", text, flags=re.S)
    text = re.sub(r"\\begin\{enumerate\}.+?\\end\{enumerate\}", "", text, flags=re.S)
    text = re.sub(r"\\begin\{.*?\}.+?\\end\{.*?\}", "", text, flags=re.S)
    text = re.sub(r"\\titlepage|\\frametitle\{.*?\}", "", text, flags=re.S)
    return text


def _clean_latex_text(text: str | None) -> str:
    if not text:
        return ""
    clean = text
    clean = re.sub(r"\\text(?:bf|it|color)?(?:\{[^{}]*\})?\{([^{}]*)\}", r"\1", clean)
    clean = re.sub(r"\\[a-zA-Z]+\*?(?:\[[^\]]*\])?(?:\{([^{}]*)\})?", lambda m: m.group(1) or "", clean)
    clean = clean.replace("\\%", "%").replace("\\&", "&").replace("\\_", "_")
    clean = re.sub(r"[{}]", "", clean)
    clean = re.sub(r"\s+", " ", clean)
    return clean.strip()


def _add_top_line(slide: Any) -> None:
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0.95), Inches(13.333), Pt(2.5))
    line.fill.solid()
    line.fill.fore_color.rgb = RGBColor(0, 116, 112)
    line.line.fill.background()


def _add_title_slide(slide: Any, meta: dict[str, Any], slide_data: dict[str, Any]) -> None:
    title_box = slide.shapes.add_textbox(Inches(1.4), Inches(2.2), Inches(10.5), Inches(1.2))
    p = title_box.text_frame.paragraphs[0]
    p.text = slide_data.get("title") or meta.get("title") or "Presentation"
    p.font.size = Pt(34)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    subtitle = slide_data.get("subtitle") or meta.get("subtitle") or ""
    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(2), Inches(3.55), Inches(9.4), Inches(0.7))
        p = sub_box.text_frame.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(18)
        p.alignment = PP_ALIGN.CENTER


def _add_toc_slide(slide: Any, slide_data: dict[str, Any]) -> None:
    title_box = slide.shapes.add_textbox(Inches(0.55), Inches(0.35), Inches(12), Inches(0.55))
    p = title_box.text_frame.paragraphs[0]
    p.text = slide_data.get("title") or "Contents"
    p.font.size = Pt(22)
    p.font.bold = True

    items = slide_data.get("items") or []
    body = slide.shapes.add_textbox(Inches(1.0), Inches(1.55), Inches(11.2), Inches(5.2))
    tf = body.text_frame
    tf.word_wrap = True
    for index, item in enumerate(items):
        p = tf.paragraphs[0] if index == 0 else tf.add_paragraph()
        p.text = f"{index + 1}. {item}"
        p.font.size = Pt(17)
        p.space_after = Pt(10)


def _add_content_slide(slide: Any, slide_data: dict[str, Any]) -> None:
    title_box = slide.shapes.add_textbox(Inches(0.55), Inches(0.35), Inches(12), Inches(0.55))
    p = title_box.text_frame.paragraphs[0]
    p.text = slide_data.get("title") or ""
    p.font.size = Pt(22)
    p.font.bold = True

    body = slide.shapes.add_textbox(Inches(0.85), Inches(1.55), Inches(11.6), Inches(5.4))
    tf = body.text_frame
    tf.word_wrap = True
    items = [str(item) for item in slide_data.get("items") or []]
    equations = [str(item) for item in slide_data.get("equations") or []]
    lines = items + [f"${eq}$" for eq in equations]
    if not lines:
        lines = [str(slide_data.get("notes") or "")]
    for index, line in enumerate(lines):
        p = tf.paragraphs[0] if index == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = Pt(17 if index < len(items) else 15)
        p.space_after = Pt(8)


def _infer_meta(content: str) -> tuple[str, str]:
    lines = [line.strip() for line in (content or "").splitlines() if line.strip()]
    if not lines:
        return "Presentation", ""
    first = _clean_heading(lines[0])
    subtitle = ""
    if len(lines) > 1:
        subtitle = _clean_heading(lines[1])[:80]
    if len(first) > 60:
        first = first[:60]
    return first or "Presentation", subtitle


def _clean_heading(text: str) -> str:
    text = re.sub(r"^\s*#+\s*", "", text)
    text = re.sub(r"^\s*(?:[一二三四五六七八九十]+[、.])\s*", "", text)
    text = text.strip(" -:：")
    return _clean_latex_text(text)


def _target_section_count(content: str, slide_count: int) -> int:
    if slide_count > 0:
        return max(1, min(6, slide_count // 3 or 1))
    length = len(content or "")
    if length < 400:
        return 2
    if length < 1200:
        return 3
    if length < 2500:
        return 4
    return 5


@dataclass
class _Section:
    title: str
    points: list[str]


def _split_into_sections(content: str, *, max_sections: int) -> list[dict[str, Any]]:
    text = (content or "").strip()
    if not text:
        return []

    explicit_sections = _extract_explicit_sections(text)
    if explicit_sections:
        sections = explicit_sections
    else:
        sections = _chunk_into_sections(_normalize_sentences(text), max_sections=max_sections)

    result: list[dict[str, Any]] = []
    for section in sections[:max_sections]:
        title = _clean_heading(section.title) or "核心内容"
        points = [p for p in (_normalize_point(point) for point in section.points) if p]
        if not points:
            continue
        result.append({"title": title, "points": points})
    return result


def _extract_explicit_sections(text: str) -> list[_Section]:
    sections: list[_Section] = []
    current_title = ""
    current_points: list[str] = []

    def flush() -> None:
        nonlocal current_title, current_points
        if current_points:
            sections.append(_Section(title=current_title or current_points[0][:20], points=current_points[:]))
        current_title = ""
        current_points = []

    for raw_line in text.splitlines():
        line = raw_line.strip()
        if not line:
            continue
        if _looks_like_section_heading(line):
            flush()
            current_title = _clean_heading(line)
            continue
        if _looks_like_bullet(line):
            current_points.append(_clean_heading(_strip_bullet(line)))
        else:
            current_points.append(_clean_heading(line))

    flush()
    return [section for section in sections if section.points]


def _chunk_into_sections(points: list[str], *, max_sections: int) -> list[_Section]:
    if not points:
        return []
    if len(points) <= max_sections:
        return [_Section(title=point[:20] or f"部分 {index + 1}", points=[point]) for index, point in enumerate(points)]

    chunk_size = max(1, len(points) // max_sections)
    sections: list[_Section] = []
    for index in range(0, len(points), chunk_size):
        chunk = points[index : index + chunk_size]
        if not chunk:
            continue
        title = chunk[0][:20] or f"部分 {len(sections) + 1}"
        sections.append(_Section(title=title, points=chunk))
    return sections


def _normalize_sentences(text: str) -> list[str]:
    text = text.replace("\r", "\n")
    paragraphs = [part.strip() for part in re.split(r"\n{2,}", text) if part.strip()]
    points: list[str] = []
    for paragraph in paragraphs:
        paragraph = _strip_heading(paragraph)
        parts = re.split(r"[。！？!?；;\n]+", paragraph)
        for part in parts:
            clean = _clean_heading(part)
            if clean:
                points.append(clean)
    if not points:
        points = [_clean_heading(line) for line in text.splitlines() if _clean_heading(line)]
    return points


def _normalize_point(point: str) -> str:
    clean = _clean_heading(point)
    if not clean:
        return ""
    if len(clean) > 120:
        clean = clean[:117] + "..."
    return clean


def _strip_heading(text: str) -> str:
    return re.sub(r"^\s*(?:#+\s*|[一二三四五六七八九十]+[、.]\s*)", "", text).strip()


def _looks_like_section_heading(line: str) -> bool:
    if len(line) > 40:
        return False
    return bool(
        re.match(r"^\s*#+\s*\S+", line)
        or re.match(r"^\s*[一二三四五六七八九十]+[、.]\s*\S+", line)
        or re.match(r"^\s*(?:第[一二三四五六七八九十0-9]+[章节部分篇节])\s*\S*", line)
    )


def _looks_like_bullet(line: str) -> bool:
    return bool(re.match(r"^\s*(?:[-*•]|[0-9]+[.)]|[一二三四五六七八九十]+[.)])\s+\S+", line))


def _strip_bullet(line: str) -> str:
    return re.sub(r"^\s*(?:[-*•]|[0-9]+[.)]|[一二三四五六七八九十]+[.)])\s*", "", line).strip()


def _format_point(point: str) -> str:
    clean = _escape_tex(point)
    if _looks_math_like(point):
        return f"\\[{point.strip()}\\]"
    return clean


def _looks_math_like(text: str) -> bool:
    return bool(re.search(r"[=<>^_]|\\frac|\\sum|\\int|\$.*\$", text))


def _escape_tex(text: str) -> str:
    s = str(text or "")
    replacements = {
        "\\": r"\textbackslash{}",
        "&": r"\&",
        "%": r"\%",
        "$": r"\$",
        "#": r"\#",
        "_": r"\_",
        "{": r"\{",
        "}": r"\}",
        "~": r"\textasciitilde{}",
        "^": r"\textasciicircum{}",
    }
    for src, dst in replacements.items():
        s = s.replace(src, dst)
    return s


def _build_extra_frames(title: str, remaining_points: list[str]) -> list[str]:
    frames: list[str] = []
    chunk_size = 4
    for index in range(0, len(remaining_points), chunk_size):
        chunk = remaining_points[index : index + chunk_size]
        if not chunk:
            continue
        frames.append(f"\\begin{{frame}}{{{_escape_tex(title)}（续）}}")
        frames.append("  \\begin{itemize}")
        frames.append("    \\setlength{\\itemsep}{0.28\\baselineskip}")
        for point in chunk:
            frames.append(f"    \\item {_format_point(point)}")
        frames.append("  \\end{itemize}")
        frames.append("\\end{frame}")
    return frames

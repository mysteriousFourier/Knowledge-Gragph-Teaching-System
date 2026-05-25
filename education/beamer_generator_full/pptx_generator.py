"""PPTX 生成器 — 将结构化幻灯片 JSON 转为 .pptx 文件"""
import io
import os
import re
from pathlib import Path
from html import escape as html_escape
from html.parser import HTMLParser
from typing import Optional

import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.enum.shapes import MSO_SHAPE

try:
    from PIL import Image as PILImage
except Exception:  # pragma: no cover - optional dependency fallback
    PILImage = None


# 主色调（与 LaTeX 模板一致）
COLOR_MYLINE = RGBColor(0, 116, 112)
COLOR_MYBLUE = RGBColor(40, 100, 180)
COLOR_BLACK = RGBColor(0, 0, 0)
COLOR_WHITE = RGBColor(255, 255, 255)
COLOR_GRAY = RGBColor(153, 153, 153)
COLOR_LIGHT_GRAY = RGBColor(240, 242, 245)

# 16:9 尺寸
SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)
BODY_LEFT = Inches(0.72)
BODY_RIGHT = Inches(0.72)
BODY_TOP = Inches(1.62)
BODY_BOTTOM = Inches(0.46)


def _prepare_text_frame(tf, font_size=None, color=None, bold=False, italic=False, align=None, vertical_anchor=None):
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    if vertical_anchor is not None:
        tf.vertical_anchor = vertical_anchor
    tf.margin_left = Inches(0.04)
    tf.margin_right = Inches(0.04)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    p = tf.paragraphs[0]
    if font_size is not None:
        p.font.size = Pt(font_size)
    if color is not None:
        p.font.color.rgb = color
    p.font.bold = bold
    p.font.italic = italic
    if align is not None:
        p.alignment = align
    return p


def _set_text_or_math(p, text: str, font_size=None, color=None, bold=False, italic=False, align=None):
    """Set paragraph text while preserving inline LaTeX-style math as rendered images are not available in PPTX."""
    text = _repair_latex_artifacts(text)
    p.clear()
    p.font.size = Pt(font_size) if font_size is not None else p.font.size
    if color is not None:
        p.font.color.rgb = color
    p.font.bold = bold
    p.font.italic = italic
    if align is not None:
        p.alignment = align
    for idx, chunk in enumerate(_split_text_and_math(text)):
        run = p.add_run()
        run.text = chunk
        if idx == 0:
            continue


def _split_text_and_math(text: str) -> list[str]:
    parts = []
    last = 0
    pattern = re.compile(r"\\\[[\s\S]*?\\\]|\\\([\s\S]*?\\\)|\$\$[\s\S]*?\$\$|(^|[^\\])\$[^$\n]+?\$")
    for m in pattern.finditer(text):
        start = m.start()
        if m.group(1):
            start += len(m.group(1))
        if start > last:
            parts.append(text[last:start])
        parts.append(text[start:m.end()])
        last = m.end()
    if last < len(text):
        parts.append(text[last:])
    return [p for p in parts if p]


def _style_paragraph(p, font_size=None, color=None, bold=False, italic=False, align=None):
    if font_size is not None:
        p.font.size = Pt(font_size)
    if color is not None:
        p.font.color.rgb = color
    p.font.bold = bold
    p.font.italic = italic
    if align is not None:
        p.alignment = align


def _parse_inline_style(style_text: str) -> dict[str, str]:
    result: dict[str, str] = {}
    for part in str(style_text or "").split(";"):
        if ":" not in part:
            continue
        key, value = part.split(":", 1)
        key = key.strip().lower()
        value = value.strip()
        if key and value:
            result[key] = value
    return result


def _rich_html_looks_safe(html: str) -> bool:
    if not html:
      return False
    lowered = str(html).lower()
    return not any(token in lowered for token in ("katex", "data-latex"))


class _RichHtmlFragmentParser(HTMLParser):
    def __init__(self):
        super().__init__()
        self.paragraphs: list[list[dict[str, object]]] = [[]]
        self._stack: list[dict[str, object]] = [self._base_style()]

    def _base_style(self) -> dict[str, object]:
        return {"bold": False, "italic": False, "color": None}

    def _current_style(self) -> dict[str, object]:
        return dict(self._stack[-1])

    def _push_style(self, style: dict[str, object]) -> None:
        self._stack.append(style)

    def _pop_style(self) -> None:
        if len(self._stack) > 1:
            self._stack.pop()

    def _new_paragraph(self) -> None:
        if self.paragraphs and self.paragraphs[-1]:
            self.paragraphs.append([])

    def handle_starttag(self, tag, attrs):
        tag = (tag or "").lower()
        attrs = dict(attrs or [])
        if tag in {"br"}:
            self.paragraphs.append([])
            return
        if tag in {"p", "div"}:
            self._new_paragraph()
            return

        style = self._current_style()
        if tag in {"b", "strong"}:
            style["bold"] = True
        if tag in {"i", "em"}:
            style["italic"] = True
        if tag == "font" and attrs.get("color"):
            style["color"] = attrs.get("color")
        inline = _parse_inline_style(attrs.get("style", ""))
        if inline.get("color"):
            style["color"] = inline.get("color")
        if inline.get("font-weight") and inline.get("font-weight").lower() not in {"normal", "400"}:
            style["bold"] = True
        if inline.get("font-style") and inline.get("font-style").lower() == "italic":
            style["italic"] = True
        if tag in {"span", "font", "b", "strong", "i", "em"}:
            self._push_style(style)

    def handle_endtag(self, tag):
        tag = (tag or "").lower()
        if tag in {"p", "div"}:
            self._new_paragraph()
            return
        if tag in {"span", "font", "b", "strong", "i", "em"}:
            self._pop_style()

    def handle_data(self, data):
        if not data:
            return
        current = self.paragraphs[-1]
        style = self._current_style()
        pieces = data.splitlines(True)
        for piece in pieces:
            if not piece:
                continue
            if piece in {"\n", "\r\n", "\r"}:
                self.paragraphs.append([])
                current = self.paragraphs[-1]
                continue
            current.append({"text": piece, "style": style})


def _write_rich_html_to_text_frame(tf, html: str, fallback_text: str = "", font_size=None, color=None, bold=False, italic=False, align=None, vertical_anchor=None) -> bool:
    html = str(html or "")
    if not html.strip() or not _rich_html_looks_safe(html):
        return False
    parser = _RichHtmlFragmentParser()
    try:
        parser.feed(html)
        parser.close()
    except Exception:
        return False
    paragraphs = [p for p in parser.paragraphs if any(str(f.get("text") or "").strip() for f in p)]
    if not paragraphs:
        return False
    try:
        tf.clear()
    except Exception:
        try:
            tf._element.clear_content()
        except Exception:
            pass
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    if vertical_anchor is not None:
        tf.vertical_anchor = vertical_anchor
    tf.margin_left = Inches(0.04)
    tf.margin_right = Inches(0.04)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    for idx, fragments in enumerate(paragraphs):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        if font_size is not None:
            p.font.size = Pt(font_size)
        if color is not None:
            p.font.color.rgb = color
        p.font.bold = bold
        p.font.italic = italic
        if align is not None:
            p.alignment = align
        for fragment in fragments:
            text = str(fragment.get("text") or "")
            if not text:
                continue
            run = p.add_run()
            run.text = text
            style = fragment.get("style") or {}
            run.font.size = Pt(font_size) if font_size is not None else run.font.size
            run.font.color.rgb = _parse_rgb_color(str(style.get("color") or "")) or color or COLOR_BLACK
            run.font.bold = bool(style.get("bold", False)) or bold
            run.font.italic = bool(style.get("italic", False)) or italic
    return True


def _has_math_like(text: str) -> bool:
    s = _repair_latex_artifacts(text)
    return bool(re.search(r"\\|\$|[_^]|\\bar|\\frac|\\sqrt|\\begin\{", s))


def _repair_latex_artifacts(text: str) -> str:
    """Undo LaTeX escaping artifacts introduced by editor round-trips."""
    s = str(text or "")
    replacements = {
        r"\textbackslash\{\}": "\\",
        r"\textbackslash{}": "\\",
        r"\textbackslash ": "\\",
        r"\textasciicircum{}": "^",
        r"\textasciitilde{}": "~",
        r"\_": "_",
        r"\$": "$",
        r"\&": "&",
        r"\%": "%",
        r"\{": "{",
        r"\}": "}",
    }
    for old, new in replacements.items():
        s = s.replace(old, new)
    return s


def _render_text_png(text: str, font_size: float = 16, color: str = "#000000", align: str = "center") -> io.BytesIO:
    s = _repair_latex_artifacts(text)
    lines = s.splitlines() or [""]
    width = max(1.6, min(8.5, max(len(line) for line in lines) * max(font_size, 10) * 0.06 + 0.7))
    height = max(0.45, 0.32 * len(lines) + 0.25)
    fig = plt.figure(figsize=(width, height), dpi=220)
    fig.patch.set_alpha(0)
    try:
        ax = fig.add_axes([0, 0, 1, 1])
        ax.axis("off")
        x = 0.5 if align == "center" else (0.03 if align == "left" else 0.97)
        ha = "center" if align == "center" else ("left" if align == "left" else "right")
        for idx, line in enumerate(lines):
            y = 0.5 if len(lines) == 1 else 1 - ((idx + 0.5) / len(lines))
            fig.text(x, y, line, ha=ha, va="center", fontsize=font_size, color=color)
        out = io.BytesIO()
        fig.savefig(out, format="png", transparent=True, bbox_inches="tight", pad_inches=0.03)
        out.seek(0)
        return out
    finally:
        plt.close(fig)


def _add_text_or_image(slide, left, top, width, height, text: str, font_size=14, color=COLOR_BLACK, bold=False, italic=False, align=PP_ALIGN.LEFT):
    s = _repair_latex_artifacts(text)
    align_name = "center" if align == PP_ALIGN.CENTER else ("right" if align == PP_ALIGN.RIGHT else "left")
    if _has_math_like(s):
        img = _render_text_png(s, font_size=font_size, color="#2864B4" if color == COLOR_MYBLUE else "#111111", align=align_name)
        pic = slide.shapes.add_picture(img, left, top, width=width, height=height)
        return pic
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    p = _prepare_text_frame(tf, font_size=font_size, color=color, bold=bold, italic=italic, align=align)
    p.text = s
    return box


def _add_callout(slide, callout_data: dict):
    if not isinstance(callout_data, dict):
        return
    text = _repair_latex_artifacts(callout_data.get("text", "") or "").strip()
    if not text:
        return
    x = _editor_px_to_slide_x(callout_data.get("x", 130))
    y = _editor_px_to_slide_y(callout_data.get("y", 180))
    w = _editor_px_to_slide_x(callout_data.get("width", 250))
    h = _editor_px_to_slide_y(callout_data.get("height", 92))
    w = max(Inches(0.9), min(w, SLIDE_WIDTH - x))
    h = max(Inches(0.45), min(h, SLIDE_HEIGHT - y))

    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGULAR_CALLOUT, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLOR_WHITE
    shape.line.color.rgb = COLOR_MYBLUE
    shape.line.width = Pt(1.5)
    shape.adjustments[0] = 0.22
    shape.adjustments[1] = 0.22

    inner_left = x + Inches(0.08)
    inner_top = y + Inches(0.05)
    inner_w = max(Inches(0.5), w - Inches(0.16))
    inner_h = max(Inches(0.3), h - Inches(0.1))
    _add_text_or_image(
        slide,
        inner_left,
        inner_top,
        inner_w,
        inner_h,
        text,
        font_size=_safe_float(callout_data.get("fontSize", 12), 12),
        color=COLOR_BLACK,
        align=PP_ALIGN.CENTER if str(callout_data.get("align", "center")).lower() == "center" else PP_ALIGN.LEFT,
    )


def _add_textbox_or_image(slide, left, top, width, height, text: str, font_size=14, color=COLOR_BLACK, bold=False, italic=False, align=PP_ALIGN.LEFT, bg=None, border=None):
    text = _repair_latex_artifacts(text)
    if _has_math_like(text):
        return _add_text_or_image(slide, left, top, width, height, text, font_size=font_size, color=color, bold=bold, italic=italic, align=align)
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    tf.margin_left = Inches(0.06)
    tf.margin_right = Inches(0.06)
    tf.margin_top = Inches(0.03)
    tf.margin_bottom = Inches(0.03)
    if bg:
        box.fill.solid()
        box.fill.fore_color.rgb = bg
    if border:
        box.line.color.rgb = border
        box.line.width = Pt(0.5)
    p = _prepare_text_frame(tf, font_size=font_size, color=color, bold=bold, italic=italic, align=align)
    p.text = text
    return box


def _add_callout(slide, callout_data: dict):
    if not isinstance(callout_data, dict):
        return
    text = _repair_latex_artifacts(callout_data.get("text", "") or "").strip()
    if not text:
        return
    x = _editor_px_to_slide_x(callout_data.get("x", 130))
    y = _editor_px_to_slide_y(callout_data.get("y", 180))
    w = _editor_px_to_slide_x(callout_data.get("width", 250))
    h = _editor_px_to_slide_y(callout_data.get("height", 92))
    w = max(Inches(0.9), min(w, SLIDE_WIDTH - x))
    h = max(Inches(0.45), min(h, SLIDE_HEIGHT - y))

    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGULAR_CALLOUT, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLOR_WHITE
    shape.line.color.rgb = COLOR_MYBLUE
    shape.line.width = Pt(1.5)
    shape.adjustments[0] = 0.22
    shape.adjustments[1] = 0.22

    tf = shape.text_frame
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    tf.margin_left = Inches(0.06)
    tf.margin_right = Inches(0.06)
    tf.margin_top = Inches(0.05)
    tf.margin_bottom = Inches(0.05)
    p = tf.paragraphs[0]
    _set_text_or_math(
        p,
        text,
        font_size=_safe_float(callout_data.get("fontSize", 12), 12),
        color=COLOR_BLACK,
        align=PP_ALIGN.CENTER if str(callout_data.get("align", "center")).lower() == "center" else PP_ALIGN.LEFT,
    )


def _add_math_textbox(slide, left, top, width, height, text: str, font_size=14, color=COLOR_MYBLUE, align=PP_ALIGN.CENTER, italic=False):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    p = _prepare_text_frame(tf, font_size=font_size, color=color, italic=italic, align=align)
    _set_text_or_math(p, text, font_size=font_size, color=color, italic=italic, align=align)
    return box


def _image_size(image_source):
    if PILImage is None:
        return None

    try:
        if isinstance(image_source, (str, bytes, os.PathLike)):
            with PILImage.open(image_source) as img:
                return img.size

        pos = None
        if hasattr(image_source, "tell"):
            try:
                pos = image_source.tell()
            except Exception:
                pos = None
        if hasattr(image_source, "seek"):
            image_source.seek(0)
        with PILImage.open(image_source) as img:
            size = img.size
        if pos is not None and hasattr(image_source, "seek"):
            image_source.seek(pos)
        return size
    except Exception:
        return None


def generate_pptx(slides_data: dict, upload_dir: str = "") -> bytes:
    """
    将结构化幻灯片数据生成 PPTX 二进制内容。

    参数:
        slides_data: parse_latex_to_slides() 返回的字典
        upload_dir: 上传图片的存储目录

    返回: PPTX 文件的 bytes
    """
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    meta = slides_data
    slides = slides_data.get("slides", [])
    figure_asset_paths = _build_figure_asset_lookup(slides_data.get("figure_assets", {}), upload_dir)

    for slide_data in slides:
        slide_type = slide_data.get("type", "content")

        if slide_type == "title":
            _add_title_slide(prs, meta, slide_data)
        elif slide_type == "toc":
            _add_toc_slide(prs, slide_data)
        else:
            _add_content_slide(prs, slide_data, upload_dir, figure_asset_paths)

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf.read()


# ================================================================
#  标题页
# ================================================================
def _add_title_slide(prs: Presentation, meta: dict, slide_data: dict):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # 空白布局

    # 顶部横线
    _add_top_line(slide)

    # 左上角机构标签
    tag = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(3.5), Inches(0.55)
    )
    tag.fill.solid()
    tag.fill.fore_color.rgb = COLOR_LIGHT_GRAY
    tag.line.fill.background()
    tf = tag.text_frame
    p = _prepare_text_frame(tf, font_size=11, color=COLOR_BLACK, align=PP_ALIGN.CENTER, vertical_anchor=MSO_ANCHOR.MIDDLE)
    p.text = "Presentation"

    # 主标题
    title_box = slide.shapes.add_textbox(
        Inches(1.35), Inches(2.05), Inches(10.6), Inches(1.8)
    )
    tf = title_box.text_frame
    title_rich = slide_data.get("titleRichHtml") or meta.get("titleRichHtml") or ""
    if not _write_rich_html_to_text_frame(tf, title_rich, meta.get("title", ""), font_size=34, color=COLOR_BLACK, bold=True, align=PP_ALIGN.CENTER):
        p = _prepare_text_frame(tf, font_size=34, color=COLOR_BLACK, bold=True, align=PP_ALIGN.CENTER)
        p.text = meta.get("title", "")

    # 副标题
    if meta.get("subtitle"):
        subtitle_box = slide.shapes.add_textbox(
            Inches(1.9), Inches(3.7), Inches(9.5), Inches(0.8)
        )
        tf = subtitle_box.text_frame
        subtitle_rich = slide_data.get("subtitleRichHtml") or meta.get("subtitleRichHtml") or ""
        if not _write_rich_html_to_text_frame(tf, subtitle_rich, meta["subtitle"], font_size=19, color=COLOR_BLACK, align=PP_ALIGN.CENTER):
            p = _prepare_text_frame(tf, font_size=19, color=COLOR_BLACK, align=PP_ALIGN.CENTER)
            p.text = meta["subtitle"]

    # 作者 + 日期
    info_parts = []
    if meta.get("author"):
        info_parts.append(meta["author"])
    if meta.get("date"):
        info_parts.append(meta["date"])
    if info_parts:
        info_box = slide.shapes.add_textbox(
            Inches(4.9), Inches(5.0), Inches(5.1), Inches(0.5)
        )
        tf = info_box.text_frame
        p = _prepare_text_frame(tf, font_size=13, color=COLOR_GRAY, align=PP_ALIGN.RIGHT)
        p.text = "  |  ".join(info_parts)


# ================================================================
#  目录导航页
# ================================================================
def _add_toc_slide(prs: Presentation, slide_data: dict):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    _add_top_line(slide)
    _add_frame_title(slide, slide_data.get("title", "Contents"), slide_data.get("titleRichHtml", ""))

    items = slide_data.get("items", [])
    if not items:
        return

    body_box = slide.shapes.add_textbox(
        Inches(2.9), Inches(2.0), Inches(7.6), Inches(4.8)
    )
    tf = body_box.text_frame
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE

    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"{i + 1}.  {item}"
        _style_paragraph(p, font_size=17, color=COLOR_BLACK, align=PP_ALIGN.LEFT)
        p.space_after = Pt(10)


# ================================================================
#  内容页
# ================================================================
def _add_content_slide(prs: Presentation, slide_data: dict, upload_dir: str, figure_asset_paths: dict[str, str]):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    if slide_data.get("reviewBackground"):
        _add_review_background(slide)
    _add_top_line(slide)
    _add_frame_title(slide, slide_data.get("title", ""), slide_data.get("titleRichHtml", ""))
    body_left = BODY_LEFT
    body_right = BODY_RIGHT
    body_top = BODY_TOP
    body_bottom = BODY_BOTTOM
    body_width = SLIDE_WIDTH - body_left - body_right
    body_height = SLIDE_HEIGHT - body_top - body_bottom

    # 副标题条（蓝线下方描述）
    if slide_data.get("subtitle"):
        sub_box = slide.shapes.add_textbox(
            Inches(0.6), Inches(1.08), Inches(12.0), Inches(0.4)
        )
        tf = sub_box.text_frame
        subtitle_rich = slide_data.get("subtitleRichHtml", "")
        if not _write_rich_html_to_text_frame(tf, subtitle_rich, f"  {slide_data['subtitle']}", font_size=11, color=COLOR_GRAY):
            p = _prepare_text_frame(tf, font_size=11, color=COLOR_GRAY)
            p.text = f"  {slide_data['subtitle']}"

    # 计算正文区域起始 Y
    body_top = BODY_TOP
    current_y = body_top

    # 表格
    if slide_data.get("table"):
        current_y = _add_table(slide, slide_data["table"], current_y)

    # 列表项
    items = slide_data.get("items", [])
    if items:
        equation_reserve = Inches(0.8 + 0.35 * max(1, len(slide_data.get("equations", []))))
        available_h = max(Inches(0.9), body_top + body_height - current_y - equation_reserve)
        estimated_h = Inches(min(4.0, 0.34 * max(1, len(items)) + 0.25))
        item_box = slide.shapes.add_textbox(
            body_left, current_y, body_width, min(available_h, estimated_h)
        )
        tf = item_box.text_frame
        tf.word_wrap = True
        tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE

        item_rich_html = slide_data.get("itemRichHtml", []) or []
        if item_rich_html and any(item_rich_html):
            rich_body_html = "".join(
                "<div>&bull; " + (item_rich_html[i] if i < len(item_rich_html) and item_rich_html[i] else html_escape(str(item or ""))) + "</div>"
                for i, item in enumerate(items)
            )
            if _write_rich_html_to_text_frame(
                tf,
                rich_body_html,
                "\n".join("• " + str(item or "") for item in items),
                font_size=15,
                color=COLOR_BLACK,
                align=PP_ALIGN.LEFT,
            ):
                items = []

        for i, item in enumerate(items):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = f"•  {item}"
            _style_paragraph(p, font_size=15, color=COLOR_BLACK)
            p.space_after = Pt(8)
            p.level = 0
        current_y = current_y + min(available_h, estimated_h) + Inches(0.06)

    # 公式（渲染为图片插入，避免在 PPT 中显示 LaTeX 源码）
    equations = slide_data.get("equations", [])
    if equations:
        eq_top = current_y
        for i, eq in enumerate(equations):
            img_stream = _render_equation_image(eq)
            if img_stream:
                top = eq_top + Inches(i * 0.72)
                _add_picture_fit(slide, img_stream, body_left, top, body_width, Inches(0.72))
            else:
                fallback_stream = _render_formula_source_image(eq)
                if fallback_stream:
                    top = eq_top + Inches(i * 0.72)
                    _add_picture_fit(slide, fallback_stream, body_left, top, body_width, Inches(0.72))
                else:
                    eq_box = slide.shapes.add_textbox(
                        body_left + Inches(0.05), eq_top + Inches(i * 0.55), body_width - Inches(0.1), Inches(0.5)
                    )
                    tf = eq_box.text_frame
                    p = _prepare_text_frame(tf, font_size=14, color=COLOR_MYBLUE, italic=True, align=PP_ALIGN.CENTER)
                    p.text = _formula_to_plain_text(eq)

    # 备注（callout 标注内容放在 PPT 备注区）
    if slide_data.get("notes"):
        notes_slide = slide.notes_slide
        notes_tf = notes_slide.notes_text_frame
        notes_tf.text = slide_data["notes"]

    # 图片占位
    images = slide_data.get("images", [])
    for img in images:
        img_path = img.get("path", "")
        if upload_dir:
            img_path = os.path.join(upload_dir, os.path.basename(img_path))
        if os.path.exists(img_path):
            x = _editor_px_to_slide_x(img.get("x", 0))
            y = _editor_px_to_slide_y(img.get("y", 0))
            w = _editor_px_to_slide_x(img.get("width", 200))
            h_px = _safe_float(img.get("height", 0), 0)
            h = _editor_px_to_slide_y(h_px) if h_px > 0 else Inches(1.8)
            _add_picture_fit(slide, img_path, x, y, max(Inches(0.8), w), max(Inches(0.6), h))

    _add_image_placeholders(slide, slide_data.get("placeholders", []), figure_asset_paths, upload_dir)
    _add_textboxes(slide, slide_data.get("textboxes", []))
    for callout in slide_data.get("callouts", []) or []:
        _add_callout(slide, callout)


# ================================================================
#  辅助函数
# ================================================================
def _add_review_background(slide):
    back = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_WIDTH, SLIDE_HEIGHT)
    back.fill.solid()
    back.fill.fore_color.rgb = RGBColor(217, 217, 217)
    back.line.fill.background()

    front = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.92),
        Inches(1.2),
        SLIDE_WIDTH - Inches(1.84),
        SLIDE_HEIGHT - Inches(1.85),
    )
    front.fill.solid()
    front.fill.fore_color.rgb = RGBColor(255, 255, 255)
    front.line.fill.background()


def _add_top_line(slide):
    """顶部青色横线"""
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0.95),
        SLIDE_WIDTH, Pt(2.5),
    )
    line.fill.solid()
    line.fill.fore_color.rgb = COLOR_MYLINE
    line.line.fill.background()


def _add_frame_title(slide, title: str, rich_html: str = ""):
    """页面标题（横线上方）"""
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.4), Inches(12), Inches(0.5)
    )
    tf = title_box.text_frame
    if not _write_rich_html_to_text_frame(tf, rich_html, title, font_size=22, color=COLOR_BLACK, bold=True):
        p = _prepare_text_frame(tf, font_size=22, color=COLOR_BLACK, bold=True)
        p.text = title


def _add_table(slide, table_data: dict, top_y) -> float:
    """添加表格，返回表格底部的 Y 坐标"""
    headers = table_data.get("headers", [])
    rows_data = table_data.get("rows", [])
    if not headers:
        return top_y

    n_cols = len(headers)
    n_rows = len(rows_data) + 1  # +1 表头

    table_width = Inches(11)
    table_height = Inches(0.4 * n_rows)
    left = Inches(1)

    table_shape = slide.shapes.add_table(
        n_rows, n_cols, left, top_y, table_width, table_height
    )
    table = table_shape.table

    # 表头
    for j, h in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = h
        tf = cell.text_frame
        tf.word_wrap = True
        tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        _prepare_text_frame(tf, font_size=13, color=COLOR_WHITE, bold=True, vertical_anchor=MSO_ANCHOR.MIDDLE)
        cell.fill.solid()
        cell.fill.fore_color.rgb = COLOR_MYLINE

    # 数据行
    for i, row in enumerate(rows_data):
        for j, val in enumerate(row):
            if j < n_cols:
                cell = table.cell(i + 1, j)
                cell.text = val
                tf = cell.text_frame
                tf.word_wrap = True
                tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
                _prepare_text_frame(tf, font_size=12, color=COLOR_BLACK, vertical_anchor=MSO_ANCHOR.MIDDLE)
                # 交替行背景
                if i % 2 == 1:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = RGBColor(245, 247, 250)

    return top_y + table_height + Inches(0.3)


def _add_image_placeholders(slide, placeholders: list, figure_asset_paths: dict[str, str], upload_dir: str = ""):
    """Add editable image placeholder boxes from parsed or manually edited data."""
    if not placeholders:
        return

    for ph in placeholders:
        if not isinstance(ph, dict):
            continue
        x = _editor_px_to_slide_x(ph.get("x", 570))
        y = _editor_px_to_slide_y(ph.get("y", 150))
        w = _editor_px_to_slide_x(ph.get("width", 245))
        h = _editor_px_to_slide_y(ph.get("height", 230))
        w = max(Inches(0.8), min(w, SLIDE_WIDTH - x))
        h = max(Inches(0.6), min(h, SLIDE_HEIGHT - y))

        figure_path = _resolve_figure_placeholder_path(ph, figure_asset_paths, upload_dir)
        if figure_path and os.path.exists(figure_path):
            _add_picture_fit(slide, figure_path, x, y, w, h)
            continue

        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(232, 246, 245)
        shape.line.color.rgb = COLOR_MYLINE
        shape.line.width = Pt(1.5)

        tf = shape.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.text = str(ph.get("label") or "图片占位")
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = COLOR_MYLINE


def _add_textboxes(slide, textboxes: list):
    """Add free-position text boxes saved by the browser PPT editor."""
    if not textboxes:
        return

    for tb in textboxes:
        text = _repair_latex_artifacts(tb.get("text", "") if isinstance(tb, dict) else "")
        rich_html = str(tb.get("richHtml", "") if isinstance(tb, dict) else "")
        if not text.strip():
            continue

        x = _editor_px_to_slide_x(tb.get("x", 40))
        y = _editor_px_to_slide_y(tb.get("y", 190))
        w = _editor_px_to_slide_x(tb.get("width", 260))
        h = _editor_px_to_slide_y(tb.get("height", 96))
        w = max(Inches(0.6), min(w, SLIDE_WIDTH - x))
        h = max(Inches(0.25), min(h, SLIDE_HEIGHT - y))

        align = str(tb.get("align", "left") or "left").lower()
        ppt_align = PP_ALIGN.CENTER if align == "center" else (PP_ALIGN.RIGHT if align == "right" else PP_ALIGN.LEFT)
        if _has_math_like(text):
            _add_textbox_or_image(
                slide,
                x,
                y,
                w,
                h,
                text,
                font_size=_safe_float(tb.get("fontSize", 14), 14),
                color=_parse_rgb_color(tb.get("color", "")) or COLOR_BLACK,
                bold=bool(tb.get("bold", False)),
                italic=bool(tb.get("italic", False)),
                align=ppt_align,
                bg=_parse_rgb_color(tb.get("bg", "")),
                border=RGBColor(200, 205, 214),
            )
            continue

        box = slide.shapes.add_textbox(x, y, w, h)
        tf = box.text_frame
        tf.word_wrap = True
        tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        tf.margin_left = Inches(0.06)
        tf.margin_right = Inches(0.06)
        tf.margin_top = Inches(0.03)
        tf.margin_bottom = Inches(0.03)

        bg = _parse_rgb_color(tb.get("bg", ""))
        if bg:
            box.fill.solid()
            box.fill.fore_color.rgb = bg
        else:
            box.fill.background()
        box.line.color.rgb = RGBColor(200, 205, 214)
        box.line.width = Pt(0.5)

        if rich_html and _write_rich_html_to_text_frame(
            tf,
            rich_html,
            text,
            font_size=_safe_float(tb.get("fontSize", 14), 14),
            color=_parse_rgb_color(tb.get("color", "")) or COLOR_BLACK,
            bold=bool(tb.get("bold", False)),
            italic=bool(tb.get("italic", False)),
            align=ppt_align,
        ):
            continue

        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(_safe_float(tb.get("fontSize", 14), 14))
        p.font.color.rgb = _parse_rgb_color(tb.get("color", "")) or COLOR_BLACK
        p.font.bold = bool(tb.get("bold", False))
        p.font.italic = bool(tb.get("italic", False))
        p.alignment = ppt_align


def _build_figure_asset_lookup(figure_assets: dict, upload_dir: str) -> dict[str, str]:
    lookup: dict[str, str] = {}
    if not isinstance(figure_assets, dict):
        return lookup

    for label, asset in figure_assets.items():
        key = _normalize_figure_label(label)
        if not key:
            continue
        path = _resolve_uploaded_asset_path(str(asset or ""), upload_dir)
        if path and os.path.exists(path):
            lookup[key] = path
    return lookup


def _resolve_figure_placeholder_path(ph: dict, figure_asset_paths: dict[str, str], upload_dir: str = "") -> Optional[str]:
    for direct_key in ("asset", "url", "path"):
        direct_path = _resolve_uploaded_asset_path(str(ph.get(direct_key, "") or ""), upload_dir)
        if direct_path and os.path.exists(direct_path):
            return direct_path

    label = _normalize_figure_label(ph.get("figure") or ph.get("label", ""))
    if not label:
        return None
    return figure_asset_paths.get(label)


def _normalize_figure_label(value: str) -> str:
    text = str(value or "")
    match = re.search(r"figure\s*\d+(?:\.\d+)?", text, re.IGNORECASE)
    if match:
        return re.sub(r"\s+", " ", match.group(0)).strip().lower()
    return re.sub(r"\s+", " ", text).strip().lower()


def _resolve_uploaded_asset_path(value: str, upload_dir: str) -> str:
    raw = str(value or "").strip().strip('"').strip("'")
    if not raw:
        return ""
    if raw.startswith("data:"):
        return ""

    normalized = raw.replace("\\", "/")
    if normalized.startswith("http://") or normalized.startswith("https://"):
        marker = "/beamer-generator/uploads/"
        idx = normalized.find(marker)
        if idx != -1:
            normalized = normalized[idx + len(marker):]
        else:
            return ""

    if normalized.startswith("/beamer-generator/uploads/"):
        normalized = normalized.split("/beamer-generator/uploads/", 1)[1]
    elif normalized.startswith("/uploads/"):
        normalized = normalized.split("/uploads/", 1)[1]
    elif os.path.isabs(raw) and os.path.exists(raw):
        return raw

    normalized = normalized.lstrip("./")
    if not normalized or normalized.startswith(".."):
        return ""

    candidate = Path(upload_dir) / normalized if upload_dir else Path(normalized)
    return str(candidate)


def _add_picture_fit(slide, image_path: str, left, top, box_width, box_height):
    """Insert a picture and keep it contained inside the target box."""
    try:
        size = _image_size(image_path)
        if not size or not size[0] or not size[1]:
            pic = slide.shapes.add_picture(image_path, left, top, width=box_width)
            pic.left = left
            pic.top = top
            return pic

        img_w, img_h = size
        scale = min(float(box_width) / float(img_w), float(box_height) / float(img_h))
        width = Emu(int(img_w * scale))
        height = Emu(int(img_h * scale))
        pic = slide.shapes.add_picture(
            image_path,
            left + Emu(int((int(box_width) - int(width)) / 2)),
            top + Emu(int((int(box_height) - int(height)) / 2)),
            width=width,
            height=height,
        )
        return pic
    except Exception:
        return None


def _safe_float(value, default: float) -> float:
    try:
        return float(value)
    except (TypeError, ValueError):
        return default


def _editor_px_to_slide_x(value) -> Emu:
    return Emu(int(_safe_float(value, 0) / 860.0 * int(SLIDE_WIDTH)))


def _editor_px_to_slide_y(value) -> Emu:
    return Emu(int(_safe_float(value, 0) / 484.0 * int(SLIDE_HEIGHT)))


def _parse_rgb_color(value: str) -> Optional[RGBColor]:
    s = str(value or "").strip()
    if not s or s in {"transparent", "rgba(0, 0, 0, 0)"}:
        return None
    if s.startswith("#"):
        raw = s[1:]
        if len(raw) == 3:
            raw = "".join(ch * 2 for ch in raw)
        if len(raw) == 6:
            try:
                return RGBColor(int(raw[0:2], 16), int(raw[2:4], 16), int(raw[4:6], 16))
            except ValueError:
                return None
    match = re.match(r"rgba?\((\d+),\s*(\d+),\s*(\d+)(?:,\s*([0-9.]+))?\)", s)
    if match:
        if match.group(4) is not None and _safe_float(match.group(4), 1) <= 0:
            return None
        return RGBColor(
            max(0, min(255, int(match.group(1)))),
            max(0, min(255, int(match.group(2)))),
            max(0, min(255, int(match.group(3)))),
        )
    return None


def _render_equation_image(equation: str) -> Optional[io.BytesIO]:
    """将 LaTeX 公式渲染为透明 PNG，供 python-pptx 插入。"""
    lines = _split_equation_lines(equation)
    if not lines:
        return None

    fig_height = max(0.42 * len(lines) + 0.16, 0.52)
    fig = plt.figure(figsize=(10.4, fig_height), dpi=220)
    fig.patch.set_alpha(0)

    try:
      for idx, line in enumerate(lines):
          y = 1 - ((idx + 0.5) / len(lines))
          fig.text(
              0.5,
              y,
              f"${line}$",
              ha="center",
              va="center",
              fontsize=18,
              color="#2864B4",
          )

      out = io.BytesIO()
      fig.savefig(out, format="png", transparent=True, bbox_inches="tight", pad_inches=0.05)
      out.seek(0)
      return out
    except Exception:
      return _render_formula_source_image(equation)
    finally:
      plt.close(fig)


def _render_formula_source_image(equation: str) -> Optional[io.BytesIO]:
    """Fallback formula box when mathtext cannot render a source equation."""
    text = _formula_to_plain_text(equation)
    if not text:
        return None
    try:
        return _render_text_png(text, font_size=13, color="#2864B4", align="center")
    except Exception:
        return None


def _split_equation_lines(equation: str) -> list[str]:
    s = _repair_latex_artifacts(equation).strip()
    if not s:
        return []

    s = re.sub(r"\\begin\{(?:aligned|align|equation|gather|split)\*?\}", "", s)
    s = re.sub(r"\\end\{(?:aligned|align|equation|gather|split)\*?\}", "", s)
    s = re.sub(r"\\label\{[^}]*\}", "", s)
    s = re.sub(r"\\nonumber\b", "", s)
    s = re.sub(r"\\\\\[[^\]]*\]", r"\\\\", s)

    raw_lines = re.split(r"\\\\", s)
    result: list[str] = []
    for raw in raw_lines:
        line = raw.strip()
        if not line:
            continue
        tag = ""
        tag_match = re.search(r"\\tag\{([^}]*)\}", line)
        if tag_match:
            tag = tag_match.group(1).strip()
            line = re.sub(r"\\tag\{[^}]*\}", "", line)

        line = line.replace("&", "")
        line = _normalize_fraction_shorthand(line)
        line = _normalize_mathtext_unsupported(line)
        line = re.sub(r"\s+", " ", line).strip()
        if tag:
            line = f"{line}\\quad \\mathrm{{({tag})}}"
        if line:
            result.append(line)

    return result


def _normalize_fraction_shorthand(text: str) -> str:
    # LaTeX 允许 \frac18，matplotlib mathtext 更稳定地接受 \frac{1}{8}。
    return re.sub(r"\\frac\s*([A-Za-z0-9])\s*([A-Za-z0-9])", r"\\frac{\1}{\2}", text)


def _normalize_mathtext_unsupported(text: str) -> str:
    s = text
    s = s.replace(r"\,", r"\ ")
    s = s.replace(r"\;", r"\ ")
    s = s.replace(r"\!", "")
    s = s.replace(r"\quad", r"\ ")
    s = s.replace(r"\qquad", r"\ ")
    s = re.sub(r"\\text\{([^{}]*)\}", r"\\mathrm{\1}", s)
    s = re.sub(r"\\operatorname\{([^{}]*)\}", r"\\mathrm{\1}", s)
    s = re.sub(r"\\mathrm\{\\text\{([^{}]*)\}\}", r"\\mathrm{\1}", s)
    # mathtext 对 || 支持有限，按普通竖线显示即可。
    s = s.replace(r"\|", "|")
    return s


def _formula_to_plain_text(equation: str) -> str:
    return "  ".join(_split_equation_lines(equation)) or str(equation or "")

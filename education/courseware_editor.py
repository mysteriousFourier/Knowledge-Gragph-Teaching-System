"""Structured courseware editing helpers.

This module keeps the edit model independent from the legacy slide preview
payload.  The preview payload is still useful for lecture generation, while the
editable model is stable enough for canvas editing, project persistence, and
export.
"""
from __future__ import annotations

import base64
import hashlib
import html
import io
import json
import mimetypes
import posixpath
import re
import uuid
import zipfile
from datetime import datetime
from pathlib import Path
from typing import Any, Dict, Iterable, List, Optional

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches, Pt

from KGTS.core.bridge import RUNTIME_DIR
from KGTS.education.ppt_parser import IMAGE_EXTENSIONS, MAX_INLINE_IMAGE_BYTES
from KGTS.education.ppt_tex_generator import ARTIFACT_DIR


CANVAS_WIDTH = 1000.0
CANVAS_HEIGHT = 562.5
PROJECT_DIR = RUNTIME_DIR / "courseware" / "projects"


def build_editable_model(parse_result: Dict[str, Any], prompt_data: Optional[Dict[str, Any]] = None) -> Dict[str, Any]:
    """Build an editable slide model from parser output."""
    if prompt_data is None:
        from KGTS.education.ppt_parser import build_ppt_lecture_prompt_data

        prompt_data = build_ppt_lecture_prompt_data(parse_result)

    slides = prompt_data.get("slide_details") or []
    asset_map = build_asset_map(slides)
    editable_slides = [
        _editable_slide_from_detail(slide, asset_map)
        for slide in slides
        if isinstance(slide, dict)
    ]
    return {
        "version": 1,
        "title": prompt_data.get("chapter_title") or "未命名课件",
        "slide_count": len(editable_slides),
        "canvas": {"width": CANVAS_WIDTH, "height": CANVAS_HEIGHT, "unit": "px"},
        "layout": {"canvas": {"width": CANVAS_WIDTH, "height": CANVAS_HEIGHT}},
        "source_tex": parse_result.get("tex_content") or "",
        "source_tex_file": parse_result.get("tex_source_file") or "",
        "assets": asset_map,
        "slides": editable_slides,
        "updated_at": datetime.now().isoformat(),
    }


def build_editable_model_from_slide_details(
    slides: List[Dict[str, Any]],
    *,
    title: str = "未命名课件",
    source_tex: str = "",
    tex_source_file: str = "",
) -> Dict[str, Any]:
    asset_map = build_asset_map(slides)
    editable_slides = [
        _editable_slide_from_detail(slide, asset_map)
        for slide in slides
        if isinstance(slide, dict)
    ]
    return {
        "version": 1,
        "title": title,
        "slide_count": len(editable_slides),
        "canvas": {"width": CANVAS_WIDTH, "height": CANVAS_HEIGHT, "unit": "px"},
        "layout": {"canvas": {"width": CANVAS_WIDTH, "height": CANVAS_HEIGHT}},
        "source_tex": source_tex,
        "source_tex_file": tex_source_file,
        "assets": asset_map,
        "slides": editable_slides,
        "updated_at": datetime.now().isoformat(),
    }


def build_asset_map(slides: Iterable[Dict[str, Any]]) -> Dict[str, Dict[str, Any]]:
    assets: Dict[str, Dict[str, Any]] = {}
    slide_list = [slide for slide in slides if isinstance(slide, dict)]
    for slide in slide_list:
        for image in _collect_slide_images(slide):
            source_path = str(image.get("source_path") or image.get("tex_ref") or "").strip()
            data_uri = image.get("data_uri")
            key = source_path or str(image.get("tex_ref") or "") or str(data_uri or "")
            if not key:
                continue
            asset_id = _stable_id("asset", key)
            asset = assets.setdefault(
                asset_id,
                {
                    "id": asset_id,
                    "name": posixpath.basename(source_path) or source_path or asset_id,
                    "source_path": source_path,
                    "tex_ref": image.get("tex_ref") or source_path,
                    "mime_type": _mime_type_for_asset(source_path, data_uri),
                    "data_uri": data_uri,
                    "aliases": _asset_aliases(source_path, image.get("tex_ref")),
                    "slide_indices": [],
                    "figure_refs": [],
                    "oversized": bool(image.get("oversized")),
                },
            )
            if not asset.get("data_uri") and data_uri:
                asset["data_uri"] = data_uri
            if image.get("oversized"):
                asset["oversized"] = True
            index = slide.get("index")
            if index not in asset["slide_indices"]:
                asset["slide_indices"].append(index)
            for figure_ref in _figure_refs_for_slide(slide):
                if _figure_ref_matches_asset(figure_ref, source_path) or len(_collect_slide_images(slide)) == 1:
                    if figure_ref not in asset["figure_refs"]:
                        asset["figure_refs"].append(figure_ref)
    return assets


def assets_from_upload(file_bytes: bytes, filename: str) -> Dict[str, Dict[str, Any]]:
    lower_name = filename.lower()
    if lower_name.endswith(".zip"):
        return _assets_from_zip(file_bytes)
    if posixpath.splitext(lower_name)[1] not in IMAGE_EXTENSIONS:
        return {}
    return _asset_map_from_named_bytes([(filename, file_bytes)])


def serialize_editable_model_to_tex(model: Dict[str, Any], *, title: Optional[str] = None) -> str:
    source_tex = str(model.get("source_tex") or model.get("sourceTex") or "")
    slides = [slide for slide in model.get("slides") or [] if isinstance(slide, dict)]
    deck_title = title or str(model.get("title") or "KGTS 课件")
    rendered_frames = [_slide_model_to_tex(slide) for slide in slides]

    replacements: List[tuple[int, int, str]] = []
    unsourced_frames: List[str] = []
    for slide, frame_tex in zip(slides, rendered_frames):
        start = slide.get("source_start", slide.get("sourceStart"))
        end = slide.get("source_end", slide.get("sourceEnd"))
        if isinstance(start, int) and isinstance(end, int) and source_tex and 0 <= start < end <= len(source_tex):
            replacements.append((start, end, frame_tex))
        else:
            unsourced_frames.append(frame_tex)

    if source_tex and replacements:
        output = source_tex
        for start, end, frame_tex in sorted(replacements, key=lambda item: item[0], reverse=True):
            output = f"{output[:start]}{frame_tex}{output[end:]}"
        if unsourced_frames:
            insert = "\n\n".join(unsourced_frames)
            if "\\end{document}" in output:
                output = output.replace("\\end{document}", f"{insert}\n\\end{{document}}", 1)
            else:
                output = f"{output.rstrip()}\n\n{insert}\n"
        return output

    lines = [
        r"\documentclass[aspectratio=169]{beamer}",
        r"\usetheme{Madrid}",
        r"\usepackage{amsmath,amssymb}",
        r"\usepackage{graphicx}",
        r"\usepackage{ctex}",
        r"\title{" + _tex_escape(deck_title) + "}",
        r"\date{}",
        r"\begin{document}",
        r"\frame{\titlepage}",
        *rendered_frames,
        r"\end{document}",
    ]
    return "\n".join(lines)


def build_pptx_artifact_from_editable_model(
    title: str,
    model: Dict[str, Any],
    *,
    source_node_ids: Optional[List[str]] = None,
) -> Dict[str, Any]:
    ARTIFACT_DIR.mkdir(parents=True, exist_ok=True)
    slug = hashlib.md5(f"{title}:{datetime.now().isoformat()}".encode("utf-8")).hexdigest()[:12]
    pptx_path = ARTIFACT_DIR / f"courseware_{slug}.pptx"
    tex_path = ARTIFACT_DIR / f"courseware_{slug}.tex"
    tex_content = serialize_editable_model_to_tex(model, title=title)
    tex_path.write_text(tex_content, encoding="utf-8")

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]
    assets = model.get("assets") if isinstance(model.get("assets"), dict) else {}
    for slide_model in model.get("slides") or []:
        if not isinstance(slide_model, dict):
            continue
        ppt_slide = prs.slides.add_slide(blank_layout)
        for obj in sorted(_slide_objects(slide_model), key=lambda item: int(item.get("z", 0))):
            _add_object_to_pptx_slide(ppt_slide, obj, assets)
    prs.save(pptx_path)
    return {
        "kind": "editable_courseware",
        "pptx_path": str(pptx_path),
        "tex_path": str(tex_path),
        "pptx_url": f"/api/education/artifacts/{pptx_path.name}",
        "tex_url": f"/api/education/artifacts/{tex_path.name}",
        "tex_content_hash": hashlib.md5(tex_content.encode("utf-8")).hexdigest(),
        "slide_count": len(model.get("slides") or []),
        "source_node_ids": source_node_ids or [],
        "generated_at": datetime.now().isoformat(),
    }


def save_courseware_project(payload: Dict[str, Any]) -> Dict[str, Any]:
    PROJECT_DIR.mkdir(parents=True, exist_ok=True)
    project_id = _safe_project_id(str(payload.get("project_id") or payload.get("id") or "")) or f"cw_{uuid.uuid4().hex[:12]}"
    path = PROJECT_DIR / f"{project_id}.json"
    previous: Dict[str, Any] = {}
    if path.exists():
        try:
            previous = json.loads(path.read_text(encoding="utf-8"))
        except Exception:
            previous = {}
    now = datetime.now().isoformat()
    model = payload.get("editable_model") or {}
    tex_content = payload.get("tex_content") or (model.get("source_tex") if isinstance(model, dict) else "") or ""
    if isinstance(model, dict) and tex_content and not model.get("source_tex"):
        model = {**model, "source_tex": tex_content}
    record = {
        "id": project_id,
        "title": str(payload.get("title") or model.get("title") or "未命名课件"),
        "editable_model": model,
        "asset_map": payload.get("asset_map") or model.get("assets") or {},
        "slides": payload.get("slides") or [],
        "tex_content": tex_content,
        "ppt_artifact": payload.get("ppt_artifact"),
        "source_node_ids": payload.get("source_node_ids") or [],
        "created_at": previous.get("created_at") or now,
        "updated_at": now,
    }
    path.write_text(json.dumps(record, ensure_ascii=False, indent=2), encoding="utf-8")
    return record


def list_courseware_projects() -> List[Dict[str, Any]]:
    PROJECT_DIR.mkdir(parents=True, exist_ok=True)
    projects: List[Dict[str, Any]] = []
    for path in PROJECT_DIR.glob("*.json"):
        try:
            record = json.loads(path.read_text(encoding="utf-8"))
        except Exception:
            continue
        model = record.get("editable_model") if isinstance(record.get("editable_model"), dict) else {}
        projects.append(
            {
                "id": record.get("id") or path.stem,
                "title": record.get("title") or model.get("title") or path.stem,
                "slide_count": len(model.get("slides") or record.get("slides") or []),
                "created_at": record.get("created_at"),
                "updated_at": record.get("updated_at"),
            }
        )
    return sorted(projects, key=lambda item: str(item.get("updated_at") or ""), reverse=True)


def load_courseware_project(project_id: str) -> Optional[Dict[str, Any]]:
    safe_id = _safe_project_id(project_id)
    if not safe_id:
        return None
    path = PROJECT_DIR / f"{safe_id}.json"
    if not path.exists() or not path.is_file():
        return None
    return json.loads(path.read_text(encoding="utf-8"))


def _editable_slide_from_detail(slide: Dict[str, Any], asset_map: Dict[str, Dict[str, Any]]) -> Dict[str, Any]:
    index = int(slide.get("index") or 1)
    content = str(slide.get("content") or "\n".join(slide.get("body_texts") or []) or "")
    objects: List[Dict[str, Any]] = []
    canvas_items = _canvas_items_by_kind(slide)
    z = 1
    content_canvas_ordinal = 1

    title = str(slide.get("title") or "").strip()
    cursor_y = 34.0
    if title:
        title_bbox = _bbox_from_canvas(canvas_items, "title") or _default_bbox("title", 0, slide)
        objects.append(
            _object_payload(
                slide_index=index,
                kind="title",
                ordinal=1,
                bbox=title_bbox,
                z=z,
                text=title,
                rich_html=_rich_html_from_text(title),
                role="title",
                style={"fontSize": 28, "lineHeight": 1.12},
            )
        )
        cursor_y = max(cursor_y, title_bbox["y"] + title_bbox["height"] + 22)
        z += 1

    body_without_display_math = _remove_display_equations(content).strip()
    if body_without_display_math:
        body_canvas_bbox = _bbox_from_canvas(canvas_items, "content", ordinal=content_canvas_ordinal)
        body_bbox = body_canvas_bbox or _default_bbox("richText", 0, slide)
        if body_canvas_bbox:
            body_bbox["y"] = body_bbox["y"]
        elif not _layout_prefers_image_first(slide):
            body_bbox["y"] = cursor_y
        estimated_height = _estimate_text_bbox_height(body_without_display_math, 18, body_bbox["width"])
        body_bbox["height"] = (
            max(body_bbox["height"], estimated_height)
            if body_canvas_bbox
            else _fit_bbox_height(body_bbox["y"], estimated_height, min_height=72)
        )
        objects.append(
            _object_payload(
                slide_index=index,
                kind="richText",
                ordinal=1,
                bbox=body_bbox,
                z=z,
                text=body_without_display_math,
                rich_html=_rich_html_from_markdown(body_without_display_math),
                role="body",
                style={"fontSize": 18, "lineHeight": 1.32},
            )
        )
        cursor_y = max(cursor_y, body_bbox["y"] + body_bbox["height"] + 20)
        content_canvas_ordinal += 1
        z += 1

    for ordinal, equation in enumerate(_extract_equations(content), start=1):
        equation_canvas_bbox = _bbox_from_canvas(canvas_items, "content", ordinal=content_canvas_ordinal)
        equation_bbox = equation_canvas_bbox or _default_bbox("equation", ordinal - 1, slide)
        equation_bbox["y"] = equation_bbox["y"] if equation_canvas_bbox else cursor_y
        estimated_height = _estimate_text_bbox_height(equation, 24, equation_bbox["width"], line_height=1.25, min_height=54)
        equation_bbox["height"] = (
            max(equation_bbox["height"], estimated_height)
            if equation_canvas_bbox
            else _fit_bbox_height(equation_bbox["y"], estimated_height, min_height=54)
        )
        objects.append(
            _object_payload(
                slide_index=index,
                kind="equation",
                ordinal=ordinal,
                bbox=equation_bbox,
                z=z,
                latex=equation,
                text=equation,
                style={"fontSize": 24, "lineHeight": 1.25},
            )
        )
        cursor_y = max(cursor_y, equation_bbox["y"] + equation_bbox["height"] + 20)
        content_canvas_ordinal += 1
        z += 1

    for ordinal, table in enumerate(slide.get("tables") or [], start=1):
        rows = table.get("rows") if isinstance(table, dict) else None
        if rows:
            table_canvas_bbox = _bbox_from_canvas(canvas_items, "content", ordinal=content_canvas_ordinal)
            table_bbox = table_canvas_bbox or _default_bbox("table", ordinal - 1, slide)
            table_bbox["y"] = table_bbox["y"] if table_canvas_bbox else cursor_y
            estimated_height = max(72, len(rows) * 27 + 24)
            table_bbox["height"] = (
                max(table_bbox["height"], estimated_height)
                if table_canvas_bbox
                else _fit_bbox_height(table_bbox["y"], estimated_height, min_height=72)
            )
            objects.append(
                _object_payload(
                    slide_index=index,
                    kind="table",
                    ordinal=ordinal,
                    bbox=table_bbox,
                    z=z,
                    rows=rows,
                    style={"fontSize": 14, "lineHeight": 1.35},
                )
            )
            cursor_y = max(cursor_y, table_bbox["y"] + table_bbox["height"] + 20)
            content_canvas_ordinal += 1
            z += 1

    for ordinal, image in enumerate(_collect_slide_images(slide), start=1):
        asset_id = _asset_id_for_image(image, asset_map)
        kind = "image" if asset_id and asset_map.get(asset_id, {}).get("data_uri") else "placeholder"
        objects.append(
            _object_payload(
                slide_index=index,
                kind=kind,
                ordinal=ordinal,
                bbox=_bbox_from_canvas(canvas_items, "image", image=image, ordinal=ordinal) or _default_image_bbox(slide, ordinal - 1, cursor_y),
                z=z,
                asset_id=asset_id,
                source_path=image.get("source_path") or image.get("tex_ref") or "",
                tex_ref=image.get("tex_ref") or image.get("source_path") or "",
                width_ratio=image.get("width_ratio"),
                label=posixpath.basename(str(image.get("source_path") or image.get("tex_ref") or "")) or f"图片 {ordinal}",
            )
        )
        z += 1

    for ordinal, callout in enumerate(_extract_callouts(str(slide.get("source_tex") or slide.get("source_body_tex") or "")), start=1):
        callout_bbox = callout.get("bbox") or _default_bbox("callout", ordinal - 1, slide)
        estimated_callout_height = _estimate_text_bbox_height(callout.get("text") or "", 16, callout_bbox["width"], min_height=72) + 18
        callout_bbox["height"] = (
            max(callout_bbox["height"], min(estimated_callout_height, CANVAS_HEIGHT - callout_bbox["y"] - 24))
            if callout.get("bbox")
            else _fit_bbox_height(callout_bbox["y"], estimated_callout_height, min_height=80)
        )
        callout_bbox = _avoid_bbox_overlap(callout_bbox, objects)
        objects.append(
            _object_payload(
                slide_index=index,
                kind="callout",
                ordinal=ordinal,
                bbox=callout_bbox,
                z=z,
                text=callout.get("text") or "",
                title=callout.get("title") or "提示",
                rich_html=_rich_html_from_text(callout.get("text") or ""),
                style={"fontSize": 16, "lineHeight": 1.32},
            )
        )
        z += 1

    if not objects:
        objects.append(
            _object_payload(
                slide_index=index,
                kind="placeholder",
                ordinal=1,
                bbox=_default_bbox("placeholder", 0, slide),
                z=1,
                text="空白页",
                label="空白页",
            )
        )

    return {
        "id": f"slide-{index}",
        "index": index,
        "title": title,
        "items": objects,
        "objects": objects,
        "layout": slide.get("layout") or {},
        "source_tex": slide.get("source_tex") or "",
        "source_body_tex": slide.get("source_body_tex") or "",
        "source_start": slide.get("source_start"),
        "source_end": slide.get("source_end"),
        "notes": slide.get("notes") or "",
    }


def _pack_slide_objects(objects: List[Dict[str, Any]]) -> List[Dict[str, Any]]:
    if not objects:
        return objects

    image_objects = [obj for obj in objects if obj.get("type") in {"image", "placeholder"}]
    text_objects = [obj for obj in objects if obj.get("type") not in {"title", "image", "placeholder"}]
    has_images = bool(image_objects)
    if not has_images and not text_objects:
        return objects

    title_bottom = 92.0
    for obj in objects:
        if obj.get("type") == "title":
            bbox = obj.get("bbox") if isinstance(obj.get("bbox"), dict) else {}
            bbox["x"] = _finite_float(bbox.get("x"), 48)
            bbox["y"] = _finite_float(bbox.get("y"), 34)
            bbox["width"] = min(max(_finite_float(bbox.get("width"), 904), 260), CANVAS_WIDTH - bbox["x"] - 24)
            bbox["height"] = min(max(_finite_float(bbox.get("height"), 58), 42), 72)
            obj["bbox"] = bbox
            title_bottom = max(title_bottom, bbox["y"] + bbox["height"])

    content_x = 64.0 if has_images else 72.0
    content_width = 500.0 if has_images else 856.0
    cursor_y = max(112.0, title_bottom + 20.0)
    gap = 16.0
    for obj in sorted(text_objects, key=lambda item: int(item.get("z", 0))):
        bbox = obj.get("bbox") if isinstance(obj.get("bbox"), dict) else {}
        bbox["x"] = content_x
        bbox["y"] = cursor_y
        bbox["width"] = content_width
        estimated_height = _object_intrinsic_height(obj, content_width)
        bbox["height"] = _fit_bbox_height(cursor_y, estimated_height, min_height=_min_height_for_object(obj))
        obj["bbox"] = bbox
        cursor_y = bbox["y"] + bbox["height"] + gap

    image_width = 380.0 if text_objects else 460.0
    image_height = 118.0 if len(image_objects) > 1 else 240.0
    image_x = 570.0 if text_objects else (CANVAS_WIDTH - image_width) / 2
    for index, obj in enumerate(sorted(image_objects, key=lambda item: int(item.get("z", 0)))):
        bbox = obj.get("bbox") if isinstance(obj.get("bbox"), dict) else {}
        bbox["x"] = image_x
        bbox["y"] = min(max(126.0 + index * (132.0 if len(image_objects) > 1 else 156.0), 86.0), CANVAS_HEIGHT - image_height - 24.0)
        bbox["width"] = image_width
        bbox["height"] = image_height
        obj["bbox"] = bbox

    return objects


def _object_intrinsic_height(obj: Dict[str, Any], width: float) -> float:
    kind = str(obj.get("type") or "")
    style = obj.get("style") if isinstance(obj.get("style"), dict) else {}
    font_size = int(_finite_float(style.get("fontSize"), 18))
    line_height = _finite_float(style.get("lineHeight"), 1.32)
    if kind == "equation":
        return _estimate_text_bbox_height(str(obj.get("latex") or obj.get("text") or ""), font_size, width, line_height=line_height, min_height=54) + 8
    if kind == "table":
        rows = obj.get("rows") if isinstance(obj.get("rows"), list) else []
        return max(72.0, len(rows or [None]) * font_size * line_height + 34.0)
    if kind == "callout":
        return _estimate_text_bbox_height(str(obj.get("text") or ""), font_size, width, line_height=line_height, min_height=80) + 18
    return _estimate_text_bbox_height(str(obj.get("text") or obj.get("label") or ""), font_size, width, line_height=line_height, min_height=48)


def _min_height_for_object(obj: Dict[str, Any]) -> float:
    kind = str(obj.get("type") or "")
    if kind == "equation":
        return 54.0
    if kind == "table":
        return 72.0
    if kind == "callout":
        return 80.0
    return 48.0


def _collect_slide_images(slide: Dict[str, Any]) -> List[Dict[str, Any]]:
    images: List[Dict[str, Any]] = []
    seen: set[str] = set()

    def add(image: Any) -> None:
        if not isinstance(image, dict):
            return
        key = str(image.get("source_path") or image.get("tex_ref") or image.get("data_uri") or len(images))
        if key in seen:
            return
        seen.add(key)
        images.append(image)

    for image in slide.get("images") or []:
        add(image)
    layout = slide.get("layout") if isinstance(slide.get("layout"), dict) else {}
    for column in layout.get("columns") or []:
        if isinstance(column, dict):
            for image in column.get("images") or []:
                add(image)
    return images


def _object_payload(slide_index: int, kind: str, ordinal: int, bbox: Dict[str, float], z: int, **extra: Any) -> Dict[str, Any]:
    payload = {
        "id": _stable_id(f"s{slide_index}-{kind}", str(extra.get("source_path") or extra.get("text") or extra.get("latex") or ordinal)),
        "type": kind,
        "bbox": bbox,
        "z": z,
        "locked": False,
        "style": extra.pop("style", {}),
    }
    payload.update({key: value for key, value in extra.items() if value is not None})
    return payload


def _default_bbox(kind: str, ordinal: int, slide: Dict[str, Any]) -> Dict[str, float]:
    image_count = len(_collect_slide_images(slide))
    has_body = bool(str(slide.get("content") or "").strip())
    image_first = _layout_prefers_image_first(slide)
    layout = slide.get("layout") if isinstance(slide.get("layout"), dict) else {}
    is_title_layout = str(layout.get("mode") or "") == "title"
    if kind == "title":
        if is_title_layout:
            return {"x": 96, "y": 152, "width": 808, "height": 112}
        return {"x": 48, "y": 34, "width": 904, "height": 58}
    if kind == "richText":
        if is_title_layout:
            return {"x": 180, "y": 284, "width": 640, "height": 128}
        if image_count and image_first:
            y = _image_stack_bottom(slide) + 18
            return {"x": 72, "y": y, "width": 856, "height": max(72, CANVAS_HEIGHT - y - 24)}
        return {"x": 64, "y": 116, "width": 500 if image_count else 872, "height": 96}
    if kind in {"image", "placeholder"}:
        width = 420 if image_count > 1 else 460
        if image_count and image_first:
            width = 700 if image_count <= 1 else 420
            return {"x": (CANVAS_WIDTH - width) / 2, "y": 120 + ordinal * 132, "width": width, "height": 210 if image_count <= 1 else 118}
        x = 570 if has_body else (CANVAS_WIDTH - width) / 2
        return {"x": x, "y": 126 + ordinal * 128, "width": width, "height": 240 if image_count <= 1 else 118}
    if kind == "equation":
        if image_count and image_first:
            return {"x": 96, "y": 392 + ordinal * 70, "width": 808, "height": 60}
        return {"x": 72 if image_count else 96, "y": 212 + ordinal * 82, "width": 500 if image_count else 808, "height": 68}
    if kind == "table":
        return {"x": 72 if image_count else 84, "y": 250 + ordinal * 36, "width": 500 if image_count else 832, "height": 180}
    if kind == "callout":
        return {"x": 620, "y": 150 + ordinal * 120, "width": 304, "height": 112}
    return {"x": 96, "y": 160 + ordinal * 72, "width": 808, "height": 96}


def _default_image_bbox(slide: Dict[str, Any], ordinal: int, cursor_y: float) -> Dict[str, float]:
    image_count = len(_collect_slide_images(slide))
    content = str(slide.get("content") or "").strip()
    has_body = bool(_remove_display_equations(content).strip())
    layout = slide.get("layout") if isinstance(slide.get("layout"), dict) else {}
    image_first = _layout_prefers_image_first(slide)
    width = 420 if image_count > 1 else 460
    height = 126 if image_count > 1 else 240
    if image_count and image_first:
        width = min(max(_finite_float((slide.get("layout") or {}).get("max_image_width"), 0.7), 0.2), 0.95) * CANVAS_WIDTH
        height = 210 if image_count <= 1 else 118
        x = (CANVAS_WIDTH - width) / 2
        y = 120 + ordinal * 132
    elif has_body:
        x = 570
        y = 126 + ordinal * 132
    else:
        x = (CANVAS_WIDTH - width) / 2
        y = max(126 + ordinal * 132, min(cursor_y, CANVAS_HEIGHT - height - 28))
    y = min(max(y, 86), CANVAS_HEIGHT - height - 24)
    return {"x": x, "y": y, "width": width, "height": height}


def _layout_prefers_image_first(slide: Dict[str, Any]) -> bool:
    layout = slide.get("layout") if isinstance(slide.get("layout"), dict) else {}
    return bool(layout.get("image_first")) or str(layout.get("mode") or "") in {"image_text", "image_only", "title"}


def _image_stack_bottom(slide: Dict[str, Any]) -> float:
    image_count = len(_collect_slide_images(slide))
    if image_count <= 0:
        return 112.0
    height = 210.0 if image_count <= 1 else 118.0
    return 120.0 + (image_count - 1) * 132.0 + height


def _avoid_bbox_overlap(bbox: Dict[str, float], objects: List[Dict[str, Any]]) -> Dict[str, float]:
    next_bbox = dict(bbox)
    if not _bbox_overlaps_objects(next_bbox, objects):
        return next_bbox
    original_x = next_bbox["x"]
    original_y = next_bbox["y"]
    bottom_margin = 12.0
    candidates: List[Dict[str, float]] = []
    object_bottoms = [
        _finite_float((obj.get("bbox") or {}).get("y"), 0) + _finite_float((obj.get("bbox") or {}).get("height"), 0) + 18.0
        for obj in objects
        if obj.get("type") != "title" and isinstance(obj.get("bbox"), dict)
    ]
    y_values = [original_y, 86.0, 116.0, 150.0, 220.0, 300.0, 390.0, *object_bottoms]
    x_values = [original_x, 24.0, 48.0, 620.0, CANVAS_WIDTH - next_bbox["width"] - 48.0]
    for y in y_values:
        for x in x_values:
            candidates.append(
                {
                    **next_bbox,
                    "x": min(max(x, 24.0), CANVAS_WIDTH - next_bbox["width"] - 24.0),
                    "y": min(max(y, 86.0), CANVAS_HEIGHT - next_bbox["height"] - bottom_margin),
                }
            )
    for candidate in sorted(candidates, key=lambda item: abs(item["x"] - original_x) + abs(item["y"] - original_y)):
        if not _bbox_overlaps_objects(candidate, objects):
            return candidate
    return next_bbox


def _bbox_overlaps_objects(bbox: Dict[str, float], objects: List[Dict[str, Any]]) -> bool:
    for obj in objects:
        if obj.get("type") == "title":
            continue
        other = obj.get("bbox") if isinstance(obj.get("bbox"), dict) else {}
        if _bbox_overlap_area(bbox, other) > 800.0:
            return True
    return False


def _bbox_overlap_area(a: Dict[str, Any], b: Dict[str, Any]) -> float:
    ax1 = _finite_float(a.get("x"), 0)
    ay1 = _finite_float(a.get("y"), 0)
    ax2 = ax1 + _finite_float(a.get("width"), 0)
    ay2 = ay1 + _finite_float(a.get("height"), 0)
    bx1 = _finite_float(b.get("x"), 0)
    by1 = _finite_float(b.get("y"), 0)
    bx2 = bx1 + _finite_float(b.get("width"), 0)
    by2 = by1 + _finite_float(b.get("height"), 0)
    return max(0.0, min(ax2, bx2) - max(ax1, bx1)) * max(0.0, min(ay2, by2) - max(ay1, by1))


def _estimate_text_bbox_height(
    text: str,
    font_size: int,
    width: float,
    *,
    line_height: float = 1.32,
    min_height: float = 48,
) -> float:
    clean = str(text or "").strip()
    if not clean:
        return min_height
    usable_width = max(float(width or 320) - 16, 80)
    chars_per_line = max(int(usable_width / max(font_size * 0.58, 7)), 8)
    line_count = 0
    for raw_line in clean.splitlines():
        line = re.sub(r"^\s*(?:[-*•]|\\item)\s*", "", raw_line).strip()
        line_count += max(1, (len(line) + chars_per_line - 1) // chars_per_line)
    padding = 16
    return max(min_height, line_count * font_size * line_height + padding)


def _fit_bbox_height(y: float, height: float, *, min_height: float = 36) -> float:
    available = max(min_height, CANVAS_HEIGHT - float(y or 0) - 24)
    return min(max(float(height or min_height), min_height), available)


def _canvas_items_by_kind(slide: Dict[str, Any]) -> Dict[str, List[Dict[str, Any]]]:
    layout = slide.get("layout") if isinstance(slide.get("layout"), dict) else {}
    canvas = layout.get("canvas") if isinstance(layout.get("canvas"), dict) else {}
    result: Dict[str, List[Dict[str, Any]]] = {}
    for item in canvas.get("items") or []:
        if not isinstance(item, dict):
            continue
        kind = str(item.get("type") or "")
        result.setdefault(kind, []).append(item)
    return result


def _bbox_from_canvas(
    canvas_items: Dict[str, List[Dict[str, Any]]],
    kind: str,
    *,
    image: Optional[Dict[str, Any]] = None,
    ordinal: int = 1,
) -> Optional[Dict[str, float]]:
    candidates = canvas_items.get(kind) or []
    if kind == "image" and image:
        refs = set(_asset_aliases(str(image.get("source_path") or ""), image.get("tex_ref")))
        for item in candidates:
            ref = str(item.get("ref") or "")
            if ref and ref in refs:
                return _bbox_from_canvas_item(item)
    if len(candidates) >= ordinal:
        return _bbox_from_canvas_item(candidates[ordinal - 1])
    if candidates and ordinal <= 1:
        return _bbox_from_canvas_item(candidates[0])
    return None


def _bbox_from_canvas_item(item: Dict[str, Any]) -> Dict[str, float]:
    return {
        "x": _finite_float(item.get("x"), 40),
        "y": _finite_float(item.get("y"), 40),
        "width": max(_finite_float(item.get("width"), 240), 24),
        "height": max(_finite_float(item.get("height"), 80), 20),
    }


def _finite_float(value: Any, fallback: float) -> float:
    try:
        number = float(value)
    except (TypeError, ValueError):
        return fallback
    return number if number == number and abs(number) < 1_000_000 else fallback


def _asset_id_for_image(image: Dict[str, Any], asset_map: Dict[str, Dict[str, Any]]) -> str:
    aliases = set(_asset_aliases(str(image.get("source_path") or ""), image.get("tex_ref")))
    for asset_id, asset in asset_map.items():
        asset_aliases = set(asset.get("aliases") or [])
        if aliases.intersection(asset_aliases):
            return asset_id
    key = str(image.get("source_path") or image.get("tex_ref") or image.get("data_uri") or "")
    return _stable_id("asset", key) if key else ""


def _assets_from_zip(file_bytes: bytes) -> Dict[str, Dict[str, Any]]:
    try:
        archive = zipfile.ZipFile(io.BytesIO(file_bytes))
    except zipfile.BadZipFile:
        return {}
    try:
        entries: List[tuple[str, bytes]] = []
        for name in archive.namelist():
            if name.endswith("/") or name.startswith("__MACOSX/"):
                continue
            if posixpath.splitext(name)[1].lower() in IMAGE_EXTENSIONS:
                entries.append((name, archive.read(name)))
        return _asset_map_from_named_bytes(entries)
    finally:
        archive.close()


def _asset_map_from_named_bytes(entries: List[tuple[str, bytes]]) -> Dict[str, Dict[str, Any]]:
    assets: Dict[str, Dict[str, Any]] = {}
    for name, data in entries:
        source_path = name.replace("\\", "/").lstrip("/")
        asset_id = _stable_id("asset", source_path)
        mime_type = mimetypes.guess_type(source_path)[0] or "application/octet-stream"
        data_uri = None
        oversized = len(data) > MAX_INLINE_IMAGE_BYTES
        if not oversized:
            data_uri = f"data:{mime_type};base64,{base64.b64encode(data).decode('utf-8')}"
        assets[asset_id] = {
            "id": asset_id,
            "name": posixpath.basename(source_path),
            "source_path": source_path,
            "tex_ref": source_path,
            "mime_type": mime_type,
            "data_uri": data_uri,
            "aliases": _asset_aliases(source_path, source_path),
            "slide_indices": [],
            "figure_refs": [],
            "oversized": oversized,
        }
    return assets


def _mime_type_for_asset(source_path: str, data_uri: Any) -> str:
    if isinstance(data_uri, str) and data_uri.startswith("data:") and ";" in data_uri:
        return data_uri[5:data_uri.index(";")] or "application/octet-stream"
    return mimetypes.guess_type(source_path)[0] or "application/octet-stream"


def _asset_aliases(source_path: str, tex_ref: Any = None) -> List[str]:
    candidates = [source_path, str(tex_ref or "")]
    aliases: List[str] = []
    for value in candidates:
        clean = value.strip().replace("\\", "/").lstrip("./")
        if not clean:
            continue
        basename = posixpath.basename(clean)
        root, extension = posixpath.splitext(clean)
        base_root, _base_ext = posixpath.splitext(basename)
        for item in [clean, clean.lower(), root, root.lower(), basename, basename.lower(), base_root, base_root.lower()]:
            if item and item not in aliases:
                aliases.append(item)
        if not extension:
            for ext in IMAGE_EXTENSIONS:
                item = f"{clean}{ext}"
                if item not in aliases:
                    aliases.append(item)
    return aliases


def _figure_refs_for_slide(slide: Dict[str, Any]) -> List[str]:
    text = "\n".join(
        str(part or "")
        for part in [slide.get("title"), slide.get("content"), slide.get("raw_text"), slide.get("source_tex")]
    )
    refs: List[str] = []
    for match in re.finditer(r"\b(?:Figure|Fig\.?)\s*([0-9]+(?:\.[0-9]+)*[A-Za-z]?)|\b图\s*([0-9]+(?:\.[0-9]+)*[A-Za-z]?)", text, flags=re.I):
        value = (match.group(1) or match.group(2) or "").strip()
        if value and value not in refs:
            refs.append(value)
    return refs


def _figure_ref_matches_asset(figure_ref: str, source_path: str) -> bool:
    normalized_ref = figure_ref.replace(".", "")
    normalized_path = re.sub(r"[^0-9A-Za-z]+", "", source_path)
    return bool(normalized_ref and normalized_ref in normalized_path)


def _extract_equations(content: str) -> List[str]:
    equations: List[str] = []
    for match in re.finditer(r"\$\$\s*([\s\S]*?)\s*\$\$", content or ""):
        value = match.group(1).strip()
        if value and value not in equations:
            equations.append(value)
    return equations


def _remove_display_equations(content: str) -> str:
    return re.sub(r"\n?\$\$\s*[\s\S]*?\s*\$\$\n?", "\n", content or "").strip()


def _extract_callouts(source: str) -> List[Dict[str, Any]]:
    callouts: List[Dict[str, Any]] = []
    for match in re.finditer(r"\\begin\{(alertblock|exampleblock|block)\}\{([^}]*)\}([\s\S]*?)\\end\{\1\}", source or ""):
        text = _plain_latex_to_text(match.group(3))
        if text:
            callouts.append({"title": _plain_latex_to_text(match.group(2)) or "提示", "text": text})
    for match in re.finditer(r"\\node\[([^\]]*callout[^\]]*)\]([^{]*)\{([\s\S]*?)\}\s*;", source or "", flags=re.I):
        text = _plain_latex_to_text(match.group(3))
        if text:
            callout: Dict[str, Any] = {"title": "标注", "text": text}
            bbox = _tikz_callout_bbox(match.group(1), match.group(2), text)
            if bbox:
                callout["bbox"] = bbox
            callouts.append(callout)
    return callouts


def _tikz_callout_bbox(options: str, placement: str, text: str) -> Optional[Dict[str, float]]:
    placement_text = str(placement or "")
    coord_match = re.search(
        r"at\s*\(\s*\[\s*xshift\s*=\s*([-+]?\d+(?:\.\d+)?)cm\s*,\s*yshift\s*=\s*([-+]?\d+(?:\.\d+)?)cm\s*\]\s*current page\.north west\s*\)",
        placement_text,
        flags=re.I,
    )
    if not coord_match:
        return None
    x_cm = _finite_float(coord_match.group(1), 0)
    y_cm = _finite_float(coord_match.group(2), 0)
    width_cm = 3.2
    width_match = re.search(r"text width\s*=\s*([-+]?\d+(?:\.\d+)?)cm", str(options or ""), flags=re.I)
    if width_match:
        width_cm = max(_finite_float(width_match.group(1), width_cm), 1.2)
    px_per_cm_x = CANVAS_WIDTH / 16.0
    px_per_cm_y = CANVAS_HEIGHT / 9.0
    width = min(max(width_cm * px_per_cm_x, 120.0), 360.0)
    height = _estimate_text_bbox_height(text, 16, width, min_height=72) + 16
    x = x_cm * px_per_cm_x - width / 2
    y = abs(y_cm) * px_per_cm_y - height / 2
    return {
        "x": min(max(x, 24.0), CANVAS_WIDTH - width - 24.0),
        "y": min(max(y, 86.0), CANVAS_HEIGHT - height - 24.0),
        "width": width,
        "height": min(max(height, 72.0), CANVAS_HEIGHT - 110.0),
    }


def _rich_html_from_markdown(value: str) -> str:
    lines = [line.strip() for line in (value or "").splitlines() if line.strip()]
    if not lines:
        return ""
    if sum(1 for line in lines if line.startswith("- ")) >= max(1, len(lines) // 2):
        items = "".join(f"<li>{html.escape(line[2:].strip())}</li>" for line in lines)
        return f"<ul>{items}</ul>"
    return "".join(f"<p>{html.escape(line)}</p>" for line in lines)


def _rich_html_from_text(value: str) -> str:
    return f"<p>{html.escape(value or '')}</p>" if value else ""


def _plain_latex_to_text(value: str) -> str:
    text = re.sub(r"\\item(?:<[^>]*>)?(?:\[[^\]]*\])?\s*", "- ", value or "")
    text = re.sub(r"\\(?:textbf|textit|emph|alert|textcolor)(?:\[[^\]]*\])?\{([^{}]*)\}", r"\1", text)
    text = re.sub(r"\\[A-Za-z@]+(?:\*|<[^>]*>)?(?:\[[^\]]*\])?", " ", text)
    text = text.replace("\\\\", "\n")
    text = text.replace(r"\%", "%").replace(r"\&", "&").replace(r"\_", "_")
    return re.sub(r"\s+", " ", text).strip(" {};")


def _slide_objects(slide: Dict[str, Any]) -> List[Dict[str, Any]]:
    objects = slide.get("objects") or slide.get("items") or []
    return [obj for obj in objects if isinstance(obj, dict)]


def _slide_model_to_tex(slide: Dict[str, Any]) -> str:
    title = _tex_escape(str(slide.get("title") or _first_title_object_text(slide) or f"第 {slide.get('index', '')} 页"))
    layout_items = []
    for obj in _slide_objects(slide):
        bbox = obj.get("bbox") if isinstance(obj.get("bbox"), dict) else {}
        layout_items.append(
            {
                "id": obj.get("id"),
                "type": _canvas_type_for_object(obj),
                "ref": obj.get("tex_ref") or obj.get("source_path") or obj.get("asset_id"),
                "x": round(_finite_float(bbox.get("x"), 0), 1),
                "y": round(_finite_float(bbox.get("y"), 0), 1),
                "width": round(_finite_float(bbox.get("width"), 100), 1),
                "height": round(_finite_float(bbox.get("height"), 40), 1),
            }
        )
    metadata = json.dumps({"items": layout_items}, ensure_ascii=False, separators=(",", ":"))
    lines = [r"\begin{frame}{" + title + "}", f"% KGTS_LAYOUT {metadata}"]
    for obj in sorted(_slide_objects(slide), key=lambda item: int(item.get("z", 0))):
        obj_type = str(obj.get("type") or "")
        if obj_type == "title":
            continue
        rendered = _object_to_tex(obj)
        if rendered:
            lines.append(rendered)
    notes = str(slide.get("notes") or "").strip()
    if notes:
        lines.append(r"\note{" + _tex_escape(notes) + "}")
    lines.append(r"\end{frame}")
    return "\n".join(lines)


def _first_title_object_text(slide: Dict[str, Any]) -> str:
    for obj in _slide_objects(slide):
        if obj.get("type") == "title" and obj.get("text"):
            return str(obj.get("text"))
    return ""


def _canvas_type_for_object(obj: Dict[str, Any]) -> str:
    obj_type = str(obj.get("type") or "")
    if obj_type in {"image", "placeholder"}:
        return "image"
    if obj_type == "title":
        return "title"
    return "content"


def _object_to_tex(obj: Dict[str, Any]) -> str:
    obj_type = str(obj.get("type") or "")
    if obj_type in {"richText", "textbox"}:
        return _text_to_tex(str(obj.get("text") or _html_to_text(str(obj.get("rich_html") or ""))))
    if obj_type == "callout":
        title = _tex_escape(str(obj.get("title") or "提示"))
        body = _text_to_tex(str(obj.get("text") or _html_to_text(str(obj.get("rich_html") or ""))), allow_itemize=False)
        return f"\\begin{{block}}{{{title}}}\n{body}\n\\end{{block}}"
    if obj_type == "equation":
        latex = str(obj.get("latex") or obj.get("text") or "").strip()
        return f"\\[\n{latex}\n\\]" if latex else ""
    if obj_type == "table":
        return _table_to_tex(obj.get("rows") or [])
    if obj_type in {"image", "placeholder"}:
        ref = str(obj.get("tex_ref") or obj.get("source_path") or "").strip()
        if not ref:
            label = _tex_escape(str(obj.get("label") or "图片占位符"))
            return r"\begin{center}\fbox{\parbox{0.6\textwidth}{" + label + r"}}\end{center}"
        bbox = obj.get("bbox") if isinstance(obj.get("bbox"), dict) else {}
        ratio = _finite_float(obj.get("width_ratio"), _finite_float(bbox.get("width"), 560) / CANVAS_WIDTH)
        ratio = min(max(ratio, 0.08), 1.0)
        return f"\\begin{{center}}\n\\includegraphics[width={ratio:.2f}\\textwidth]{{{ref}}}\n\\end{{center}}"
    return ""


def _text_to_tex(value: str, *, allow_itemize: bool = True) -> str:
    lines = [line.strip() for line in (value or "").splitlines() if line.strip()]
    if not lines:
        return ""
    bullet_lines = [re.sub(r"^\s*[-*•]\s*", "", line).strip() for line in lines]
    looks_like_list = allow_itemize and (len(lines) > 1 or any(line.startswith(("- ", "* ", "•")) for line in lines))
    if looks_like_list:
        body = "\n".join(r"\item " + _tex_escape(line) for line in bullet_lines if line)
        return "\\begin{itemize}\n" + body + "\n\\end{itemize}"
    return "\n\n".join(_tex_escape(line) for line in lines)


def _table_to_tex(rows: Any) -> str:
    if not isinstance(rows, list) or not rows:
        return ""
    column_count = max((len(row) for row in rows if isinstance(row, list)), default=0)
    if column_count <= 0:
        return ""
    spec = "|" + "|".join(["l"] * column_count) + "|"
    lines = [r"\begin{center}", rf"\begin{{tabular}}{{{spec}}}", r"\hline"]
    for row in rows:
        cells = row if isinstance(row, list) else [row]
        padded = [*cells, *([""] * (column_count - len(cells)))]
        lines.append(" & ".join(_tex_escape(str(cell)) for cell in padded[:column_count]) + r" \\ \hline")
    lines.extend([r"\end{tabular}", r"\end{center}"])
    return "\n".join(lines)


def _html_to_text(value: str) -> str:
    text = re.sub(r"</(?:p|li|div|tr|h[1-6])>", "\n", value or "", flags=re.I)
    text = re.sub(r"<br\s*/?>", "\n", text, flags=re.I)
    text = re.sub(r"<[^>]+>", "", text)
    return html.unescape(text).strip()


def _tex_escape(value: str) -> str:
    protected: List[str] = []

    def protect(match: re.Match[str]) -> str:
        protected.append(match.group(0))
        return f"@@MATH{len(protected) - 1}@@"

    text = re.sub(r"\$\$[\s\S]*?\$\$|\$[^$\n]+\$", protect, value or "")
    replacements = {
        "\\": r"\textbackslash{}",
        "&": r"\&",
        "%": r"\%",
        "$": r"\$",
        "#": r"\#",
        "_": r"\_",
        "{": r"\{",
        "}": r"\}",
    }
    escaped = "".join(replacements.get(char, char) for char in text)
    for index, math in enumerate(protected):
        escaped = escaped.replace(f"@@MATH{index}@@", math)
    return escaped


def _add_object_to_pptx_slide(ppt_slide: Any, obj: Dict[str, Any], assets: Dict[str, Any]) -> None:
    bbox = obj.get("bbox") if isinstance(obj.get("bbox"), dict) else {}
    left, top, width, height = _pptx_box(bbox)
    obj_type = str(obj.get("type") or "")
    if obj_type in {"title", "richText", "textbox", "callout", "equation"}:
        text = str(obj.get("text") or obj.get("latex") or _html_to_text(str(obj.get("rich_html") or "")) or obj.get("label") or "")
        if obj_type == "equation" and text:
            text = f"${text}$"
        shape = ppt_slide.shapes.add_textbox(left, top, width, height)
        frame = shape.text_frame
        frame.clear()
        paragraph = frame.paragraphs[0]
        paragraph.text = text
        paragraph.font.size = Pt(_object_font_size(obj, 28 if obj_type == "title" else 18))
        paragraph.font.bold = obj_type == "title"
        return
    if obj_type == "table":
        rows = obj.get("rows") if isinstance(obj.get("rows"), list) else []
        row_count = len(rows)
        col_count = max((len(row) for row in rows if isinstance(row, list)), default=0)
        if row_count and col_count:
            table_shape = ppt_slide.shapes.add_table(row_count, col_count, left, top, width, height)
            table = table_shape.table
            for r_index, row in enumerate(rows):
                for c_index, cell in enumerate((row if isinstance(row, list) else [row])[:col_count]):
                    table.cell(r_index, c_index).text = str(cell)
        return
    if obj_type in {"image", "placeholder"}:
        asset = assets.get(obj.get("asset_id")) if isinstance(assets, dict) else None
        data_uri = asset.get("data_uri") if isinstance(asset, dict) else None
        mime_type = str((asset or {}).get("mime_type") or "")
        if isinstance(data_uri, str) and data_uri.startswith("data:") and ";base64," in data_uri and mime_type != "image/svg+xml":
            try:
                image_bytes = base64.b64decode(data_uri.split(";base64,", 1)[1])
                ppt_slide.shapes.add_picture(io.BytesIO(image_bytes), left, top, width=width, height=height)
                return
            except Exception:
                pass
        shape = ppt_slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
        shape.text = str(obj.get("label") or "图片占位符")


def _pptx_box(bbox: Dict[str, Any]) -> tuple[Any, Any, Any, Any]:
    left = Inches((_finite_float(bbox.get("x"), 0) / CANVAS_WIDTH) * 13.333)
    top = Inches((_finite_float(bbox.get("y"), 0) / CANVAS_HEIGHT) * 7.5)
    width = Inches((_finite_float(bbox.get("width"), 120) / CANVAS_WIDTH) * 13.333)
    height = Inches((_finite_float(bbox.get("height"), 80) / CANVAS_HEIGHT) * 7.5)
    return left, top, width, height


def _object_font_size(obj: Dict[str, Any], fallback: int) -> int:
    style = obj.get("style") if isinstance(obj.get("style"), dict) else {}
    try:
        value = int(float(style.get("fontSize", fallback)))
    except (TypeError, ValueError):
        return fallback
    return min(max(value, 8), 64)


def _stable_id(prefix: str, value: str) -> str:
    return f"{prefix}-{hashlib.md5(value.encode('utf-8')).hexdigest()[:12]}"


def _safe_project_id(value: str) -> str:
    return re.sub(r"[^A-Za-z0-9_-]+", "", value or "")[:80]

"""Courseware parsing helpers for extracting slide content used by lecture generation."""
from __future__ import annotations

import base64
import csv
import io
import json
import mimetypes
import posixpath
import re
import zipfile
from html.parser import HTMLParser
from pathlib import Path
from typing import Any, Dict, List
from xml.etree import ElementTree

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Pt


MAX_INLINE_IMAGE_BYTES = 800 * 1024
IMAGE_EXTENSIONS = {".png", ".jpg", ".jpeg", ".gif", ".webp", ".bmp", ".tif", ".tiff", ".svg"}
LATEX_IMAGE_COMMAND_DEFAULT_OPTIONS = {
    "safecontentimage": r"width=0.7\textwidth",
}
TEXT_COURSEWARE_EXTENSIONS = {".tex", ".md", ".markdown", ".txt", ".rst", ".csv", ".json", ".html", ".htm", ".rtf"}
PPT_COURSEWARE_EXTENSIONS = {".pptx", ".ppt"}
DOCX_COURSEWARE_EXTENSIONS = {".docx"}
PDF_COURSEWARE_EXTENSIONS = {".pdf"}
ZIP_COURSEWARE_EXTENSIONS = {".zip"}
SUPPORTED_COURSEWARE_EXTENSIONS = (
    PPT_COURSEWARE_EXTENSIONS
    .union(TEXT_COURSEWARE_EXTENSIONS)
    .union(DOCX_COURSEWARE_EXTENSIONS)
    .union(PDF_COURSEWARE_EXTENSIONS)
    .union(ZIP_COURSEWARE_EXTENSIONS)
)
SUPPORTED_COURSEWARE_FORMATS_TEXT = ", ".join(sorted(SUPPORTED_COURSEWARE_EXTENSIONS))
_STRUCTURED_FIGURE_INDEX: Dict[str, Dict[str, Any]] | None = None
_STRUCTURED_FIGURE_ASSET_CACHE: Dict[str, Dict[str, Any]] = {}
_STRUCTURED_IMAGE_ASSET_CACHE: Dict[str, Dict[str, Any]] = {}


def parse_ppt(file_bytes: bytes) -> Dict[str, Any]:
    """Parse PPT bytes into per-slide structured content."""
    try:
        prs = Presentation(io.BytesIO(file_bytes))
    except Exception as exc:
        return {"success": False, "error": f"PPT 文件解析失败: {exc}"}

    slides: List[Dict[str, Any]] = []
    all_texts: List[str] = []

    for idx, slide in enumerate(prs.slides, 1):
        slide_data = _parse_slide(slide, idx)
        slides.append(slide_data)
        raw_text = str(slide_data.get("raw_text") or "").strip()
        if raw_text:
            all_texts.append(raw_text)

    return {
        "success": True,
        "slide_count": len(slides),
        "slides": slides,
        "full_text": "\n\n---\n\n".join(all_texts),
    }


def parse_courseware(file_bytes: bytes, filename: str = "") -> Dict[str, Any]:
    lower_name = (filename or "").lower()
    if lower_name.endswith(tuple(ZIP_COURSEWARE_EXTENSIONS)):
        return parse_zip_courseware(file_bytes, filename=filename)
    if lower_name.endswith(tuple(PPT_COURSEWARE_EXTENSIONS)):
        return parse_ppt(file_bytes)
    if lower_name.endswith(tuple(DOCX_COURSEWARE_EXTENSIONS)):
        return parse_docx_courseware(file_bytes, filename=filename)
    if lower_name.endswith(tuple(PDF_COURSEWARE_EXTENSIONS)):
        return parse_pdf_courseware(file_bytes, filename=filename)
    if lower_name.endswith(tuple(TEXT_COURSEWARE_EXTENSIONS)):
        return parse_text_courseware(file_bytes, filename=filename)
    return {"success": False, "error": f"仅支持 {SUPPORTED_COURSEWARE_FORMATS_TEXT} 格式"}


def parse_text_courseware(file_bytes: bytes, filename: str = "") -> Dict[str, Any]:
    decode_result = _decode_text_bytes(file_bytes)
    if not decode_result["success"]:
        return {"success": False, "error": f"文本课件解析失败: {decode_result['error']}"}
    text = decode_result["text"]
    lower_name = (filename or "").lower()
    if lower_name.endswith((".html", ".htm")):
        text = _html_to_text(text)
    elif lower_name.endswith(".rtf"):
        text = _rtf_to_text(text)
    elif lower_name.endswith(".csv"):
        text = _csv_to_text(text)
    elif lower_name.endswith(".json"):
        text = _json_to_text(text)
    if lower_name.endswith(".tex"):
        text = _normalize_text_newlines(text)
    slides = _slides_from_text(text, filename)
    missing_image_refs = _collect_missing_image_refs(slides)
    return {
        "success": True,
        "slide_count": len(slides),
        "slides": slides,
        "full_text": "\n\n---\n\n".join(str(slide.get("raw_text") or "") for slide in slides),
        "tex_content": text if lower_name.endswith(".tex") else "",
        "missing_image_refs": missing_image_refs,
    }


def parse_zip_courseware(file_bytes: bytes, filename: str = "") -> Dict[str, Any]:
    try:
        archive = zipfile.ZipFile(io.BytesIO(file_bytes))
    except zipfile.BadZipFile as exc:
        return {"success": False, "error": f"ZIP 课件解析失败: {exc}"}

    try:
        names = [name for name in archive.namelist() if not name.endswith("/") and not name.startswith("__MACOSX/")]
        tex_name = _select_tex_entry(names)
        if not tex_name:
            return {"success": False, "error": "ZIP 中未找到 .tex 主文件"}

        decode_result = _decode_text_bytes(archive.read(tex_name))
        if not decode_result["success"]:
            return {"success": False, "error": f"ZIP 中的 TeX 文件解析失败: {decode_result['error']}"}

        tex_text = _normalize_text_newlines(decode_result["text"])
        asset_map = _build_zip_image_asset_map(archive, names)
        slides = _slides_from_text(
            tex_text,
            filename=f"{filename}:{tex_name}" if filename else tex_name,
            image_assets=asset_map,
            tex_base_dir=posixpath.dirname(tex_name),
        )
        missing_image_refs = _collect_missing_image_refs(slides)
        return {
            "success": True,
            "slide_count": len(slides),
            "slides": slides,
            "full_text": "\n\n---\n\n".join(str(slide.get("raw_text") or "") for slide in slides),
            "tex_content": tex_text,
            "tex_source_file": tex_name,
            "missing_image_refs": missing_image_refs,
        }
    finally:
        archive.close()


def parse_docx_courseware(file_bytes: bytes, filename: str = "") -> Dict[str, Any]:
    try:
        with zipfile.ZipFile(io.BytesIO(file_bytes)) as archive:
            document_xml = archive.read("word/document.xml")
    except Exception as exc:
        return {"success": False, "error": f"DOCX 文件解析失败: {exc}"}

    try:
        root = ElementTree.fromstring(document_xml)
    except ElementTree.ParseError as exc:
        return {"success": False, "error": f"DOCX XML 解析失败: {exc}"}

    namespace = {"w": "http://schemas.openxmlformats.org/wordprocessingml/2006/main"}
    paragraphs: List[str] = []
    for paragraph in root.findall(".//w:p", namespace):
        text_parts = [node.text or "" for node in paragraph.findall(".//w:t", namespace)]
        text = "".join(text_parts).strip()
        if text:
            paragraphs.append(text)
    slides = _slides_from_text("\n\n".join(paragraphs), filename)
    return {
        "success": True,
        "slide_count": len(slides),
        "slides": slides,
        "full_text": "\n\n---\n\n".join(str(slide.get("raw_text") or "") for slide in slides),
    }


def parse_pdf_courseware(file_bytes: bytes, filename: str = "") -> Dict[str, Any]:
    try:
        from pypdf import PdfReader  # type: ignore
    except Exception:
        return {"success": False, "error": "PDF 课件解析需要安装 pypdf 依赖"}

    try:
        reader = PdfReader(io.BytesIO(file_bytes))
        chunks = []
        for index, page in enumerate(reader.pages, 1):
            text = (page.extract_text() or "").strip()
            if text:
                chunks.append(f"# 第 {index} 页\n{text}")
    except Exception as exc:
        return {"success": False, "error": f"PDF 文件解析失败: {exc}"}

    slides = _slides_from_text("\n\n---\n\n".join(chunks), filename)
    return {
        "success": True,
        "slide_count": len(slides),
        "slides": slides,
        "full_text": "\n\n---\n\n".join(str(slide.get("raw_text") or "") for slide in slides),
    }


def _parse_slide(slide: Any, index: int) -> Dict[str, Any]:
    title = ""
    body_texts: List[str] = []
    tables: List[Dict[str, Any]] = []
    notes = ""
    image_count = 0
    images: List[Dict[str, Any]] = []

    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            image_count += 1
            image_data = _extract_image(shape)
            if image_data:
                images.append(image_data)
            continue

        if getattr(shape, "has_text_frame", False):
            text = str(shape.text or "").strip()
            if text:
                if _is_title_shape(shape) and not title:
                    title = text
                else:
                    body_texts.append(text)

        if getattr(shape, "has_table", False):
            tables.append(_parse_table(shape.table))

    if getattr(slide, "has_notes_slide", False):
        notes_frame = slide.notes_slide.notes_text_frame
        if notes_frame:
            notes = str(notes_frame.text or "").strip()

    raw_text = _build_raw_text(title, body_texts, tables, notes)
    return {
        "index": index,
        "title": title,
        "body_texts": body_texts,
        "tables": tables,
        "notes": notes,
        "image_count": image_count,
        "images": images,
        "raw_text": raw_text,
    }


def _extract_image(shape: Any) -> Dict[str, Any]:
    try:
        image = shape.image
        image_bytes = image.blob
        data_uri = None
        oversized = len(image_bytes) > MAX_INLINE_IMAGE_BYTES
        if not oversized:
            encoded = base64.b64encode(image_bytes).decode("utf-8")
            data_uri = f"data:{image.content_type};base64,{encoded}"
        result = {
            "data_uri": data_uri,
            "width_emu": int(shape.width),
            "height_emu": int(shape.height),
            "left_emu": int(shape.left),
            "top_emu": int(shape.top),
        }
        if oversized:
            result["oversized"] = True
        return result
    except Exception:
        return {}


def _is_title_shape(shape: Any) -> bool:
    if not getattr(shape, "has_text_frame", False):
        return False
    try:
        ph_type = shape.placeholder_format.type
        if ph_type in {1, 3, 4}:  # TITLE, CENTER_TITLE, SUBTITLE
            return True
    except Exception:
        pass

    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            if run.font.size and run.font.size >= Pt(24):
                return True
            if run.font.bold:
                return True
    return False


def _parse_table(table: Any) -> Dict[str, Any]:
    rows = []
    for row in table.rows:
        rows.append([str(cell.text or "").strip() for cell in row.cells])
    return {"rows": rows}


def _decode_text_bytes(file_bytes: bytes) -> Dict[str, Any]:
    try:
        return {"success": True, "text": file_bytes.decode("utf-8-sig")}
    except UnicodeDecodeError:
        try:
            return {"success": True, "text": file_bytes.decode("gb18030")}
        except UnicodeDecodeError as exc:
            return {"success": False, "error": str(exc)}


def _normalize_text_newlines(text: str) -> str:
    return text.replace("\r\n", "\n").replace("\r", "\n")


def _select_tex_entry(names: List[str]) -> str:
    tex_names = [name for name in names if name.lower().endswith(".tex")]
    if not tex_names:
        return ""
    preferred = {"main.tex", "slides.tex", "presentation.tex", "lecture.tex"}
    for name in tex_names:
        if posixpath.basename(name).lower() in preferred:
            return name
    return tex_names[0]


def _build_zip_image_asset_map(archive: zipfile.ZipFile, names: List[str]) -> Dict[str, Dict[str, Any]]:
    assets: Dict[str, Dict[str, Any]] = {}
    for name in names:
        extension = posixpath.splitext(name)[1].lower()
        if extension not in IMAGE_EXTENSIONS:
            continue
        try:
            data = archive.read(name)
        except Exception:
            continue
        mime_type = mimetypes.guess_type(name)[0] or "application/octet-stream"
        asset = {"name": name, "bytes": data, "mime_type": mime_type}
        normalized = _normalize_archive_path(name)
        without_extension = _strip_archive_extension(normalized)
        assets.setdefault(normalized, asset)
        assets.setdefault(without_extension, asset)
        assets.setdefault(posixpath.basename(normalized), asset)
        assets.setdefault(_strip_archive_extension(posixpath.basename(normalized)), asset)
    return assets


def _looks_like_latex(text: str) -> bool:
    return "\\begin{frame" in text or "\\documentclass" in text or "\\begin{document}" in text


def _slides_from_text(
    text: str,
    filename: str,
    image_assets: Dict[str, Dict[str, Any]] | None = None,
    tex_base_dir: str = "",
) -> List[Dict[str, Any]]:
    normalized = _normalize_text_newlines(text)
    if not normalized.strip():
        return []
    if _looks_like_latex(normalized):
        slides = _slides_from_latex(normalized, filename, image_assets or {}, tex_base_dir)
        if slides:
            return slides
    chunks = _split_text_slides(normalized.strip())
    slides = []
    for index, chunk in enumerate(chunks, 1):
        title, body = _title_and_body_from_chunk(chunk, index)
        body_texts = _body_texts_from_text(body)
        images = _images_from_text_chunk(chunk, image_assets or {}, tex_base_dir)
        missing_image_refs = _missing_image_refs_from_text_chunk(chunk, image_assets or {}, tex_base_dir)
        raw_text = _build_raw_text(title, body_texts or [body], [], "")
        slides.append(
            {
                "index": index,
                "title": title,
                "body_texts": body_texts or [body],
                "tables": [],
                "notes": "",
                "image_count": len(images),
                "images": images,
                "missing_image_refs": missing_image_refs,
                "raw_text": raw_text,
                "source_file": filename,
            }
        )
    return slides


def _slides_from_latex(
    text: str,
    filename: str,
    image_assets: Dict[str, Dict[str, Any]],
    tex_base_dir: str,
) -> List[Dict[str, Any]]:
    frames = _split_latex_frames(text)
    if not frames:
        return []

    metadata = _extract_latex_metadata(text)
    template_text = _extract_latex_template_text(text, {"title page", "footline"})
    slides: List[Dict[str, Any]] = []

    for index, frame in enumerate(frames, 1):
        raw_body = str(frame.get("body") or "")
        title = _clean_latex_text(str(frame.get("title") or "")).strip()
        is_titlepage = "\\titlepage" in raw_body
        source_tex = str(frame.get("source") or "").strip()

        if is_titlepage:
            title = title or metadata.get("title") or f"第 {index} 页"
            body = _latex_titlepage_body(metadata)
            image_chunk = "\n".join(part for part in [template_text, raw_body] if part)
        else:
            inline_title, raw_body = _extract_frametitle(raw_body)
            title = title or inline_title or f"第 {index} 页"
            body = _latex_body_to_markdown(raw_body)
            image_chunk = raw_body

        body = body.strip()
        body_texts = [body] if body else []
        images = _images_from_text_chunk(image_chunk, image_assets, tex_base_dir)
        missing_image_refs = _missing_image_refs_from_text_chunk(image_chunk, image_assets, tex_base_dir)
        layout = _infer_latex_frame_layout(
            raw_body,
            images,
            content_markdown=body,
            image_assets=image_assets,
            tex_base_dir=tex_base_dir,
            is_titlepage=is_titlepage,
        )
        canvas_layout = _extract_kgts_canvas_layout(source_tex)
        if canvas_layout:
            layout["canvas"] = canvas_layout
        raw_text = _build_raw_text(title, body_texts, [], "")
        slides.append(
            {
                "index": index,
                "title": title,
                "body_texts": body_texts,
                "tables": [],
                "notes": "",
                "image_count": len(images),
                "images": images,
                "missing_image_refs": missing_image_refs,
                "raw_text": raw_text,
                "source_file": filename,
                "source_tex": source_tex,
                "source_body_tex": raw_body.strip(),
                "source_start": frame.get("source_start"),
                "source_end": frame.get("source_end"),
                "layout": layout,
            }
        )
    return slides


def _collect_missing_image_refs(slides: List[Dict[str, Any]]) -> List[str]:
    missing: List[str] = []
    for slide in slides:
        for ref in slide.get("missing_image_refs") or []:
            value = str(ref or "").strip()
            if value and value not in missing:
                missing.append(value)
    return missing


def _extract_kgts_canvas_layout(source: str) -> Dict[str, Any]:
    match = re.search(r"^\s*%\s*KGTS_LAYOUT\s+({.*})\s*$", source or "", re.MULTILINE)
    if not match:
        return {}
    try:
        value = json.loads(match.group(1))
    except json.JSONDecodeError:
        return {}
    return value if isinstance(value, dict) else {}


def _images_from_text_chunk(chunk: str, image_assets: Dict[str, Dict[str, Any]], tex_base_dir: str) -> List[Dict[str, Any]]:
    images: List[Dict[str, Any]] = []
    seen: set[str] = set()
    refs = _latex_image_refs(chunk)
    for options, ref in refs:
        asset = _resolve_image_asset(image_assets, ref, tex_base_dir) or _resolve_structured_figure_asset(ref)
        if not asset:
            continue
        source_path = str(asset.get("name") or ref)
        if source_path in seen:
            continue
        seen.add(source_path)
        image_bytes = asset.get("bytes") or b""
        oversized = len(image_bytes) > MAX_INLINE_IMAGE_BYTES
        data_uri = None
        if image_bytes and not oversized:
            encoded = base64.b64encode(image_bytes).decode("utf-8")
            data_uri = f"data:{asset.get('mime_type') or 'application/octet-stream'};base64,{encoded}"
        image_info = {
            "data_uri": data_uri,
            "width_emu": 0,
            "height_emu": 0,
            "left_emu": 0,
            "top_emu": 0,
            "source_path": source_path,
            "tex_options": options,
            "tex_ref": asset.get("tex_ref") or ref,
        }
        width_match = re.search(r"width\s*=\s*([0-9.]+)\\textwidth", options or "")
        if width_match:
            image_info["width_ratio"] = float(width_match.group(1))
        height_match = re.search(r"height\s*=\s*([0-9.]+)\\textheight", options or "")
        if height_match:
            image_info["height_ratio"] = float(height_match.group(1))
        if oversized:
            image_info["oversized"] = True
        images.append(image_info)
    return images


def _missing_image_refs_from_text_chunk(chunk: str, image_assets: Dict[str, Dict[str, Any]], tex_base_dir: str) -> List[str]:
    missing: List[str] = []
    for _options, ref in _latex_image_refs(chunk):
        clean_ref = ref.strip().strip("\"'").replace("\\", "/")
        if not clean_ref:
            continue
        asset = _resolve_image_asset(image_assets, clean_ref, tex_base_dir) or _resolve_structured_figure_asset(clean_ref)
        if not asset and clean_ref not in missing:
            missing.append(clean_ref)
    return missing


def _latex_image_refs(chunk: str) -> List[tuple[str, str]]:
    refs = re.findall(r"\\includegraphics(?:\[([^\]]*)\])?\{([^}]+)\}", chunk or "")
    for command, default_options in LATEX_IMAGE_COMMAND_DEFAULT_OPTIONS.items():
        pattern = re.compile(rf"\\{re.escape(command)}(?![A-Za-z@])(?:\[[^\]]*\])?\{{([^{{}}]+)\}}")
        for match in pattern.finditer(chunk or ""):
            refs.append((default_options, match.group(1)))
    return refs


def _resolve_image_asset(image_assets: Dict[str, Dict[str, Any]], ref: str, tex_base_dir: str) -> Dict[str, Any] | None:
    clean_ref = ref.strip().strip("\"'").replace("\\", "/")
    if not clean_ref:
        return None
    candidates = [
        clean_ref,
        posixpath.join(tex_base_dir, clean_ref) if tex_base_dir else clean_ref,
        clean_ref.lstrip("./"),
        posixpath.join(tex_base_dir, clean_ref.lstrip("./")) if tex_base_dir else clean_ref.lstrip("./"),
        posixpath.basename(clean_ref),
    ]
    extension = posixpath.splitext(clean_ref)[1]
    if not extension:
        candidates.extend(f"{candidate}{image_extension}" for candidate in list(candidates) for image_extension in IMAGE_EXTENSIONS)
    for candidate in candidates:
        normalized = _normalize_archive_path(candidate)
        asset = image_assets.get(normalized) or image_assets.get(_strip_archive_extension(normalized))
        if asset:
            return asset
    return None


def _resolve_structured_figure_asset(ref: str) -> Dict[str, Any] | None:
    index = _load_structured_figure_index()
    if index:
        for figure_id in _structured_figure_id_candidates(ref):
            entry = index.get(figure_id)
            if not entry:
                continue
            if figure_id in _STRUCTURED_FIGURE_ASSET_CACHE:
                return _STRUCTURED_FIGURE_ASSET_CACHE[figure_id]
            asset_path = str(entry.get("asset_path") or "").replace("\\", "/").lstrip("/")
            asset = _read_structured_image_asset(asset_path, tex_ref=asset_path, cache_key=f"figure:{figure_id}")
            if asset:
                asset["figure_id"] = figure_id
                _STRUCTURED_FIGURE_ASSET_CACHE[figure_id] = asset
                return asset
    return _resolve_structured_image_file_asset(ref)


def _resolve_structured_image_file_asset(ref: str) -> Dict[str, Any] | None:
    clean_ref = str(ref or "").strip().strip("\"'").replace("\\", "/").lstrip("/")
    if not clean_ref:
        return None
    candidates = [clean_ref]
    basename = posixpath.basename(clean_ref)
    if basename and basename != clean_ref:
        candidates.append(basename)
    for prefix in ("figures", "textbook/figures"):
        if basename:
            candidates.append(f"{prefix}/{basename}")
    extension = posixpath.splitext(clean_ref)[1]
    if not extension:
        for value in list(candidates):
            candidates.extend(f"{value}{image_extension}" for image_extension in IMAGE_EXTENSIONS)
    for candidate in candidates:
        asset = _read_structured_image_asset(candidate, tex_ref=clean_ref, cache_key=f"file:{candidate}")
        if asset:
            return asset
    return None


def _read_structured_image_asset(asset_path: str, *, tex_ref: str, cache_key: str) -> Dict[str, Any] | None:
    clean_path = _normalize_archive_path(asset_path)
    if not clean_path or posixpath.splitext(clean_path)[1].lower() not in IMAGE_EXTENSIONS:
        return None
    if cache_key in _STRUCTURED_IMAGE_ASSET_CACHE:
        return _STRUCTURED_IMAGE_ASSET_CACHE[cache_key]
    root_dir = _structured_root_dir().resolve()
    absolute_path = (root_dir / clean_path).resolve()
    try:
        absolute_path.relative_to(root_dir)
    except ValueError:
        return None
    try:
        image_bytes = absolute_path.read_bytes()
    except OSError:
        return None
    asset = {
        "name": clean_path,
        "bytes": image_bytes,
        "mime_type": mimetypes.guess_type(clean_path)[0] or "application/octet-stream",
        "tex_ref": tex_ref,
    }
    _STRUCTURED_IMAGE_ASSET_CACHE[cache_key] = asset
    return asset


def _load_structured_figure_index() -> Dict[str, Dict[str, Any]]:
    global _STRUCTURED_FIGURE_INDEX
    if _STRUCTURED_FIGURE_INDEX is not None:
        return _STRUCTURED_FIGURE_INDEX
    library_path = _structured_root_dir() / "figure_library.json"
    try:
        payload = json.loads(library_path.read_text(encoding="utf-8"))
    except (OSError, json.JSONDecodeError):
        _STRUCTURED_FIGURE_INDEX = {}
        return _STRUCTURED_FIGURE_INDEX
    figures = payload.get("figures") if isinstance(payload, dict) else None
    _STRUCTURED_FIGURE_INDEX = figures if isinstance(figures, dict) else {}
    return _STRUCTURED_FIGURE_INDEX


def _structured_root_dir() -> Path:
    return Path(__file__).resolve().parents[1] / "structured"


def _structured_figure_id_candidates(ref: str) -> List[str]:
    clean = str(ref or "").strip().strip("\"'").replace("\\", "/")
    basename = posixpath.basename(clean)
    root, _extension = posixpath.splitext(basename)
    values = [clean, basename, root]
    candidates: List[str] = []

    def add(value: str) -> None:
        normalized = value.strip()
        if not normalized:
            return
        normalized = re.sub(r"^(?:figure|fig|图)\s*[-_ ]*", "", normalized, flags=re.I)
        normalized = re.sub(r"^figure\s*", "", normalized, flags=re.I)
        if normalized.startswith("a") and re.match(r"a\d+\.", normalized, flags=re.I):
            normalized = "A" + normalized[1:]
        if normalized not in candidates:
            candidates.append(normalized)
        if re.match(r"^(?:A\d+|\d+)\.\d+[a-z]$", normalized):
            parent = normalized[:-1]
            if parent not in candidates:
                candidates.append(parent)

    for value in values:
        add(value)
        for match in re.finditer(r"\b(A\d+|\d+)\.\d+[a-z]?\b", value, flags=re.I):
            add(match.group(0))
    return candidates


def _image_appears_before_text(body: str) -> bool:
    text = _strip_latex_comments(body)
    image_positions = [match.start() for match in re.finditer(r"\\includegraphics(?![A-Za-z@])", text)]
    for command in LATEX_IMAGE_COMMAND_DEFAULT_OPTIONS:
        image_positions.extend(match.start() for match in re.finditer(rf"\\{re.escape(command)}(?![A-Za-z@])", text))
    first_image = min(image_positions) if image_positions else -1
    if first_image < 0:
        return False
    first_item = text.find("\\item")
    first_formula = min([pos for pos in [text.find("\\["), text.find("\\begin{align")] if pos >= 0] or [-1])
    content_positions = [pos for pos in [first_item, first_formula] if pos >= 0]
    return not content_positions or first_image < min(content_positions)


def _latex_width_ratio(value: str) -> float | None:
    text = str(value or "").strip()
    match = re.search(r"([0-9]*\.?[0-9]+)\s*\\(?:textwidth|linewidth|paperwidth)", text)
    if not match:
        return None
    try:
        return float(match.group(1))
    except ValueError:
        return None


def _extract_latex_columns(
    body: str,
    image_assets: Dict[str, Dict[str, Any]],
    tex_base_dir: str,
) -> List[Dict[str, Any]]:
    columns: List[Dict[str, Any]] = []
    pattern = re.compile(r"\\begin\{column\}")
    position = 0

    while True:
        match = pattern.search(body, position)
        if not match:
            break
        cursor = _skip_latex_modifiers(body, match.end())
        width = ""
        if cursor < len(body) and body[cursor] == "{":
            width, cursor = _read_latex_braced_group(body, cursor)
        end_match = re.search(r"\\end\{column\}", body[cursor:])
        if not end_match:
            break
        content_end = cursor + end_match.start()
        source_end = cursor + end_match.end()
        column_body = body[cursor:content_end]
        column_images = _images_from_text_chunk(column_body, image_assets, tex_base_dir)
        image_widths = [
            image.get("width_ratio")
            for image in column_images
            if isinstance(image.get("width_ratio"), (int, float))
        ]
        columns.append(
            {
                "width_ratio": _latex_width_ratio(width),
                "content": _latex_body_to_markdown(column_body).strip(),
                "images": column_images,
                "image_count": len(column_images),
                "image_first": _image_appears_before_text(column_body),
                "align": "center" if "\\centering" in column_body or "\\begin{center}" in column_body else "left",
                "max_image_width": max(image_widths) if image_widths else None,
                "source_tex": body[match.start():source_end].strip(),
            }
        )
        position = source_end

    return columns


def _remove_latex_columns_blocks(body: str) -> str:
    result: List[str] = []
    pattern = re.compile(r"\\begin\{columns\}")
    position = 0
    while True:
        match = pattern.search(body, position)
        if not match:
            result.append(body[position:])
            break
        result.append(body[position:match.start()])
        end_match = re.search(r"\\end\{columns\}", body[match.end():])
        if not end_match:
            position = match.end()
            continue
        position = match.end() + end_match.end()
    return "".join(result)


def _infer_latex_frame_layout(
    body: str,
    images: List[Dict[str, Any]],
    *,
    content_markdown: str = "",
    image_assets: Dict[str, Dict[str, Any]] | None = None,
    tex_base_dir: str = "",
    is_titlepage: bool = False,
) -> Dict[str, Any]:
    has_columns = "\\begin{columns" in body or "\\begin{column" in body
    has_center = "\\centering" in body or "\\begin{center}" in body
    image_first = _image_appears_before_text(body) if images else False
    columns = _extract_latex_columns(body, image_assets or {}, tex_base_dir) if has_columns else []
    column_count = len(columns)
    outside_content = _latex_body_to_markdown(_remove_latex_columns_blocks(body)).strip() if columns else ""

    image_widths = [
        image.get("width_ratio")
        for image in images
        if isinstance(image.get("width_ratio"), (int, float))
    ]
    max_image_width = max(image_widths) if image_widths else None

    if is_titlepage:
        mode = "title"
    elif column_count > 1:
        mode = "columns"
    elif images and not content_markdown.strip():
        mode = "image_only"
    elif images:
        mode = "image_text" if image_first else "text_image"
    else:
        mode = "text"

    return {
        "mode": mode,
        "has_columns": has_columns,
        "column_count": column_count,
        "columns": columns,
        "outside_content": outside_content,
        "align": "center" if has_center else "left",
        "image_first": image_first,
        "image_count": len(images),
        "max_image_width": max_image_width,
    }


def _normalize_archive_path(path: str) -> str:
    normalized = posixpath.normpath(path.replace("\\", "/")).lstrip("/")
    if normalized == ".":
        return ""
    return normalized.lower()


def _strip_archive_extension(path: str) -> str:
    root, extension = posixpath.splitext(path)
    return root if extension.lower() in IMAGE_EXTENSIONS else path


class _TextOnlyHtmlParser(HTMLParser):
    def __init__(self) -> None:
        super().__init__()
        self.parts: List[str] = []

    def handle_starttag(self, tag: str, attrs: List[tuple[str, str | None]]) -> None:
        if tag.lower() in {"br", "p", "div", "li", "section", "article", "h1", "h2", "h3", "tr"}:
            self.parts.append("\n")

    def handle_data(self, data: str) -> None:
        text = data.strip()
        if text:
            self.parts.append(text)

    def text(self) -> str:
        return re.sub(r"\n{3,}", "\n\n", " ".join(self.parts)).strip()


def _html_to_text(text: str) -> str:
    parser = _TextOnlyHtmlParser()
    parser.feed(text)
    return parser.text() or text


def _rtf_to_text(text: str) -> str:
    stripped = re.sub(r"\\'[0-9a-fA-F]{2}", " ", text)
    stripped = re.sub(r"\\[a-zA-Z]+\d* ?", " ", stripped)
    stripped = stripped.replace("{", " ").replace("}", " ")
    return re.sub(r"\s+", " ", stripped).strip()


def _csv_to_text(text: str) -> str:
    rows = list(csv.reader(io.StringIO(text)))
    lines = [" | ".join(cell.strip() for cell in row if cell.strip()) for row in rows]
    return "\n".join(line for line in lines if line)


def _json_to_text(text: str) -> str:
    try:
        payload = json.loads(text)
    except json.JSONDecodeError:
        return text
    if isinstance(payload, list):
        chunks = [_json_item_to_text(item, index + 1) for index, item in enumerate(payload)]
        return "\n\n---\n\n".join(chunk for chunk in chunks if chunk.strip())
    if isinstance(payload, dict):
        slides = payload.get("slides") or payload.get("pages")
        if isinstance(slides, list):
            chunks = [_json_item_to_text(item, index + 1) for index, item in enumerate(slides)]
            return "\n\n---\n\n".join(chunk for chunk in chunks if chunk.strip())
        return _json_item_to_text(payload, 1)
    return str(payload)


def _json_item_to_text(item: Any, index: int) -> str:
    if isinstance(item, dict):
        title = str(item.get("title") or item.get("name") or f"第 {index} 页").strip()
        content = item.get("content") or item.get("text") or item.get("body") or item.get("bullets") or item
        if isinstance(content, list):
            body = "\n".join(str(value) for value in content)
        elif isinstance(content, dict):
            body = "\n".join(f"{key}: {value}" for key, value in content.items() if key not in {"title", "name"})
        else:
            body = str(content or "")
        return f"# {title}\n{body}".strip()
    return f"# 第 {index} 页\n{item}"


def _skip_latex_whitespace(text: str, cursor: int) -> int:
    while cursor < len(text) and text[cursor].isspace():
        cursor += 1
    return cursor


def _skip_latex_delimited(text: str, cursor: int, open_char: str, close_char: str) -> int:
    if cursor >= len(text) or text[cursor] != open_char:
        return cursor
    depth = 1
    cursor += 1
    while cursor < len(text) and depth > 0:
        if text[cursor] == "\\":
            cursor += 2
            continue
        if text[cursor] == open_char:
            depth += 1
        elif text[cursor] == close_char:
            depth -= 1
        cursor += 1
    return cursor


def _read_latex_braced_group(text: str, cursor: int) -> tuple[str, int]:
    if cursor >= len(text) or text[cursor] != "{":
        return "", cursor
    end = _skip_latex_delimited(text, cursor, "{", "}")
    if end <= cursor:
        return "", cursor
    return text[cursor + 1:end - 1], end


def _skip_latex_modifiers(text: str, cursor: int) -> int:
    while True:
        cursor = _skip_latex_whitespace(text, cursor)
        if cursor < len(text) and text[cursor] == "<":
            cursor = _skip_latex_delimited(text, cursor, "<", ">")
            continue
        if cursor < len(text) and text[cursor] == "[":
            cursor = _skip_latex_delimited(text, cursor, "[", "]")
            continue
        return cursor


def _split_latex_frames(text: str) -> List[Dict[str, str]]:
    frames: List[Dict[str, str]] = []
    begin_pattern = re.compile(r"\\begin\{frame\}")
    position = 0

    while True:
        begin_match = begin_pattern.search(text, position)
        if not begin_match:
            break

        cursor = _skip_latex_modifiers(text, begin_match.end())
        title = ""
        if cursor < len(text) and text[cursor] == "{":
            title, cursor = _read_latex_braced_group(text, cursor)

        end_match = re.search(r"\\end\{frame\}", text[cursor:])
        if not end_match:
            break
        body_end = cursor + end_match.start()
        source_end = cursor + end_match.end()
        frames.append({
            "title": title,
            "body": text[cursor:body_end],
            "source": text[begin_match.start():source_end],
            "source_start": begin_match.start(),
            "source_end": source_end,
        })
        position = cursor + end_match.end()

    return frames


def _extract_latex_command_arg(text: str, command: str) -> str:
    pattern = re.compile(rf"\\{re.escape(command)}(?![A-Za-z@])")
    for match in pattern.finditer(text):
        cursor = match.end()
        while True:
            cursor = _skip_latex_whitespace(text, cursor)
            if cursor < len(text) and text[cursor] == "[":
                cursor = _skip_latex_delimited(text, cursor, "[", "]")
                continue
            break
        if cursor < len(text) and text[cursor] == "{":
            value, _end = _read_latex_braced_group(text, cursor)
            return value
    return ""


def _extract_latex_metadata(text: str) -> Dict[str, str]:
    metadata: Dict[str, str] = {}
    for command in ("title", "subtitle", "author", "date"):
        value = _extract_latex_command_arg(text, command)
        if value:
            metadata[command] = _clean_latex_text(value)
    return metadata


def _extract_latex_template_text(text: str, template_names: set[str]) -> str:
    parts: List[str] = []
    pattern = re.compile(r"\\setbeamertemplate(?![A-Za-z@])")
    position = 0
    while True:
        match = pattern.search(text, position)
        if not match:
            break
        cursor = _skip_latex_whitespace(text, match.end())
        if cursor >= len(text) or text[cursor] != "{":
            position = match.end()
            continue
        name, cursor = _read_latex_braced_group(text, cursor)
        cursor = _skip_latex_whitespace(text, cursor)
        if cursor < len(text) and text[cursor] == "[":
            cursor = _skip_latex_delimited(text, cursor, "[", "]")
            cursor = _skip_latex_whitespace(text, cursor)
        if cursor < len(text) and text[cursor] == "{":
            content, cursor = _read_latex_braced_group(text, cursor)
            if name.strip().lower() in template_names:
                parts.append(content)
        position = max(cursor, match.end() + 1)
    return "\n".join(parts)


def _extract_frametitle(body: str) -> tuple[str, str]:
    pattern = re.compile(r"\\frametitle(?![A-Za-z@])")
    match = pattern.search(body)
    if not match:
        return "", body
    cursor = _skip_latex_modifiers(body, match.end())
    if cursor >= len(body) or body[cursor] != "{":
        return "", body
    title, end = _read_latex_braced_group(body, cursor)
    cleaned = _clean_latex_text(title)
    return cleaned, f"{body[:match.start()]}\n{body[end:]}"


def _latex_titlepage_body(metadata: Dict[str, str]) -> str:
    parts = [
        metadata.get("subtitle", ""),
        metadata.get("author", ""),
        metadata.get("date", ""),
    ]
    return "\n".join(part for part in parts if part).strip()


def _strip_latex_comments(text: str) -> str:
    return re.sub(r"(?<!\\)%.*", "", text)


def _replace_latex_commands_by_arg(
    text: str,
    commands: set[str],
    arg_count: int,
    keep_arg_index: int,
) -> str:
    if not text or not commands:
        return text
    result: List[str] = []
    cursor = 0
    command_names = sorted(commands, key=len, reverse=True)

    while cursor < len(text):
        replacement_done = False
        if text[cursor] == "\\":
            for command in command_names:
                prefix = f"\\{command}"
                if not text.startswith(prefix, cursor):
                    continue
                after_name = cursor + len(prefix)
                if after_name < len(text) and (text[after_name].isalpha() or text[after_name] == "@"):
                    continue
                if after_name < len(text) and text[after_name] == "*":
                    after_name += 1
                scan = _skip_latex_modifiers(text, after_name)
                args: List[str] = []
                ok = True
                for _index in range(arg_count):
                    scan = _skip_latex_modifiers(text, scan)
                    if scan >= len(text) or text[scan] != "{":
                        ok = False
                        break
                    value, scan = _read_latex_braced_group(text, scan)
                    args.append(value)
                if ok:
                    if 0 <= keep_arg_index < len(args):
                        result.append(args[keep_arg_index])
                    cursor = scan
                    replacement_done = True
                    break
        if replacement_done:
            continue
        result.append(text[cursor])
        cursor += 1

    return "".join(result)


def _replace_tikz_nodes_with_content(text: str) -> str:
    result: List[str] = []
    position = 0
    pattern = re.compile(r"\\node(?![A-Za-z@])")

    while True:
        match = pattern.search(text, position)
        if not match:
            result.append(text[position:])
            break
        result.append(text[position:match.start()])
        semicolon = text.find(";", match.end())
        if semicolon == -1:
            position = match.end()
            continue

        best: tuple[int, int, str] | None = None
        cursor = match.end()
        while cursor < semicolon:
            brace = text.find("{", cursor, semicolon + 1)
            if brace == -1:
                break
            value, end = _read_latex_braced_group(text, brace)
            if end <= brace:
                cursor = brace + 1
                continue
            if end <= semicolon + 1:
                best = (brace, end, value)
            cursor = end

        if best is not None:
            result.append(f"\n{best[2]}\n")
        position = semicolon + 1

    return "".join(result)


def _remove_custom_latex_image_commands(text: str) -> str:
    for command in LATEX_IMAGE_COMMAND_DEFAULT_OPTIONS:
        text = re.sub(rf"\\{re.escape(command)}(?![A-Za-z@])(?:\[[^\]]*\])?\{{[^{{}}]+\}}", "", text)
    return text


def _clean_latex_math(formula: str) -> str:
    cleaned = _strip_latex_comments(formula)
    cleaned = _replace_latex_commands_by_arg(cleaned, {"textcolor", "color"}, 2, 1)
    cleaned = _replace_latex_commands_by_arg(cleaned, {"tikzmark"}, 1, -1)
    cleaned = re.sub(r"\\(?:onslide|only|uncover|visible)<[^>]*>", "", cleaned)
    cleaned = re.sub(r"[_^]\{\s*\}", "", cleaned)
    cleaned = re.sub(r"[ \t]+\n", "\n", cleaned)
    return cleaned.strip()


def _extract_latex_math_placeholders(text: str) -> tuple[str, Dict[str, str]]:
    placeholders: Dict[str, str] = {}

    def store(markdown: str) -> str:
        token = f"KGTSMATH{len(placeholders)}TOKEN"
        placeholders[token] = markdown
        return token

    def display_replacement(match: re.Match[str]) -> str:
        formula = _clean_latex_math(match.group(1))
        return store(f"\n\n$$\n{formula}\n$$\n\n") if formula else "\n"

    def environment_replacement(match: re.Match[str]) -> str:
        environment = match.group(1).replace("*", "")
        formula = _clean_latex_math(match.group(2))
        if not formula:
            return "\n"
        if environment == "align":
            formula = f"\\begin{{aligned}}\n{formula}\n\\end{{aligned}}"
        elif environment == "gather":
            formula = f"\\begin{{gathered}}\n{formula}\n\\end{{gathered}}"
        return store(f"\n\n$$\n{formula}\n$$\n\n")

    def inline_parentheses_replacement(match: re.Match[str]) -> str:
        formula = _clean_latex_math(match.group(1))
        return store(f"${formula}$") if formula else ""

    def inline_dollar_replacement(match: re.Match[str]) -> str:
        formula = _clean_latex_math(match.group(1))
        return store(f"${formula}$") if formula else ""

    text = re.sub(r"\\\[\s*([\s\S]*?)\s*\\\]", display_replacement, text)
    text = re.sub(
        r"\\begin\{(align\*?|equation\*?|gather\*?)\}([\s\S]*?)\\end\{\1\}",
        environment_replacement,
        text,
    )
    text = re.sub(r"\\\(([\s\S]*?)\\\)", inline_parentheses_replacement, text)
    text = re.sub(r"(?<!\$)\$(?!\$)([\s\S]*?)(?<!\$)\$(?!\$)", inline_dollar_replacement, text)
    return text, placeholders


def _restore_latex_math_placeholders(text: str, placeholders: Dict[str, str]) -> str:
    for token, markdown in placeholders.items():
        text = text.replace(token, markdown)
    return text


def _remove_latex_layout_commands(text: str) -> str:
    text = _replace_latex_commands_by_arg(text, {"setlength"}, 2, -1)
    text = _replace_latex_commands_by_arg(
        text,
        {"vspace", "hspace", "vskip", "hskip", "tikzmark", "label", "usebeamertemplate"},
        1,
        -1,
    )
    text = re.sub(r"\\includegraphics(?:\[[^\]]*\])?\{[^}]+\}", "", text)
    text = _remove_custom_latex_image_commands(text)
    text = re.sub(r"\\(?:draw|path|coordinate|filldraw|fill)\b[\s\S]*?;", "\n", text)
    text = re.sub(
        r"\\begin\{(?:itemize|enumerate|center|columns|column|minipage|tikzpicture|block|alertblock)\}"
        r"(?:\[[^\]]*\])?(?:\{[^{}]*\})?",
        "\n",
        text,
    )
    text = re.sub(r"\\end\{(?:itemize|enumerate|center|columns|column|minipage|tikzpicture|block|alertblock)\}", "\n", text)
    text = re.sub(
        r"\\(?:centering|raggedright|raggedleft|vfill|hfill|small|scriptsize|footnotesize|tiny|normalsize|large|Large|LARGE|huge|Huge|bfseries|itshape|leavevmode|par)\b",
        " ",
        text,
    )
    return text


def _replace_latex_escaped_chars(text: str) -> str:
    replacements = {
        r"\%": "%",
        r"\&": "&",
        r"\_": "_",
        r"\#": "#",
        r"\$": "$",
        r"\{": "{",
        r"\}": "}",
        r"\textasciitilde{}": "~",
    }
    for source, target in replacements.items():
        text = text.replace(source, target)
    return text.replace("~", " ")


def _is_noise_latex_line(line: str) -> bool:
    if not line:
        return True
    stripped = line.strip()
    if stripped in {"{", "}", ";", "};", "[", "]"}:
        return True
    noise_markers = (
        "current page",
        "remember picture",
        "anchor=",
        "inner sep",
        "rounded corners",
        "text width=",
        "callout absolute pointer",
    )
    return any(marker in stripped for marker in noise_markers)


def _normalize_latex_markdown(text: str) -> str:
    lines: List[str] = []
    in_math = False
    previous = ""
    for raw_line in text.splitlines():
        line = re.sub(r"[ \t]+", " ", raw_line).strip()
        if line == "$$":
            lines.append(line)
            in_math = not in_math
            previous = line
            continue
        if in_math:
            if line:
                lines.append(line)
            continue
        line = line.strip("{}; ")
        line = re.sub(r"\s+([,.;:!?，。；：！？])", r"\1", line)
        if _is_noise_latex_line(line):
            continue
        if line.startswith("-"):
            line = "- " + line.lstrip("- ").strip()
            if line == "-":
                continue
        if line == previous:
            continue
        lines.append(line)
        previous = line

    normalized = "\n".join(lines)
    normalized = re.sub(r"\n{3,}", "\n\n", normalized)
    normalized = re.sub(r"\n(?=- )", "\n", normalized)
    return normalized.strip()


def _latex_body_to_markdown(body: str) -> str:
    text = _strip_latex_comments(body)
    text = _replace_latex_commands_by_arg(text, {"onslide", "only", "uncover", "visible"}, 1, 0)
    text = _replace_tikz_nodes_with_content(text)
    text, math_placeholders = _extract_latex_math_placeholders(text)
    text = re.sub(r"\\item(?:<[^>]*>)?(?:\[[^\]]*\])?\s*", "\n- ", text)
    text = _remove_latex_layout_commands(text)

    for _ in range(4):
        previous = text
        text = _replace_latex_commands_by_arg(text, {"textcolor", "href", "parbox", "colorbox"}, 2, 1)
        text = _replace_latex_commands_by_arg(
            text,
            {"textbf", "textit", "emph", "alert", "underline", "structure", "textrm", "mbox", "caption"},
            1,
            0,
        )
        if text == previous:
            break

    text = text.replace("\\\\", "\n")
    text = _replace_latex_escaped_chars(text)
    text = re.sub(r"\\[A-Za-z@]+(?:\*|<[^>]*>)?(?:\[[^\]]*\])?", " ", text)
    text = _restore_latex_math_placeholders(text, math_placeholders)
    return _normalize_latex_markdown(text)


def _clean_latex_text(text: str) -> str:
    cleaned = _latex_body_to_markdown(text)
    cleaned = cleaned.replace("$$", " ").replace("$", " ")
    cleaned = re.sub(r"^-\s*", "", cleaned)
    cleaned = re.sub(r"\s+", " ", cleaned)
    return cleaned.strip()


def _split_text_slides(text: str) -> List[str]:
    if "\\begin{frame" in text:
        frames = [
            (frame.get("title", ""), frame.get("body", ""))
            for frame in _split_latex_frames(text)
        ]
        chunks = []
        for title, body in frames:
            title_line = f"# {title.strip()}\n" if title.strip() else ""
            chunks.append((title_line + body).strip())
        if chunks:
            return chunks
    if "\n---\n" in text:
        chunks = [chunk.strip() for chunk in re.split(r"\n\s*---\s*\n", text) if chunk.strip()]
        if chunks:
            return chunks
    heading_chunks = re.split(r"(?m)^(?=#{1,3}\s+)", text)
    chunks = [chunk.strip() for chunk in heading_chunks if chunk.strip()]
    if len(chunks) > 1:
        return chunks
    paragraphs = [item.strip() for item in re.split(r"\n\s*\n", text) if item.strip()]
    if len(paragraphs) <= 1:
        return [text]
    grouped: List[str] = []
    current: List[str] = []
    current_len = 0
    for paragraph in paragraphs:
        current.append(paragraph)
        current_len += len(paragraph)
        if current_len >= 900:
            grouped.append("\n\n".join(current))
            current = []
            current_len = 0
    if current:
        grouped.append("\n\n".join(current))
    return grouped[:40]


def _title_and_body_from_chunk(chunk: str, index: int) -> tuple[str, str]:
    lines = [line.strip() for line in chunk.splitlines() if line.strip()]
    if not lines:
        return f"第 {index} 页", ""
    first = lines[0]
    heading = re.match(r"^#{1,6}\s+(.+)$", first)
    tex_title = re.match(r"^\\frametitle\{(.+)\}$", first)
    if heading or tex_title:
        title = (heading or tex_title).group(1).strip()
        return title, "\n".join(lines[1:]).strip()
    if len(first) <= 80 and len(lines) > 1:
        return first, "\n".join(lines[1:]).strip()
    return f"第 {index} 页", "\n".join(lines).strip()


def _body_texts_from_text(body: str) -> List[str]:
    lines = []
    for line in body.splitlines():
        clean = re.sub(r"^\s*[-*•]\s*", "", line).strip()
        clean = re.sub(r"^\\item\s+", "", clean).strip()
        clean = re.sub(r"\\(?:begin|end)\{[^}]+\}", "", clean).strip()
        if clean:
            lines.append(clean)
    if lines:
        return lines[:12]
    compact = re.sub(r"\s+", " ", body).strip()
    return [compact] if compact else []


def _build_raw_text(title: str, body_texts: List[str], tables: List[Dict[str, Any]], notes: str) -> str:
    parts: List[str] = []
    if title:
        parts.append(f"## {title}")
    parts.extend(text for text in body_texts if text)
    for table_data in tables:
        table_text = _format_table_text(table_data)
        if table_text:
            parts.append(table_text)
    if notes:
        parts.append(f"[备注] {notes}")
    return "\n".join(parts)


def _format_table_text(table_data: Dict[str, Any]) -> str:
    rows = table_data.get("rows", [])
    if not rows:
        return ""
    lines = []
    for index, row in enumerate(rows):
        line = " | ".join(str(cell) for cell in row)
        lines.append(line)
        if index == 0:
            lines.append("-" * max(len(line), 10))
    return "\n".join(lines)


def build_ppt_lecture_prompt_data(parse_result: Dict[str, Any]) -> Dict[str, Any]:
    """Convert parser output into the payload used by lecture generation."""
    if not parse_result.get("success"):
        return {
            "chapter_title": "未命名PPT",
            "chapter_content": "",
            "slide_details": [],
            "total_slides": 0,
        }

    slides = parse_result.get("slides", [])
    first_title = next((slide.get("title") for slide in slides if slide.get("title")), "")

    slide_details = []
    for slide in slides:
        slide_details.append(
            {
                "index": slide["index"],
                "title": slide.get("title", ""),
                "content": "\n".join(slide.get("body_texts", [])),
                "notes": slide.get("notes", ""),
                "has_images": slide.get("image_count", 0) > 0,
                "image_count": slide.get("image_count", 0),
                "images": slide.get("images", []),
                "missing_image_refs": slide.get("missing_image_refs", []),
                "tables": slide.get("tables", []),
                "body_texts": slide.get("body_texts", []),
                "raw_text": slide.get("raw_text", ""),
                "source_tex": slide.get("source_tex", ""),
                "source_body_tex": slide.get("source_body_tex", ""),
                "source_start": slide.get("source_start"),
                "source_end": slide.get("source_end"),
                "layout": slide.get("layout", {}),
            }
        )

    return {
        "chapter_title": first_title or "未命名PPT",
        "chapter_content": parse_result.get("full_text", ""),
        "slide_details": slide_details,
        "total_slides": len(slides),
    }

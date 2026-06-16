from __future__ import annotations

import asyncio
import json
import logging
import os
import re
import base64
import shutil
import subprocess
import tempfile
import uuid
import zipfile
from datetime import datetime
from io import BytesIO
from pathlib import Path
from typing import List, Optional

from fastapi import APIRouter, Body, File, HTTPException, UploadFile
from fastapi.responses import FileResponse, JSONResponse, Response, StreamingResponse
from pydantic import BaseModel, ConfigDict, Field

from .beamer_generator_full import config
from .beamer_generator_full.deepseek_client import DeepSeekClient
from .beamer_generator_full.latex_parser import parse_latex_to_slides
from .beamer_generator_full.pptx_generator import generate_pptx
from .beamer_generator_full.prompt_engine import PromptEngine

try:
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE
except Exception:  # pragma: no cover - optional dependency is checked at runtime.
    Presentation = None
    MSO_SHAPE_TYPE = None


logger = logging.getLogger(__name__)

BASE_DIR = Path(__file__).resolve().parent / "beamer_generator_full"
STATIC_DIR = BASE_DIR / "static"
UPLOAD_DIR = BASE_DIR / "uploads"
UPLOAD_DIR.mkdir(parents=True, exist_ok=True)
SAVED_PROJECT_DIR = BASE_DIR / "saved_projects"
SAVED_PROJECT_DIR.mkdir(parents=True, exist_ok=True)
EQUATION_INDEX_PATH = BASE_DIR / "equation_index.json"

prompt_engine = PromptEngine(config.SYSTEM_PROMPT_PATH)
router = APIRouter(prefix="/beamer-generator", tags=["beamer-generator"])


class GenerateRequest(BaseModel):
    content: str = Field(..., min_length=1, max_length=50000)
    provider: str = Field(default="auto")
    api_key: str = Field(default="")
    style: str = Field(default="academic")
    custom_requirements: str = Field(default="", max_length=5000)
    slide_count: int = Field(default=7, ge=1, le=80)
    language: str = Field(default="title_terms_en_content_zh")
    base_url: str = Field(default="")
    model: str = Field(default="")
    figure_assets: dict[str, str] = Field(default_factory=dict)


class ParseRequest(BaseModel):
    latex: str = Field(..., min_length=1)
    filename: str = Field(default="presentation.tex")


class RenderLatexRequest(ParseRequest):
    asset_urls: dict[str, str] = Field(default_factory=dict)


class SlideImage(BaseModel):
    model_config = ConfigDict(extra="allow")

    path: str = ""
    x: float = 1.0
    y: float = 3.0
    width: float = 4.0


class SlideTextbox(BaseModel):
    model_config = ConfigDict(extra="allow")

    text: str = ""
    x: float = 0
    y: float = 0
    width: float = 260
    height: float = 96
    color: str = "#333333"
    bg: str = ""
    fontSize: float = 14
    align: str = "left"
    bold: bool = False
    italic: bool = False


class SlideCallout(BaseModel):
    model_config = ConfigDict(extra="allow")

    text: str = ""
    x: float = 130
    y: float = 180
    width: float = 250
    height: float = 92
    fontSize: float = 12
    align: str = "center"


class SlidePlaceholder(BaseModel):
    model_config = ConfigDict(extra="allow")

    type: str = "image"
    label: str = "图片占位"
    position: str = ""
    figure: str = ""
    asset: str = ""
    url: str = ""
    path: str = ""
    page: int = 0
    x: float = 570
    y: float = 150
    width: float = 245
    height: float = 230


class SlideData(BaseModel):
    model_config = ConfigDict(extra="allow")

    id: int = 0
    type: str = "content"
    title: str = ""
    subtitle: str = ""
    titleCredit: str = ""
    reviewBackground: bool = False
    backgroundMode: str = ""
    renderedBackground: str = ""
    renderedBackgroundWidth: int = 0
    renderedBackgroundHeight: int = 0
    latexRenderedPage: bool = False
    hideParsedContent: bool = False
    items: List[str] = []
    equations: List[str] = []
    missing_equations: List[dict] = Field(default_factory=list)
    table: Optional[dict] = None
    notes: str = ""
    images: List[SlideImage] = []
    placeholders: List[SlidePlaceholder] = []
    textboxes: List[SlideTextbox] = []
    callouts: List[SlideCallout] = []


class RenderedSlide(BaseModel):
    page_index: int = 0
    image: str = ""
    width: int = 0
    height: int = 0


class ExportRequest(BaseModel):
    title: str = "Presentation"
    subtitle: str = ""
    author: str = ""
    date: str = ""
    slides: List[SlideData] = []
    figure_assets: dict[str, str] = Field(default_factory=dict)
    rendered_slides: List[RenderedSlide] = Field(default_factory=list)
    missing_equations: List[dict] = Field(default_factory=list)


class SaveProjectRequest(ExportRequest):
    latex: str = ""
    chapter_id: str = ""
    chapter_title: str = ""


class OverleafPackageRequest(SaveProjectRequest):
    pass


def _safe_saved_project_id(value: str, fallback: str = "presentation") -> str:
    raw = (value or "").strip() or fallback
    safe = re.sub(r"[^A-Za-z0-9_.-]+", "-", raw).strip(".-")
    return safe[:80] or fallback


def _saved_project_path(project_id: str) -> Path:
    return SAVED_PROJECT_DIR / f"{_safe_saved_project_id(project_id)}.json"


def _saved_project_summary(path: Path) -> dict:
    try:
        data = json.loads(path.read_text(encoding="utf-8"))
    except Exception:
        return {}
    slides = data.get("slides") if isinstance(data.get("slides"), list) else []
    return {
        "chapter_id": data.get("chapter_id") or path.stem,
        "chapter_title": data.get("chapter_title") or data.get("title") or path.stem,
        "title": data.get("title") or "",
        "updated_at": data.get("updated_at") or "",
        "slide_count": len(slides),
        "missing_equation_count": len(data.get("missing_equations") or []),
        "slides": [
            {
                "page_index": idx,
                "title": (slide or {}).get("title") or f"Slide {idx + 1}",
                "type": (slide or {}).get("type") or "content",
            }
            for idx, slide in enumerate(slides)
        ],
    }


def _read_root_env_values() -> dict[str, str]:
    env_path = BASE_DIR.parents[1] / ".env"
    values: dict[str, str] = {}
    if not env_path.exists():
        return values
    for line in env_path.read_text(encoding="utf-8").splitlines():
        if not line.strip() or line.lstrip().startswith("#") or "=" not in line:
            continue
        key, value = line.split("=", 1)
        values[key.strip()] = value.strip()
    return values


def _llm_setting(name: str, fallback: str = "") -> str:
    file_values = _read_root_env_values()
    for key, value in file_values.items():
        if key.startswith(("GPT_", "OPENAI_", "DEEPSEEK_")) and value and not os.getenv(key):
            os.environ[key] = value
    return os.getenv(name) or file_values.get(name, fallback)


def _default_llm_api_key() -> str:
    return (
        _llm_setting("GPT_API_KEY")
        or _llm_setting("OPENAI_API_KEY")
    )


def _default_llm_base_url() -> str:
    return (
        _llm_setting("GPT_API_BASE")
        or _llm_setting("OPENAI_API_BASE")
        or "https://api.openai.com/v1"
    )


def _default_llm_model(kind: str = "default") -> str:
    return _llm_setting("GPT_MODEL") or "gpt-5.5"


@router.get("", include_in_schema=False)
@router.get("/", include_in_schema=False)
@router.get("/index.html", include_in_schema=False)
async def index() -> FileResponse:
    return FileResponse(
        str(STATIC_DIR / "index.html"),
        media_type="text/html; charset=utf-8",
        headers={"Cache-Control": "no-store, no-cache, must-revalidate, max-age=0"},
    )


@router.get("/api/health")
async def health() -> dict[str, str]:
    return {"status": "ok", "service": "beamer-generator"}


@router.post("/api/generate-outline")
async def generate_outline(payload: dict | None = Body(default=None)) -> JSONResponse:
    payload = payload or {}
    content = _payload_text(payload, "content").strip()
    if not content:
        raise HTTPException(status_code=400, detail="请先导入 .md/.markdown 知识图谱文件")
    api_key = _payload_text(payload, "api_key").strip() or _default_llm_api_key()
    if not api_key:
        raise HTTPException(status_code=400, detail="请输入本次生成使用的 GPT API Key")
    base_url = _payload_text(payload, "base_url").strip() or _default_llm_base_url()
    model = _payload_text(payload, "model").strip() or _default_llm_model("pro")
    style = _payload_text(payload, "style", "academic").strip() or "academic"
    custom_requirements = _payload_text(payload, "custom_requirements").strip()
    language = _payload_text(payload, "language", "title_terms_en_content_zh").strip() or "title_terms_en_content_zh"
    slide_count = max(1, min(_payload_int(payload, "slide_count", 7), 80))
    section_slide_min = max(1, min(_payload_int(payload, "section_slide_min", 1), 80))
    section_slide_max = max(section_slide_min, min(_payload_int(payload, "section_slide_max", 8), 80))
    figure_assets = _payload_map(payload, "figure_assets")
    selected_sections = _payload_selected_sections(payload)

    client = DeepSeekClient(api_key=api_key, base_url=base_url)
    try:
        markdown_sections = _filter_outline_sections_by_selection(
            _split_outline_markdown_sections(content),
            selected_sections,
        )
        if _should_partition_outline(content, slide_count, markdown_sections, selected_sections):
            outline, raw = await _generate_partitioned_outline(
                client=client,
                content=content,
                sections=markdown_sections,
                model=model,
                style=style,
                custom_requirements=custom_requirements,
                slide_count=slide_count,
                section_slide_min=section_slide_min,
                section_slide_max=section_slide_max,
                language=language,
                figure_assets=figure_assets,
            )
        else:
            outline, raw = await _generate_outline_once(
                client=client,
                content=content,
                model=model,
                style=style,
                custom_requirements=custom_requirements,
                slide_count=slide_count,
                section_slide_min=section_slide_min,
                section_slide_max=section_slide_max,
                language=language,
                figure_assets=figure_assets,
                selected_sections=selected_sections,
                timeout_seconds=180,
            )
        return JSONResponse({"success": True, "outline": outline, "raw": raw})
    except json.JSONDecodeError as exc:
        logger.exception("Outline JSON parse failed")
        raise HTTPException(status_code=502, detail=f"纪要 JSON 解析失败：{exc}") from exc
    except Exception as exc:
        logger.exception("Outline generation failed")
        raise HTTPException(status_code=502, detail=f"纪要生成失败：{exc}") from exc


@router.post("/api/regenerate-outline-section")
async def regenerate_outline_section(payload: dict | None = Body(default=None)) -> JSONResponse:
    payload = payload or {}
    content = _payload_text(payload, "content").strip()
    if not content:
        raise HTTPException(status_code=400, detail="请先导入 .md/.markdown 知识图谱文件")
    api_key = _payload_text(payload, "api_key").strip() or _default_llm_api_key()
    if not api_key:
        raise HTTPException(status_code=400, detail="请输入本次生成使用的 GPT API Key")
    section_id = _payload_text(payload, "section_id").strip()
    if not section_id:
        raise HTTPException(status_code=400, detail="缺少要刷新的大节编号")

    base_url = _payload_text(payload, "base_url").strip() or _default_llm_base_url()
    model = _payload_text(payload, "model").strip() or _default_llm_model("pro")
    style = _payload_text(payload, "style", "academic").strip() or "academic"
    custom_requirements = _payload_text(payload, "custom_requirements").strip()
    language = _payload_text(payload, "language", "title_terms_en_content_zh").strip() or "title_terms_en_content_zh"
    slide_count = max(1, min(_payload_int(payload, "slide_count", 1), 80))
    section_slide_min = max(1, min(_payload_int(payload, "section_slide_min", 1), 80))
    section_slide_max = max(section_slide_min, min(_payload_int(payload, "section_slide_max", 8), 80))
    figure_assets = _payload_map(payload, "figure_assets")

    sections = _split_outline_markdown_sections(content)
    source_section = _find_outline_source_section(sections, section_id)
    if not source_section:
        raise HTTPException(status_code=404, detail=f"未在当前 Markdown 中找到大节 {section_id}")

    client = DeepSeekClient(api_key=api_key, base_url=base_url)
    try:
        section_outline, raw = await _generate_outline_once(
            client=client,
            content=_compact_outline_section_content(source_section["content"], 5200),
            model=model,
            style=style,
            custom_requirements=_section_outline_requirements(
                custom_requirements,
                source_section,
                slide_count,
            ),
            slide_count=slide_count,
            section_slide_min=section_slide_min,
            section_slide_max=section_slide_max,
            language=language,
            figure_assets=figure_assets,
            selected_sections=[],
            max_tokens=_outline_section_max_tokens(slide_count),
            timeout_seconds=75,
        )
        generated_section = _select_generated_outline_section(section_outline, source_section["id"])
        section = _fit_outline_section_frame_count(generated_section, source_section, slide_count)
        return JSONResponse({"success": True, "section": section, "raw": raw})
    except Exception as exc:
        logger.exception("Outline section refresh failed")
        raise HTTPException(status_code=502, detail=f"大节纪要刷新失败：{exc}") from exc


async def _generate_outline_once(
    *,
    client: DeepSeekClient,
    content: str,
    model: str,
    style: str,
    custom_requirements: str,
    slide_count: int,
    section_slide_min: int,
    section_slide_max: int,
    language: str,
    figure_assets: dict[str, str],
    selected_sections: list[dict[str, str]] | None = None,
    max_tokens: int | None = None,
    timeout_seconds: int | None = None,
) -> tuple[dict, str]:
    system_prompt = prompt_engine.build_outline_system_prompt(
        style=style,
        custom_requirements=_append_section_slide_range_requirement(
            custom_requirements,
            section_slide_min,
            section_slide_max,
        ),
        slide_count=slide_count,
        language=language,
        figure_assets=figure_assets,
    )
    user_prompt = prompt_engine.build_outline_user_prompt(
        content,
        custom_requirements=_append_section_slide_range_requirement(
            custom_requirements,
            section_slide_min,
            section_slide_max,
        ),
        slide_count=slide_count,
        selected_sections=selected_sections or [],
    )
    async def collect_parts() -> list[str]:
        collected: list[str] = []
        async for chunk in client.stream_generate(
            system_prompt=system_prompt,
            user_prompt=user_prompt,
            model=model,
            max_tokens=max_tokens or min(config.MAX_TOKENS, 6000),
            temperature=0.35,
        ):
            collected.append(chunk)
        return collected

    if timeout_seconds and timeout_seconds > 0:
        parts = await asyncio.wait_for(collect_parts(), timeout=timeout_seconds)
    else:
        parts = await collect_parts()
    raw = "".join(parts)
    return _validate_outline(json.loads(_clean_json_response(raw))), raw


async def _generate_partitioned_outline(
    *,
    client: DeepSeekClient,
    content: str,
    sections: list[dict[str, str]],
    model: str,
    style: str,
    custom_requirements: str,
    slide_count: int,
    section_slide_min: int,
    section_slide_max: int,
    language: str,
    figure_assets: dict[str, str],
) -> tuple[dict, list[dict[str, object]]]:
    allocations = _allocate_outline_slide_counts(
        sections,
        slide_count,
        section_slide_min,
        section_slide_max,
    )
    merged_sections: list[dict] = []
    raw_sections: list[dict[str, object]] = []

    for section, target_count in zip(sections, allocations):
        section_requirements = _section_outline_requirements(
            _partitioned_outline_requirements(
                custom_requirements,
                [section],
                [target_count],
                section_slide_min,
                section_slide_max,
            ),
            section,
            target_count,
        )
        try:
            section_outline, raw = await _generate_outline_once(
                client=client,
                content=_compact_outline_section_content(section["content"], 5200),
                model=model,
                style=style,
                custom_requirements=section_requirements,
                slide_count=target_count,
                section_slide_min=section_slide_min,
                section_slide_max=section_slide_max,
                language=language,
                figure_assets=figure_assets,
                selected_sections=[],
                max_tokens=_outline_section_max_tokens(target_count),
                timeout_seconds=75,
            )
            generated_section = _select_generated_outline_section(section_outline, section["id"])
            raw_sections.append({
                "id": section["id"],
                "target_slide_count": target_count,
                "status": "generated",
                "raw": raw,
            })
        except Exception as exc:
            logger.warning("Outline section %s generation failed: %s", section.get("id"), exc)
            generated_section = {}
            raw_sections.append({
                "id": section["id"],
                "target_slide_count": target_count,
                "status": "fallback",
                "error": str(exc),
            })
        merged_sections.append(
            _fit_outline_section_frame_count(generated_section, section, target_count)
        )

    total = sum(len(section["frames"]) for section in merged_sections)
    outline = _validate_outline({
        "title": _infer_outline_title(content),
        "target_slide_count": total,
        "sections": merged_sections,
    })
    return outline, [{
        "mode": "section_by_section",
        "section_count": len(sections),
        "allocations": allocations,
    }, *raw_sections]


def _should_partition_outline(
    content: str,
    slide_count: int,
    sections: list[dict[str, str]],
    selected_sections: list[dict[str, str]],
) -> bool:
    return len(sections) >= 2 and (slide_count >= 20 or len(content) >= 60000)


def _split_outline_markdown_sections(content: str) -> list[dict[str, str]]:
    matches = list(re.finditer(r"^##\s+(.+?)\s*$", content or "", flags=re.MULTILINE))
    sections: list[dict[str, str]] = []
    seen_ids: set[str] = set()
    for index, match in enumerate(matches):
        heading = match.group(1).strip()
        id_match = re.search(r"(?:chapter\d+_)?(\d{3})\b", heading, flags=re.IGNORECASE)
        if not id_match:
            continue
        section_id = id_match.group(1)
        if section_id in seen_ids:
            continue
        start = match.start()
        end = matches[index + 1].start() if index + 1 < len(matches) else len(content)
        section_text = content[start:end].strip()
        title = re.sub(
            r"^(?:chapter\d+_)?\d{3}\s*(?:[-:：·•路–—]+)?\s*",
            "",
            heading,
            flags=re.IGNORECASE,
        ).strip()
        sections.append({
            "id": section_id,
            "title": title or f"Section {section_id}",
            "content": section_text,
        })
        seen_ids.add(section_id)
    return sections


def _filter_outline_sections_by_selection(
    sections: list[dict[str, str]],
    selected_sections: list[dict[str, str]],
) -> list[dict[str, str]]:
    selected_numbers = _selected_outline_section_numbers(selected_sections)
    if not selected_numbers:
        return sections
    filtered = [section for section in sections if section.get("id") in selected_numbers]
    return filtered or sections


def _selected_outline_section_numbers(selected_sections: list[dict[str, str]]) -> set[str]:
    numbers: set[str] = set()
    for section in selected_sections or []:
        source = " ".join(
            str(section.get(key) or "")
            for key in ("id", "title", "file")
        )
        for match in re.finditer(r"(?:chapter\d+_)?(\d{3})\b", source, flags=re.IGNORECASE):
            numbers.add(match.group(1))
    return numbers


def _allocate_outline_slide_counts(
    sections: list[dict[str, str]],
    target_count: int,
    min_per_section: int = 1,
    max_per_section: int = 8,
) -> list[int]:
    if not sections:
        return []
    section_count = len(sections)
    min_per_section = max(1, min(int(min_per_section or 1), 80))
    max_per_section = max(min_per_section, min(int(max_per_section or min_per_section), 80))
    effective_min = min_per_section if min_per_section * section_count <= target_count else 1
    if target_count <= section_count * effective_min:
        return [effective_min] * section_count

    allocations = [effective_min] * section_count
    remaining = target_count - sum(allocations)
    weights = [max(1, len(section.get("content") or "")) for section in sections]
    total_weight = sum(weights) or section_count
    exact = [remaining * weight / total_weight for weight in weights]
    floors = [int(value) for value in exact]
    for index, value in enumerate(floors):
        allocations[index] += min(value, max_per_section - effective_min)
    leftover = target_count - sum(allocations)
    order = sorted(range(section_count), key=lambda idx: exact[idx] - floors[idx], reverse=True)
    while leftover > 0 and any(value < max_per_section for value in allocations):
        changed = False
        for index in order:
            if allocations[index] >= max_per_section:
                continue
            allocations[index] += 1
            leftover -= 1
            changed = True
            if leftover <= 0:
                break
        if not changed:
            break
    relaxed_index = 0
    while leftover > 0:
        allocations[order[relaxed_index % section_count]] += 1
        leftover -= 1
        relaxed_index += 1
    return allocations


def _append_section_slide_range_requirement(
    custom_requirements: str,
    section_slide_min: int,
    section_slide_max: int,
) -> str:
    rule = (
        f"每个用户勾选的 Markdown 大节建议生成 {section_slide_min}-{section_slide_max} 页；"
        "若总页数目标与该范围冲突，以后端给出的具体 section.slide_count 分配为准。"
    )
    base = (custom_requirements or "").strip()
    return base + "\n" + rule if base else rule


def _build_partitioned_outline_content(
    content: str,
    sections: list[dict[str, str]],
    allocations: list[int],
) -> str:
    per_section_chars = max(320, min(700, 9000 // max(1, len(sections))))
    parts = [
        f"# {_infer_outline_title(content)}",
        "",
        "以下是按 001/002 大节压缩后的 Markdown 内容，用于先生成 PPT 纪要。",
        "最终 LaTeX 生成时会再次使用完整 Markdown 原文；本阶段只需要规划大节概要和每页 frame 主题。",
        "",
    ]
    for section, slide_count in zip(sections, allocations):
        parts.append(f"## chapter_section_{section['id']} - {section['title']}")
        parts.append(f"planned_slide_count: {slide_count}")
        parts.append(_compact_outline_section_content(section["content"], per_section_chars))
        parts.append("")
    return "\n".join(parts)


def _partitioned_outline_requirements(
    custom_requirements: str,
    sections: list[dict[str, str]],
    allocations: list[int],
    section_slide_min: int,
    section_slide_max: int,
) -> str:
    allocation_lines = [
        f"- {section['id']}：{section['title']}，{count} 个 frame"
        for section, count in zip(sections, allocations)
    ]
    rules = [
        custom_requirements.strip(),
        "本次是大文件紧凑纪要规划。必须严格按下面的大节和页数分配输出 JSON：",
        "\n".join(allocation_lines),
        f"用户设置的每个勾选大节页数范围是 {section_slide_min}-{section_slide_max} 页；下面的具体分配已经按该范围估算。",
        "sections 数量、顺序、id、title 必须与上面完全一致。",
        "每个 section 的 slide_count 必须等于对应 frame 数量。",
        "所有 frame 总数必须等于上述分配总和；每个 frame 对应后续 LaTeX 的一页。",
        "不要额外新增 section，不要省略任何 001/002 大节。",
    ]
    return "\n".join(rule for rule in rules if rule)


def _align_partitioned_outline_sections(
    generated_outline: dict,
    source_sections: list[dict[str, str]],
    allocations: list[int],
) -> list[dict]:
    return [
        _fit_outline_section_frame_count(
            _select_generated_outline_section(generated_outline, section["id"]),
            section,
            target_count,
        )
        for section, target_count in zip(source_sections, allocations)
    ]


def _section_outline_requirements(
    custom_requirements: str,
    section: dict[str, str],
    target_count: int,
) -> str:
    rules = [
        custom_requirements.strip(),
        (
            f"本次只规划 Markdown 大节 {section['id']}：{section['title']}。"
            f"输出 JSON 中只能包含这一个 section，section.id 必须是 {section['id']}。"
        ),
        f"该大节必须生成恰好 {target_count} 个 frame；slide_count 必须等于 {target_count}。",
        "每个 frame 是后续 LaTeX 的一页，必须给出具体主题、中文概要和关键要点。",
    ]
    return "\n".join(rule for rule in rules if rule)


def _compact_outline_section_content(content: str, max_chars: int = 18000) -> str:
    text = (content or "").strip()
    if len(text) <= max_chars:
        return text
    head = text[: int(max_chars * 0.72)].rstrip()
    tail = text[-int(max_chars * 0.20):].lstrip()
    return head + "\n\n[...section content truncated for outline planning...]\n\n" + tail


def _outline_section_max_tokens(target_count: int) -> int:
    return min(config.MAX_TOKENS, max(1200, 800 + target_count * 220))


def _select_generated_outline_section(outline: dict, section_id: str) -> dict:
    sections = outline.get("sections") if isinstance(outline, dict) else []
    if not isinstance(sections, list) or not sections:
        return {}
    for section in sections:
        if isinstance(section, dict) and str(section.get("id") or "").strip() == section_id:
            return section
    return sections[0] if isinstance(sections[0], dict) else {}


def _find_outline_source_section(sections: list[dict[str, str]], section_id: str) -> dict[str, str] | None:
    normalized = _normalize_outline_section_id(section_id)
    for section in sections:
        if _normalize_outline_section_id(section.get("id") or "") == normalized:
            return section
    return None


def _normalize_outline_section_id(value: str) -> str:
    text = str(value or "").strip()
    match = re.search(r"(?:chapter\d+_)?(\d{3})\b", text, flags=re.IGNORECASE)
    return match.group(1) if match else text


def _fit_outline_section_frame_count(
    generated_section: dict,
    source_section: dict[str, str],
    target_count: int,
) -> dict:
    frames = generated_section.get("frames") if isinstance(generated_section, dict) else []
    normalized_frames: list[dict] = []
    if isinstance(frames, list):
        for index, frame in enumerate(frames[:target_count], start=1):
            if not isinstance(frame, dict):
                continue
            normalized_frames.append({
                "title": str(frame.get("title") or f"{source_section['title']} {index}").strip(),
                "summary": str(frame.get("summary") or "").strip(),
                "key_points": [
                    str(point).strip()
                    for point in (frame.get("key_points") if isinstance(frame.get("key_points"), list) else [])
                    if str(point).strip()
                ],
            })

    while len(normalized_frames) < target_count:
        index = len(normalized_frames) + 1
        normalized_frames.append({
            "title": f"{source_section['title']} Expansion {index}",
            "summary": "根据该大节 Markdown 内容补充展开成一页教学 frame。",
            "key_points": [
                "围绕本大节核心概念展开",
                "结合原文中的定义、公式、图示或例子",
                "保持与前后 frame 的教学顺序一致",
            ],
        })

    return {
        "id": source_section["id"],
        "title": source_section["title"],
        "summary": str(generated_section.get("summary") or "").strip()
        or f"{source_section['title']} 的内容概要。",
        "slide_count": target_count,
        "frames": normalized_frames,
    }


def _infer_outline_title(content: str) -> str:
    match = re.search(r"^#\s+(.+?)\s*$", content or "", flags=re.MULTILINE)
    if match:
        return match.group(1).strip()
    return "Presentation"


def _infer_chapter_directory_title(content: str) -> str:
    text = content or ""
    for pattern in (
        r"^#\s*(?:Chapter\s*)?(\d+)\s*(?:[·路\.-]\s*)?(.+?)\s*$",
        r"^##\s*chapter\d+_\d+\s*(?:[·路\.-]\s*)?(.+?)\s*$",
    ):
        match = re.search(pattern, text, flags=re.MULTILINE | re.IGNORECASE)
        if not match:
            continue
        if len(match.groups()) >= 2:
            number = match.group(1).strip()
            title = re.sub(r"\s+", " ", match.group(2)).strip()
            title = re.sub(r"^[^\w\u4e00-\u9fff]+", "", title).strip()
            if title:
                return f"{number} {title}"
        title = re.sub(r"\s+", " ", match.group(1)).strip()
        title = re.sub(r"^[^\w\u4e00-\u9fff]+", "", title).strip()
        if title:
            return title
    return ""



def _stream_error(message: str) -> StreamingResponse:
    async def event_stream():
        error_data = json.dumps({"type": "error", "content": message}, ensure_ascii=False)
        yield f"data: {error_data}\n\n"

    return StreamingResponse(
        event_stream(),
        media_type="text/event-stream",
        headers={"Cache-Control": "no-cache", "X-Accel-Buffering": "no"},
    )


def _clean_json_response(text: str) -> str:
    clean = (text or "").strip()
    clean = re.sub(r"^```(?:json)?\s*", "", clean, flags=re.I)
    clean = re.sub(r"\s*```$", "", clean)
    if "{" in clean and "}" in clean:
        clean = clean[clean.find("{"): clean.rfind("}") + 1]
    return clean.strip()


def _validate_outline(outline: dict) -> dict:
    if not isinstance(outline, dict):
        raise ValueError("outline must be a JSON object")
    sections = outline.get("sections")
    if not isinstance(sections, list) or not sections:
        raise ValueError("outline.sections must be a non-empty list")
    normalized_sections = []
    total = 0
    for index, section in enumerate(sections, start=1):
        if not isinstance(section, dict):
            continue
        frames = section.get("frames")
        if not isinstance(frames, list):
            frames = []
        normalized_frames = []
        for frame_index, frame in enumerate(frames, start=1):
            if not isinstance(frame, dict):
                continue
            title = str(frame.get("title") or "").strip() or f"Frame {frame_index}"
            summary = str(frame.get("summary") or "").strip()
            key_points = frame.get("key_points") if isinstance(frame.get("key_points"), list) else []
            normalized_frames.append({
                "title": title,
                "summary": summary,
                "key_points": [str(point).strip() for point in key_points if str(point).strip()],
            })
        if not normalized_frames:
            normalized_frames.append({
                "title": str(section.get("title") or f"Section {index}").strip(),
                "summary": str(section.get("summary") or "").strip(),
                "key_points": [],
            })
        section_id = str(section.get("id") or f"{index:03d}").strip()
        normalized_section = {
            "id": section_id,
            "title": str(section.get("title") or f"Section {section_id}").strip(),
            "summary": str(section.get("summary") or "").strip(),
            "slide_count": len(normalized_frames),
            "frames": normalized_frames,
        }
        total += len(normalized_frames)
        normalized_sections.append(normalized_section)
    if not normalized_sections:
        raise ValueError("outline.sections has no valid sections")
    return {
        "title": str(outline.get("title") or "Presentation").strip(),
        "target_slide_count": int(outline.get("target_slide_count") or total),
        "sections": normalized_sections,
    }


def _payload_text(payload: dict, key: str, default: str = "") -> str:
    value = payload.get(key, default)
    return str(value if value is not None else default)


def _payload_int(payload: dict, key: str, default: int = 0) -> int:
    try:
        return int(payload.get(key, default) or default)
    except (TypeError, ValueError):
        return default


def _payload_map(payload: dict, key: str) -> dict[str, str]:
    value = payload.get(key, {})
    if not isinstance(value, dict):
        return {}
    result: dict[str, str] = {}
    for k, v in value.items():
        k_text = str(k or "").strip()
        v_text = str(v or "").strip()
        if k_text and v_text:
            result[k_text] = v_text
    return result


def _payload_selected_sections(payload: dict) -> list[dict[str, str]]:
    value = payload.get("selected_sections", [])
    if not isinstance(value, list):
        return []
    result: list[dict[str, str]] = []
    for entry in value[:40]:
        if not isinstance(entry, dict):
            continue
        file_text = str(entry.get("file") or "").strip()
        title_text = str(entry.get("title") or "").strip()
        id_text = str(entry.get("id") or "").strip()
        if file_text or title_text:
            result.append({"file": file_text, "title": title_text, "id": id_text})
    return result


def _selected_sections_requirement(sections: list[dict[str, str]]) -> str:
    if not sections:
        return ""
    lines = [
        "本次只允许使用前端已引用的 Markdown 小节生成 PPT；不要使用同一 Markdown 文件中未引用的小节内容。",
        "已引用小节：",
    ]
    for section in sections:
        label = " / ".join(part for part in (section.get("file"), section.get("title")) if part)
        if label:
            lines.append(f"- {label}")
    return "\n".join(lines)


def _equation_reference_key(value: str, is_label: bool = False) -> str:
    text = re.sub(r"\s+", " ", str(value or "")).strip()
    if not text:
        return ""
    if is_label:
        return "label:" + text
    match = re.search(r"\d+(?:\.\d+)+|\d+", text)
    return "num:" + match.group(0) if match else "label:" + text


def _extract_equation_references(latex: str) -> list[dict[str, str]]:
    refs: list[dict[str, str]] = []
    seen: set[str] = set()

    def add(key: str, label: str, raw: str) -> None:
        if not key or key in seen:
            return
        seen.add(key)
        refs.append({"key": key, "label": label, "raw": raw})

    for match in re.finditer(r"\\(eqref|ref)\{([^{}]+)\}", latex or ""):
        command, label = match.group(1), match.group(2).strip()
        low = label.lower()
        if command == "ref" and not re.search(r"(^|[:_.-])(eq|equation|formula)([:_.-]|$)", low):
            continue
        if re.search(r"(^|[:_.-])(fig|figure|tab|table|sec|section|chap|chapter)([:_.-]|$)", low):
            continue
        add(_equation_reference_key(label, is_label=True), f"Equation {label}", match.group(0))

    text_ref_pattern = re.compile(
        r"((?:公式|方程|Equation|Eq\.?)\s*[\(（]?\s*(\d+(?:\.\d+)+|\d+)\s*[\)）]?)",
        re.IGNORECASE,
    )
    for match in text_ref_pattern.finditer(latex or ""):
        number = match.group(2)
        add(_equation_reference_key(number), f"Equation {number}", match.group(1))
    return refs


def _extract_equation_definition_keys(latex: str) -> set[str]:
    keys: set[str] = set()
    for label in re.findall(r"\\label\{([^{}]+)\}", latex or ""):
        keys.add(_equation_reference_key(label, is_label=True))
        number = re.search(r"\d+(?:\.\d+)+|\d+", label)
        if number:
            keys.add(_equation_reference_key(number.group(0)))
    for tag in re.findall(r"\\tag\{([^{}]+)\}", latex or ""):
        keys.add(_equation_reference_key(tag))
    return keys


def _missing_equation_macro() -> str:
    return (
        "\\providecommand{\\kgmissingequation}[2]{"
        "\\textcolor{red}{\\textbf{缺失公式：#2。请导入包含该公式的章节后自动补全。}}}\n"
    )


def _safe_image_macros() -> str:
    return (
        "\\providecommand{\\safelogoimage}[1]{\\fbox{\\parbox[c][30pt][c]{65pt}{\\centering\\tiny Logo}}}\n"
        "\\providecommand{\\safecontentimage}[1]{\\IfFileExists{#1}{\\includegraphics[width=0.7\\textwidth]{#1}}{\\fbox{\\parbox[c][0.34\\textheight][c]{0.7\\textwidth}{\\centering\\scriptsize Missing image\\\\\\texttt{\\detokenize{#1}}}}}}\n"
        "\\providecommand{\\safeverticalimage}[1]{\\IfFileExists{#1}{\\includegraphics[width=\\textwidth]{#1}}{\\fbox{\\parbox[c][0.34\\textheight][c]{\\textwidth}{\\centering\\scriptsize Missing image\\\\\\texttt{\\detokenize{#1}}}}}}\n"
    )


def _ensure_safe_image_macros(latex: str) -> str:
    if re.search(r"\\(?:providecommand|newcommand|renewcommand)\{\\safecontentimage\}", latex or ""):
        return latex
    begin_doc = r"\begin{document}"
    idx = (latex or "").find(begin_doc)
    if idx >= 0:
        return latex[:idx] + _safe_image_macros() + latex[idx:]
    return _safe_image_macros() + (latex or "")


def _ensure_missing_equation_macro(latex: str) -> str:
    if re.search(r"\\(?:providecommand|newcommand|renewcommand)\{\\kgmissingequation\}", latex or ""):
        return latex
    begin_doc = r"\begin{document}"
    idx = latex.find(begin_doc)
    if idx != -1:
        return latex[:idx] + _missing_equation_macro() + latex[idx:]
    return _missing_equation_macro() + (latex or "")


def _missing_equation_block(ref: dict[str, str]) -> str:
    key = ref.get("key", "")
    label = ref.get("label", key)
    return (
        "\\begin{alertblock}{缺失公式}\n"
        f"\\kgmissingequation{{{key}}}{{{label}}}\n"
        "\\end{alertblock}"
    )


def _mark_missing_equations(latex: str) -> tuple[str, list[dict[str, str]]]:
    text = latex or ""
    refs = _extract_equation_references(text)
    defined = _extract_equation_definition_keys(text)
    missing = [ref for ref in refs if ref["key"] not in defined]
    if not missing:
        return text, []

    text = _ensure_missing_equation_macro(text)
    frames = re.findall(r"\\begin\{frame\}(?:\[[^\]]*\])?(?:\{[^{}]*\})?[\s\S]*?\\end\{frame\}", text)
    appended: set[str] = set()
    for frame in frames:
        frame_missing = [
            ref for ref in missing
            if ref["key"] not in appended and ref.get("raw") and ref["raw"] in frame
        ]
        if not frame_missing:
            continue
        addition = "\n\n" + "\n\n".join(_missing_equation_block(ref) for ref in frame_missing)
        replacement = frame.replace(r"\end{frame}", addition + "\n\\end{frame}", 1)
        text = text.replace(frame, replacement, 1)
        appended.update(ref["key"] for ref in frame_missing)

    remaining = [ref for ref in missing if ref["key"] not in appended]
    if remaining:
        frame = (
            "\\begin{frame}{Missing Equations}\n"
            + "\n\n".join(_missing_equation_block(ref) for ref in remaining)
            + "\n\\end{frame}"
        )
        text = _insert_frames_before_end_document(text, [frame])
    return text, missing


def _extract_equations_from_source(source: str) -> dict[str, str]:
    equations: dict[str, str] = {}
    patterns = [
        r"\\begin\{(equation|align|alignat|gather|multline)\*?\}[\s\S]*?\\end\{\1\*?\}",
        r"\\\[[\s\S]*?\\\]",
        r"\$\$[\s\S]*?\$\$",
    ]
    for pattern in patterns:
        for match in re.finditer(pattern, source or ""):
            block = match.group(0).strip()
            keys = _extract_equation_definition_keys(block)
            before = (source or "")[max(0, match.start() - 120):match.start()]
            after = (source or "")[match.end():match.end() + 120]
            context = before + " " + after
            for ref in _extract_equation_references(context):
                keys.add(ref["key"])
            for key in keys:
                equations.setdefault(key, block)
    return equations


def _resolve_missing_equations_in_latex(latex: str, source: str) -> tuple[str, list[str], list[str]]:
    equations = _extract_equations_from_source(source)
    resolved: list[str] = []
    missing: list[str] = []
    text = latex or ""
    marker_pattern = re.compile(
        r"\\begin\{alertblock\}\{缺失公式\}\s*\\kgmissingequation\{([^{}]+)\}\{([^{}]*)\}\s*\\end\{alertblock\}",
        re.DOTALL,
    )

    def replace(match: re.Match[str]) -> str:
        key = match.group(1)
        label = match.group(2) or key
        formula = equations.get(key)
        if formula:
            resolved.append(label)
            return formula
        missing.append(label)
        return match.group(0)

    return marker_pattern.sub(replace, text), resolved, missing


def _normalize_figure_label(value: str) -> str:
    text = str(value or "")
    match = re.search(r"figure\s*\d+(?:\.\d+)?", text, re.IGNORECASE)
    if match:
        return re.sub(r"\s+", " ", match.group(0)).strip().lower()
    match = re.search(r"(?:fig\.?|图)\s*\d+(?:[._]\d+)?", text, re.IGNORECASE)
    if match:
        return re.sub(r"\s+", " ", match.group(0)).strip().lower()
    return re.sub(r"\s+", " ", text).strip().lower()


def _figure_number_filename(value: str) -> str:
    text = str(value or "")
    match = re.search(r"(?:figure|fig\.?|图)\s*(\d+(?:[._]\d+)?)", text, re.IGNORECASE)
    if not match:
        return ""
    number = match.group(1).replace(".", "_")
    return f"图{number}.png"


def _fig_path_for_label(label: str, figure_assets: dict[str, str]) -> str:
    key = _normalize_figure_label(label)
    for asset_label, asset in (figure_assets or {}).items():
        if _normalize_figure_label(asset_label) == key:
            filename = Path(str(asset or "").split("?", 1)[0].split("#", 1)[0].replace("\\", "/")).name
            if filename:
                return f"fig/{filename}"
    filename = _figure_number_filename(label)
    return f"fig/{filename}" if filename else ""


def _normalize_figure_label(value: str) -> str:
    text = str(value or "")
    match = re.search(r"(?:figure|fig\.?|图)\s*(\d+(?:[._]\d+)?)", text, re.IGNORECASE)
    if match:
        return "figure " + match.group(1).replace("_", ".")
    return re.sub(r"\s+", " ", text).strip().lower()


def _figure_number_filename(value: str) -> str:
    text = str(value or "")
    match = re.search(r"(?:figure|fig\.?|图)\s*(\d+(?:[._]\d+)?)", text, re.IGNORECASE)
    if not match:
        return ""
    number = match.group(1).replace(".", "_")
    return f"图{number}.png"


def _strip_missing_equation_markers(latex: str) -> str:
    text = latex or ""
    text = re.sub(
        r"\\begin\{alertblock\}\{(?:缺失公式|Missing Equation|Missing Equations|缂哄け鍏紡)\}"
        r"[\s\S]*?\\end\{alertblock\}",
        "",
        text,
    )
    text = re.sub(r"\\kgmissingequation\{[^{}]*\}\{[^{}]*\}", "", text)
    return _strip_missing_equation_macro_definition(text)


def _strip_missing_equation_macro_definition(text: str) -> str:
    source = text or ""
    pattern = re.compile(r"\\(?:providecommand|newcommand|renewcommand)\{\\kgmissingequation\}(?:\[\d+\])?\{")
    while True:
        match = pattern.search(source)
        if not match:
            return source
        open_pos = match.end() - 1
        depth = 0
        end = None
        for idx in range(open_pos, len(source)):
            char = source[idx]
            if char == "{":
                depth += 1
            elif char == "}":
                depth -= 1
                if depth == 0:
                    end = idx + 1
                    while end < len(source) and source[end] in " \t\r\n":
                        end += 1
                    break
        if end is None:
            return source[:match.start()]
        source = source[:match.start()] + source[end:]


def _extract_equation_references(latex: str) -> list[dict[str, str]]:
    refs: list[dict[str, str]] = []
    seen: set[str] = set()
    source = _strip_missing_equation_markers(latex)

    def add(key: str, label: str, raw: str) -> None:
        if not key or key in seen:
            return
        seen.add(key)
        refs.append({"key": key, "label": label, "raw": raw})

    for match in re.finditer(r"\\(eqref|ref|autoref|cref|Cref)\{([^{}]+)\}", source):
        command, label = match.group(1), match.group(2).strip()
        low = label.lower()
        if command == "ref" and not re.search(r"(^|[:_.-])(eq|equation|formula)([:_.-]|$)", low):
            continue
        if re.search(r"(^|[:_.-])(fig|figure|tab|table|sec|section|chap|chapter)([:_.-]|$)", low):
            continue
        add(_equation_reference_key(label, is_label=True), f"Equation {label}", match.group(0))

    text_ref_pattern = re.compile(
        r"((?:公式|方程|Equation|Equations|Eq\.?|Eqs\.?)\s*[\(（]?\s*(\d+(?:\.\d+)+|\d+)\s*[\)）]?)",
        re.IGNORECASE,
    )
    for match in text_ref_pattern.finditer(source):
        number = match.group(2)
        add(_equation_reference_key(number), f"Equation {number}", match.group(1))
    return refs


def _extract_equation_definition_keys(latex: str) -> set[str]:
    source = latex or ""
    keys: set[str] = set()
    equation_blocks = re.findall(
        r"\\begin\{(?:equation|align|alignat|gather|multline)\*?\}[\s\S]*?\\end\{(?:equation|align|alignat|gather|multline)\*?\}"
        r"|\\\[[\s\S]*?\\\]|\$\$[\s\S]*?\$\$",
        source,
    )
    blocks = equation_blocks or [source]
    for block in blocks:
        for label in re.findall(r"\\label\{([^{}]+)\}", block):
            low = label.lower()
            if re.search(r"(^|[:_.-])(fig|figure|tab|table|sec|section|chap|chapter)([:_.-]|$)", low):
                continue
            keys.add(_equation_reference_key(label, is_label=True))
            number = re.search(r"\d+(?:\.\d+)+|\d+", label)
            if number:
                keys.add(_equation_reference_key(number.group(0)))
        for tag in re.findall(r"\\tag\{([^{}]+)\}", block):
            keys.add(_equation_reference_key(tag))
    return keys


def _frame_contains_display_formula(frame: str) -> bool:
    return bool(re.search(
        r"\\begin\{(?:equation|align|alignat|gather|multline)\*?\}[\s\S]*?\\end\{(?:equation|align|alignat|gather|multline)\*?\}"
        r"|\\\[[\s\S]*?\\\]|\$\$[\s\S]*?\$\$",
        frame or "",
    ))


def _reference_is_present_in_own_frame(latex: str, ref: dict[str, str]) -> bool:
    raw = str(ref.get("raw") or "")
    if not raw:
        return False
    frames = re.findall(r"\\begin\{frame\}(?:\[[^\]]*\])?(?:\{[^{}]*\})?[\s\S]*?\\end\{frame\}", latex or "")
    for frame in frames:
        if raw in frame and _frame_contains_display_formula(frame):
            return True
    return False


def _collect_unresolved_equation_references(
    latex: str,
    known_equations: dict[str, str] | None = None,
) -> tuple[list[dict[str, str]], list[str]]:
    text = _strip_missing_equation_markers(latex or "")
    refs = _extract_equation_references(text)
    defined = _extract_equation_definition_keys(text)
    known = known_equations or {}
    missing: list[dict[str, str]] = []
    resolved: list[str] = []
    seen: set[str] = set()
    for ref in refs:
        key = ref.get("key", "")
        if not key or key in seen:
            continue
        seen.add(key)
        if key in defined or _reference_is_present_in_own_frame(text, ref):
            continue
        if key in known:
            resolved.append(ref.get("label") or key)
            continue
        missing.append(ref)
    return missing, resolved


def _missing_equation_macro() -> str:
    return (
        "\\providecommand{\\kgmissingequation}[2]{\n"
        "  \\textcolor{red}{\\textbf{缺失公式：#2。请导入包含该公式的章节后自动补全。}}\n"
        "}\n"
    )


def _missing_equation_block(ref: dict[str, str]) -> str:
    key = ref.get("key", "")
    label = ref.get("label", key)
    return (
        "\\begin{alertblock}{缺失公式}\n"
        f"\\kgmissingequation{{{key}}}{{{label}}}\n"
        "\\end{alertblock}"
    )


def _missing_equation_marker_pattern() -> re.Pattern[str]:
    return re.compile(
        r"\\begin\{alertblock\}\{(?:缺失公式|Missing Equation|Missing Equations|缂哄け鍏紡)\}\s*"
        r"\\kgmissingequation\{([^{}]+)\}\{([^{}]*)\}\s*"
        r"\\end\{alertblock\}",
        re.DOTALL,
    )


def _mark_missing_equations(latex: str) -> tuple[str, list[dict[str, str]]]:
    text = latex or ""
    refs = _extract_equation_references(text)
    defined = _extract_equation_definition_keys(text)
    already_marked = {match.group(1) for match in _missing_equation_marker_pattern().finditer(text)}
    missing = [ref for ref in refs if ref["key"] not in defined and ref["key"] not in already_marked]
    if not missing:
        return text, []

    text = _ensure_missing_equation_macro(text)
    frames = re.findall(r"\\begin\{frame\}(?:\[[^\]]*\])?(?:\{[^{}]*\})?[\s\S]*?\\end\{frame\}", text)
    appended: set[str] = set()
    for frame in frames:
        frame_missing = [
            ref for ref in missing
            if ref["key"] not in appended and ref.get("raw") and ref["raw"] in frame
        ]
        if not frame_missing:
            continue
        addition = "\n\n" + "\n\n".join(_missing_equation_block(ref) for ref in frame_missing)
        replacement = frame.replace(r"\end{frame}", addition + "\n\\end{frame}", 1)
        text = text.replace(frame, replacement, 1)
        appended.update(ref["key"] for ref in frame_missing)

    remaining = [ref for ref in missing if ref["key"] not in appended]
    if remaining:
        frame = (
            "\\begin{frame}{Missing Equations}\n"
            + "\n\n".join(_missing_equation_block(ref) for ref in remaining)
            + "\n\\end{frame}"
        )
        text = _insert_frames_before_end_document(text, [frame])
    return text, missing


def _extract_equations_from_source(source: str) -> dict[str, str]:
    equations: dict[str, str] = {}
    patterns = [
        r"\\begin\{(equation|align|alignat|gather|multline)\*?\}[\s\S]*?\\end\{\1\*?\}",
        r"\\\[[\s\S]*?\\\]",
        r"\$\$[\s\S]*?\$\$",
    ]
    text = source or ""
    for pattern in patterns:
        for match in re.finditer(pattern, text):
            block = match.group(0).strip()
            keys = _extract_equation_definition_keys(block)
            before = text[max(0, match.start() - 180):match.start()]
            after = text[match.end():match.end() + 180]
            context = before + " " + after
            for ref in _extract_equation_references(context):
                keys.add(ref["key"])
            for key in keys:
                equations.setdefault(key, block)
    return equations


def _load_equation_index() -> dict[str, list[dict[str, str]]]:
    try:
        data = json.loads(EQUATION_INDEX_PATH.read_text(encoding="utf-8"))
    except Exception:
        return {}
    if not isinstance(data, dict):
        return {}
    normalized: dict[str, list[dict[str, str]]] = {}
    for key, entries in data.items():
        if not isinstance(entries, list):
            continue
        cleaned = []
        for entry in entries:
            if isinstance(entry, dict) and entry.get("latex"):
                cleaned.append({
                    "latex": str(entry.get("latex") or ""),
                    "source_id": str(entry.get("source_id") or ""),
                    "source_title": str(entry.get("source_title") or ""),
                })
        if cleaned:
            normalized[str(key)] = cleaned
    return normalized


def _write_equation_index(index: dict[str, list[dict[str, str]]]) -> None:
    EQUATION_INDEX_PATH.write_text(json.dumps(index, ensure_ascii=False, indent=2), encoding="utf-8")


def _update_equation_index_from_source(source: str, source_id: str = "", source_title: str = "") -> dict[str, str]:
    equations = _extract_equations_from_source(source)
    if not equations:
        return {}
    index = _load_equation_index()
    for key, latex in equations.items():
        entries = index.setdefault(key, [])
        signature = re.sub(r"\s+", "", latex)
        if not any(re.sub(r"\s+", "", item.get("latex", "")) == signature for item in entries):
            entries.append({
                "latex": latex,
                "source_id": source_id,
                "source_title": source_title,
            })
    _write_equation_index(index)
    return equations


def _known_equation_map(extra_source: str = "", source_id: str = "", source_title: str = "") -> dict[str, str]:
    known: dict[str, str] = {}
    for key, entries in _load_equation_index().items():
        if entries:
            known[key] = entries[0].get("latex", "")
    for path in SAVED_PROJECT_DIR.glob("*.json"):
        try:
            data = json.loads(path.read_text(encoding="utf-8"))
        except Exception:
            continue
        latex = str(data.get("latex") or "")
        if latex:
            for key, formula in _extract_equations_from_source(latex).items():
                known.setdefault(key, formula)
    if extra_source:
        known.update(_update_equation_index_from_source(extra_source, source_id, source_title))
    return known


def _attach_equation_number_to_formula(formula: str, key: str) -> str:
    text = (formula or "").strip()
    match = re.match(r"num:(\d+(?:\.\d+)*|\d+)$", key or "")
    if not text or not match or r"\tag{" in text or r"\label{" in text:
        return text
    number = match.group(1)
    display_match = re.fullmatch(r"\\\[\s*([\s\S]*?)\s*\\\]", text)
    if display_match:
        body = display_match.group(1).strip()
        return "\\begin{equation}\n" + body + f"\n\\tag{{{number}}}\n\\end{{equation}}"
    dollar_match = re.fullmatch(r"\$\$\s*([\s\S]*?)\s*\$\$", text)
    if dollar_match:
        body = dollar_match.group(1).strip()
        return "\\begin{equation}\n" + body + f"\n\\tag{{{number}}}\n\\end{{equation}}"
    return text


def _resolve_missing_equations_in_latex(
    latex: str,
    source: str = "",
    known_equations: dict[str, str] | None = None,
    source_id: str = "",
    source_title: str = "",
) -> tuple[str, list[str], list[str]]:
    equations = dict(known_equations or {})
    if source:
        equations.update(_update_equation_index_from_source(source, source_id, source_title))
    resolved: list[str] = []
    missing: list[str] = []
    text = latex or ""

    def replace(match: re.Match[str]) -> str:
        key = match.group(1)
        label = match.group(2) or key
        formula = equations.get(key)
        if formula:
            resolved.append(label)
            return _attach_equation_number_to_formula(formula, key)
        missing.append(label)
        return match.group(0)

    return _missing_equation_marker_pattern().sub(replace, text), resolved, missing


def _append_equation_to_reference_frame(latex: str, ref: dict[str, str], formula: str) -> str:
    raw = str(ref.get("raw") or "")
    if not raw or not formula:
        return latex
    frames = re.findall(r"\\begin\{frame\}(?:\[[^\]]*\])?(?:\{[^{}]*\})?[\s\S]*?\\end\{frame\}", latex or "")
    for frame in frames:
        if raw not in frame:
            continue
        equation_block = "\n\n" + _attach_equation_number_to_formula(formula, ref.get("key", "")) + "\n"
        replacement = frame.replace(r"\end{frame}", equation_block + r"\end{frame}", 1)
        return latex.replace(frame, replacement, 1)
    return latex


def _apply_equation_reference_policy(
    latex: str,
    source: str = "",
    source_id: str = "",
    source_title: str = "",
) -> tuple[str, list[dict[str, str]], list[str]]:
    known = _known_equation_map(source, source_id, source_title)
    clean_latex = _strip_missing_equation_markers(latex or "")
    refs = _extract_equation_references(clean_latex)
    resolved: list[str] = []
    seen: set[str] = set()
    for ref in refs:
        key = ref.get("key", "")
        if not key or key in seen:
            continue
        seen.add(key)
        if key in _extract_equation_definition_keys(clean_latex) or _reference_is_present_in_own_frame(clean_latex, ref):
            continue
        formula = known.get(key)
        if not formula:
            continue
        clean_latex = _append_equation_to_reference_frame(clean_latex, ref, formula)
        resolved.append(ref.get("label") or key)

    marked_latex, missing = _mark_missing_equations(clean_latex)
    return marked_latex, missing, resolved


def _attach_missing_equations_to_slides(parsed: dict, latex: str, missing: list[dict[str, str]]) -> dict:
    if not parsed or not isinstance(parsed.get("slides"), list) or not missing:
        return parsed
    frames = _extract_latex_frames(latex or "")
    assigned: set[str] = set()
    for idx, frame in enumerate(frames):
        if idx >= len(parsed["slides"]):
            break
        slide_missing = []
        for ref in missing:
            key = ref.get("key", "")
            raw = ref.get("raw", "")
            if key in assigned:
                continue
            if raw and raw in frame:
                slide_missing.append(ref)
                assigned.add(key)
        if slide_missing:
            existing = parsed["slides"][idx].get("missing_equations") or []
            parsed["slides"][idx]["missing_equations"] = existing + slide_missing
    remaining = [ref for ref in missing if ref.get("key", "") not in assigned]
    if remaining and parsed["slides"]:
        existing = parsed["slides"][-1].get("missing_equations") or []
        parsed["slides"][-1]["missing_equations"] = existing + remaining
    return parsed


def _prepare_generated_image_paths(latex: str, figure_assets: dict[str, str]) -> str:
    text = latex or ""
    placeholder_pattern = re.compile(
        r"\\(?:kgimageplaceholder|imageplaceholder|pptimageplaceholder)\s*(?:\[([^\]]*)\])?\s*\{([^{}]*)\}",
        re.DOTALL,
    )

    def replace_placeholder(match: re.Match[str]) -> str:
        options = match.group(1) or ""
        label = match.group(2) or ""
        figure_key = _figure_key_from_placeholder(options, label)
        fig_path = _fig_path_for_label(figure_key or label or options, figure_assets)
        if not fig_path:
            return match.group(0)
        return "\\begin{center}\n  \\includegraphics[width=0.7\\textwidth]{" + fig_path + "}\n\\end{center}"

    text = placeholder_pattern.sub(replace_placeholder, text)

    include_pattern = re.compile(r"(\\includegraphics\s*(?:\[[^\]]*\])?\s*)\{([^{}]+)\}", re.DOTALL)

    def replace_include(match: re.Match[str]) -> str:
        prefix = match.group(1)
        target = match.group(2).strip()
        if target.startswith("fig/"):
            return match.group(0)
        fig_path = _fig_path_for_label(target, figure_assets)
        if not fig_path:
            return match.group(0)
        return f"{prefix}{{{fig_path}}}"

    return include_pattern.sub(replace_include, text)


def _caption_text_from_image_frame_body(body: str, image_command: str) -> str:
    text = (body or "").replace(image_command, " ")
    text = re.sub(r"\\begin\{(?:center|figure|itemize|enumerate|columns)\}(?:\[[^\]]*\])?", " ", text)
    text = re.sub(r"\\end\{(?:center|figure|itemize|enumerate|columns)\}", " ", text)
    text = re.sub(r"\\begin\{column\}\{[^{}]*\}", " ", text)
    text = re.sub(r"\\end\{column\}", " ", text)
    text = re.sub(r"\\(?:centering|caption)\*?(?:\[[^\]]*\])?", " ", text)
    text = re.sub(r"\\item\b", " ", text)
    text = re.sub(r"\\vspace\{[^{}]*\}", " ", text)
    text = re.sub(r"%.*", " ", text)
    text = re.sub(r"\s+", " ", text).strip()
    return text


def _default_image_caption(frame_title: str, image_path: str) -> str:
    label = (frame_title or _figure_frame_title_for_path(image_path, {}) or "本图").strip()
    if label.lower() == "figure":
        label = "本图"
    return (
        f"{_escape_latex_text(label)} 用于说明本节中的关键关系或变化趋势。"
        "阅读时先看图中的变量、坐标或模块，再比较主要曲线、区域或节点之间的差异，"
        "最后把这些变化联系到本页对应的核心概念和公式含义。"
    )


def _is_uninformative_image_caption(caption: str, frame_title: str, image_path: str) -> bool:
    text = re.sub(r"\s+", " ", str(caption or "")).strip()
    if not text:
        return True
    title = re.sub(r"\s+", " ", str(frame_title or "")).strip()
    if title and text.lower() == title.lower():
        return True
    normalized_text = _normalize_figure_label(text)
    normalized_title = _normalize_figure_label(title)
    normalized_path = _normalize_figure_label(Path(str(image_path or "")).stem)
    if normalized_text and normalized_text in {normalized_title, normalized_path}:
        return True
    if re.fullmatch(r"(?:Figure|Fig\.?|图)\s*\d+(?:\.\d+)*", text, re.IGNORECASE):
        return True
    return False


def _image_dimensions_from_bytes(data: bytes) -> tuple[int, int] | None:
    if data.startswith(b"\x89PNG\r\n\x1a\n") and len(data) >= 24:
        return int.from_bytes(data[16:20], "big"), int.from_bytes(data[20:24], "big")
    if data.startswith(b"\xff\xd8"):
        index = 2
        while index + 9 < len(data):
            if data[index] != 0xFF:
                index += 1
                continue
            marker = data[index + 1]
            index += 2
            if marker in {0xD8, 0xD9}:
                continue
            if index + 2 > len(data):
                break
            size = int.from_bytes(data[index:index + 2], "big")
            if size < 2 or index + size > len(data):
                break
            if marker in {0xC0, 0xC1, 0xC2, 0xC3, 0xC5, 0xC6, 0xC7, 0xC9, 0xCA, 0xCB, 0xCD, 0xCE, 0xCF}:
                height = int.from_bytes(data[index + 3:index + 5], "big")
                width = int.from_bytes(data[index + 5:index + 7], "big")
                return width, height
            index += size
    return None


def _asset_dimensions(value: str) -> tuple[int, int] | None:
    raw = str(value or "").strip()
    if raw.startswith("data:image/") and "," in raw:
        try:
            return _image_dimensions_from_bytes(base64.b64decode(raw.split(",", 1)[1], validate=False))
        except Exception:
            return None
    path = _resolve_uploaded_asset_path(raw)
    if not path:
        return None
    try:
        return _image_dimensions_from_bytes(path.read_bytes())
    except Exception:
        return None


def _image_orientation_for_path(image_path: str, figure_assets: dict[str, str] | None = None) -> str:
    candidates: list[str] = []
    normalized_path = str(image_path or "").replace("\\", "/")
    filename = Path(normalized_path).name.lower()
    figure_key = _normalize_figure_label(normalized_path)
    for label, asset in (figure_assets or {}).items():
        asset_name = Path(str(asset or "").split("?", 1)[0].split("#", 1)[0].replace("\\", "/")).name.lower()
        if (
            _normalize_figure_label(label) == figure_key
            or asset_name == filename
            or _normalize_figure_label(asset_name) == figure_key
        ):
            candidates.append(str(asset or ""))
    candidates.append(normalized_path)
    for candidate in candidates:
        dims = _asset_dimensions(candidate)
        if not dims:
            continue
        width, height = dims
        return "vertical" if height > width else "horizontal"
    return "horizontal"


def _figure_frame_title_for_path(image_path: str, figure_assets: dict[str, str] | None = None) -> str:
    normalized_path = str(image_path or "").replace("\\", "/")
    filename = Path(normalized_path).name
    for label, asset in (figure_assets or {}).items():
        label_text = str(label or "").strip()
        asset_name = Path(str(asset or "").split("?", 1)[0].split("#", 1)[0].replace("\\", "/")).name
        if label_text and (
            asset_name.lower() == filename.lower()
            or _normalize_figure_label(asset_name) == _normalize_figure_label(filename)
            or _normalize_figure_label(label_text) == _normalize_figure_label(filename)
        ):
            normalized = _normalize_figure_label(label_text)
            if normalized.startswith("figure "):
                return "Figure " + normalized.split(" ", 1)[1]
            return label_text

    match = re.search(r"(?:figure|fig\.?|图)?\s*(\d+(?:[._]\d+)?)", filename, re.IGNORECASE)
    if match:
        return "Figure " + match.group(1).replace("_", ".")
    return "Figure"


def _enforce_top_image_bottom_text_layout(latex: str, figure_assets: dict[str, str] | None = None) -> str:
    frame_pattern = re.compile(r"\\begin\{frame\}(?:\[[^\]]*\])?(?:\{[^{}]*\})?[\s\S]*?\\end\{frame\}")
    image_pattern = re.compile(
        r"\\(?P<cmd>includegraphics|safecontentimage|safeverticalimage)(?:\[[^\]]*\])?\{(?P<path>[^{}]+)\}"
    )

    def replace_frame(match: re.Match[str]) -> str:
        frame = match.group(0)
        if not re.search(r"\\(?:includegraphics|safecontentimage|safeverticalimage)", frame):
            return frame
        image_match = image_pattern.search(frame)
        if not image_match:
            return frame
        image_cmd = image_match.group("cmd")
        image_path = image_match.group("path").strip()
        body = re.sub(r"^\\begin\{frame\}(?:\[[^\]]*\])?(?:\{[^{}]*\})?", "", frame)
        body = re.sub(r"\\end\{frame\}\s*$", "", body)
        caption = _caption_text_from_image_frame_body(body, image_match.group(0))
        frame_title = _figure_frame_title_for_path(image_path, figure_assets or {})
        if _is_uninformative_image_caption(caption, frame_title, image_path):
            caption = _default_image_caption(frame_title, image_path)
        orientation = "vertical" if image_cmd == "safeverticalimage" else _image_orientation_for_path(image_path, figure_assets or {})
        if orientation == "vertical":
            return (
                f"\\begin{{frame}}{{{frame_title}}}\n"
                "  \\begin{columns}[T]\n"
                "    \\begin{column}{0.45\\textwidth}\n"
                f"      \\scriptsize {caption}\n"
                "    \\end{column}\n"
                "    \\begin{column}{0.45\\textwidth}\n"
                "      \\centering\n"
                f"      \\safeverticalimage{{{image_path}}}\n"
                "    \\end{column}\n"
                "  \\end{columns}\n"
                "\\end{frame}"
            )
        return (
            f"\\begin{{frame}}{{{frame_title}}}\n"
            "  \\centering\n"
            f"  \\safecontentimage{{{image_path}}}\n"
            "  \\vspace{0.3cm}\n"
            "  \\begin{center}\n"
            f"    \\parbox{{0.95\\textwidth}}{{\\scriptsize {caption}}}\n"
            "  \\end{center}\n"
            "\\end{frame}"
        )

    return frame_pattern.sub(replace_frame, latex or "")


def _image_path_from_frame(frame: str) -> str:
    patterns = [
        r"\\safecontentimage\{([^{}]+)\}",
        r"\\safeverticalimage\{([^{}]+)\}",
        r"\\includegraphics(?:\[[^\]]*\])?\{([^{}]+)\}",
    ]
    for pattern in patterns:
        match = re.search(pattern, frame or "")
        if match:
            return match.group(1).strip()
    return ""


def _image_merge_key(image_path: str) -> str:
    normalized = str(image_path or "").replace("\\", "/").strip()
    figure = _normalize_figure_label(normalized)
    if figure.startswith("figure "):
        return figure
    return normalized.lower()


def _extract_image_frame_caption(frame: str) -> str:
    parbox = re.search(
        r"\\parbox\{0\.95\\textwidth\}\{\\scriptsize\s*([\s\S]*?)\}\s*\\end\{center\}",
        frame or "",
    )
    if parbox:
        return re.sub(r"\s+", " ", parbox.group(1)).strip()
    column = re.search(r"\\scriptsize\s+([\s\S]*?)\s*\\end\{column\}", frame or "")
    if column:
        return re.sub(r"\s+", " ", column.group(1)).strip()
    return ""


def _append_image_frame_caption(frame: str, caption: str) -> str:
    addition = re.sub(r"\s+", " ", caption or "").strip()
    if not addition or addition in frame:
        return frame

    def replace_parbox(match: re.Match[str]) -> str:
        existing = match.group(1).strip()
        merged = existing + "；" + addition if existing else addition
        return f"\\parbox{{0.95\\textwidth}}{{\\scriptsize {merged}}}\n  \\end{{center}}"

    updated = re.sub(
        r"\\parbox\{0\.95\\textwidth\}\{\\scriptsize\s*([\s\S]*?)\}\s*\\end\{center\}",
        replace_parbox,
        frame,
        count=1,
    )
    if updated != frame:
        return updated

    def replace_column(match: re.Match[str]) -> str:
        existing = match.group(1).strip()
        merged = existing + "；" + addition if existing else addition
        return f"\\scriptsize {merged}\n    \\end{{column}}"

    return re.sub(
        r"\\scriptsize\s+([\s\S]*?)\s*\\end\{column\}",
        replace_column,
        frame,
        count=1,
    )


def _merge_duplicate_image_frames(latex: str) -> str:
    source = latex or ""
    frame_pattern = re.compile(r"\\begin\{frame\}(?:\[[^\]]*\])?(?:\{[^{}]*\})?[\s\S]*?\\end\{frame\}")
    parts: list[str] = []
    first_frame_part_by_image: dict[str, int] = {}
    last = 0
    for match in frame_pattern.finditer(source):
        parts.append(source[last:match.start()])
        frame = match.group(0)
        image_path = _image_path_from_frame(frame)
        key = _image_merge_key(image_path)
        if image_path and key:
            caption = _extract_image_frame_caption(frame)
            if key in first_frame_part_by_image:
                first_idx = first_frame_part_by_image[key]
                parts[first_idx] = _append_image_frame_caption(parts[first_idx], caption)
            else:
                first_frame_part_by_image[key] = len(parts)
                parts.append(frame)
        else:
            parts.append(frame)
        last = match.end()
    parts.append(source[last:])
    return "".join(parts)


def _resolve_uploaded_asset_path(value: str) -> Optional[Path]:
    raw = str(value or "").strip().strip('"').strip("'")
    if not raw or raw.startswith("data:"):
        return None

    normalized = raw.replace("\\", "/")
    marker = "/beamer-generator/uploads/"
    if normalized.startswith(("http://", "https://")):
        idx = normalized.find(marker)
        if idx == -1:
            return None
        normalized = normalized[idx + len(marker):]
    elif normalized.startswith(marker):
        normalized = normalized.split(marker, 1)[1]
    elif normalized.startswith("/uploads/"):
        normalized = normalized.split("/uploads/", 1)[1]
    elif os.path.isabs(raw):
        candidate = Path(raw).resolve()
        try:
            candidate.relative_to(UPLOAD_DIR.resolve())
        except ValueError:
            return None
        return candidate if candidate.exists() else None

    normalized = normalized.lstrip("./")
    if not normalized or normalized.startswith(".."):
        return None
    candidate = (UPLOAD_DIR / normalized).resolve()
    try:
        candidate.relative_to(UPLOAD_DIR.resolve())
    except ValueError:
        return None
    return candidate if candidate.exists() else None


def _safe_overleaf_filename(value: str, fallback: str, used: set[str]) -> str:
    raw_name = Path(str(value or "")).name
    stem = Path(raw_name).stem or fallback
    suffix = Path(raw_name).suffix.lower() or ".png"
    stem = re.sub(r'[<>:"/\\|?*\x00-\x1F]+', "_", stem).strip(" ._") or fallback
    candidate = f"{stem}{suffix}"
    index = 2
    while candidate.lower() in used:
        candidate = f"{stem}_{index}{suffix}"
        index += 1
    used.add(candidate.lower())
    return candidate


def _figure_key_from_placeholder(options: str, label: str) -> str:
    for piece in re.split(r"[,;]", options or ""):
        if "=" not in piece:
            continue
        key, value = piece.split("=", 1)
        if key.strip().lower() in {"figure", "fig", "image"}:
            normalized = _normalize_figure_label(value)
            if normalized:
                return normalized
    return _normalize_figure_label(options) or _normalize_figure_label(label)


def _prepare_overleaf_latex(
    latex: str,
    figure_assets: dict[str, str],
    asset_paths: dict[str, tuple[Path, str]],
) -> str:
    text = _strip_markdown_code_fence(latex or "")

    for label, asset in figure_assets.items():
        key = _normalize_figure_label(label)
        if key not in asset_paths:
            continue
        _, filename = asset_paths[key]
        fig_ref = f"fig/{filename}"
        for original in {str(asset or "").strip(), str(asset or "").replace("\\", "/").strip()}:
            if original:
                text = text.replace(original, fig_ref)

    placeholder_pattern = re.compile(
        r"\\(?:kgimageplaceholder|imageplaceholder|pptimageplaceholder)\s*(?:\[([^\]]*)\])?\s*\{([^{}]*)\}",
        re.DOTALL,
    )

    def replace_placeholder(match: re.Match[str]) -> str:
        options = match.group(1) or ""
        label = match.group(2) or ""
        figure_key = _figure_key_from_placeholder(options, label)
        asset_info = asset_paths.get(figure_key)
        if not asset_info:
            return match.group(0)
        _, filename = asset_info
        return "\\begin{center}\n  \\includegraphics[width=0.7\\textwidth]{fig/" + filename + "}\n\\end{center}"

    return placeholder_pattern.sub(replace_placeholder, text)


def _build_overleaf_package(req: OverleafPackageRequest) -> tuple[bytes, str]:
    latex = (req.latex or "").strip()
    if not latex:
        raise ValueError("没有可发送到 Overleaf 的 LaTeX 内容")

    figure_assets = req.figure_assets or {}
    used_names: set[str] = set()
    asset_paths: dict[str, tuple[Path, str]] = {}
    for index, (label, asset) in enumerate(figure_assets.items(), start=1):
        key = _normalize_figure_label(label)
        if not key:
            continue
        path = _resolve_uploaded_asset_path(str(asset or ""))
        if not path:
            continue
        filename = _safe_overleaf_filename(path.name, f"image_{index}", used_names)
        asset_paths[key] = (path, filename)

    prepared_latex = _prepare_overleaf_latex(latex, figure_assets, asset_paths)
    title = req.chapter_title or req.title or "presentation"
    zip_filename = _safe_saved_project_id(title, "presentation") + "_overleaf.zip"

    buffer = BytesIO()
    with zipfile.ZipFile(buffer, "w", compression=zipfile.ZIP_DEFLATED) as archive:
        archive.writestr("main.tex", prepared_latex)
        archive.writestr(
            "README.txt",
            "This ZIP was generated by the Knowledge Graph Teaching System for Overleaf.\n"
            "Open main.tex in Overleaf to edit and compile the Beamer presentation.\n",
        )
        for path, filename in asset_paths.values():
            archive.write(path, f"fig/{filename}")
    return buffer.getvalue(), zip_filename


def _render_pdf_bytes_to_pages(pdf_bytes: bytes, dpi: int = 168) -> list[dict]:
    try:
        import fitz  # type: ignore
    except Exception as exc:  # pragma: no cover - depends on local installation
        raise RuntimeError("服务器未安装 PyMuPDF，无法把 PDF 渲染为 PPT 背景图片") from exc

    pages: list[dict] = []
    doc = fitz.open(stream=pdf_bytes, filetype="pdf")
    zoom = max(0.5, min(4.0, dpi / 72.0))
    matrix = fitz.Matrix(zoom, zoom)
    try:
        for index, page in enumerate(doc):
            pix = page.get_pixmap(matrix=matrix, alpha=False)
            png_bytes = pix.tobytes("png")
            pages.append({
                "page_index": index,
                "image": "data:image/png;base64," + base64.b64encode(png_bytes).decode("ascii"),
                "width": pix.width,
                "height": pix.height,
            })
    finally:
        doc.close()
    return pages


def _latex_compiler_command(tex_filename: str = "main.tex") -> list[str] | None:
    candidate_dirs = [
        Path(os.environ.get("LOCALAPPDATA", "")) / "Programs" / "MiKTeX" / "miktex" / "bin" / "x64",
        Path(os.environ.get("ProgramFiles", "")) / "MiKTeX" / "miktex" / "bin" / "x64",
        Path(os.environ.get("ProgramFiles(x86)", "")) / "MiKTeX" / "miktex" / "bin" / "x64",
    ]
    for name in ("xelatex", "pdflatex"):
        exe = shutil.which(name)
        if exe:
            return [exe, "-interaction=nonstopmode", "-halt-on-error", tex_filename]
        exe_name = f"{name}.exe" if os.name == "nt" else name
        for directory in candidate_dirs:
            if not str(directory):
                continue
            candidate = directory / exe_name
            if candidate.exists():
                return [str(candidate), "-interaction=nonstopmode", "-halt-on-error", tex_filename]
    exe = shutil.which("tectonic")
    if exe:
        return [exe, tex_filename]
    return None


def _latex_graphic_targets(latex: str) -> list[str]:
    targets: list[str] = []
    for match in re.finditer(
        r"\\(?:includegraphics|safecontentimage|safeverticalimage)(?:\[[^\]]*\])?\{([^{}]+)\}",
        latex or "",
    ):
        target = match.group(1).strip().replace("\\", "/")
        if target and target not in targets:
            targets.append(target)
    return targets


def _materialize_latex_assets(temp_dir: Path, latex: str, asset_urls: dict[str, str] | None = None) -> None:
    lookup = {str(key).replace("\\", "/").lower(): str(value) for key, value in (asset_urls or {}).items()}
    for target in _latex_graphic_targets(latex):
        normalized = target.lstrip("./")
        candidates = [
            normalized,
            Path(normalized).name,
            "fig/" + Path(normalized).name,
        ]
        source_value = ""
        for candidate in candidates:
            source_value = lookup.get(candidate.lower(), "")
            if source_value:
                break
        source_path = _resolve_uploaded_asset_path(source_value) if source_value else _resolve_uploaded_asset_path(normalized)
        if not source_path:
            continue
        output_path = (temp_dir / normalized).resolve()
        try:
            output_path.relative_to(temp_dir.resolve())
        except ValueError:
            continue
        output_path.parent.mkdir(parents=True, exist_ok=True)
        shutil.copyfile(source_path, output_path)


def _compile_latex_to_pdf_bytes(latex: str, asset_urls: dict[str, str] | None = None) -> bytes:
    compiler = _latex_compiler_command("main.tex")
    if not compiler:
        raise RuntimeError("当前服务器未安装 xelatex / pdflatex / tectonic，不能直接把 .tex 编译成高保真页面。请先安装 TeX Live/MiKTeX，或在此处导入对应 PDF。")

    with tempfile.TemporaryDirectory(prefix="kg-beamer-render-") as temp_name:
        temp_dir = Path(temp_name)
        tex_path = temp_dir / "main.tex"
        tex_path.write_text(latex or "", encoding="utf-8")
        _materialize_latex_assets(temp_dir, latex or "", asset_urls)
        last_output = ""
        for _ in range(2):
            proc = subprocess.run(
                compiler,
                cwd=str(temp_dir),
                stdout=subprocess.PIPE,
                stderr=subprocess.STDOUT,
                text=True,
                encoding="utf-8",
                errors="replace",
                timeout=120,
                check=False,
            )
            last_output = proc.stdout or ""
            if proc.returncode != 0:
                tail = "\n".join(last_output.splitlines()[-24:])
                raise RuntimeError("LaTeX 编译失败：\n" + tail)
        pdf_path = temp_dir / "main.pdf"
        if not pdf_path.exists():
            raise RuntimeError("LaTeX 编译未生成 PDF 文件")
        return pdf_path.read_bytes()


def _compile_latex_file_to_pdf_bytes(tex_path: Path) -> bytes:
    tex_path = tex_path.resolve()
    compiler = _latex_compiler_command(tex_path.name)
    if not compiler:
        raise RuntimeError("当前服务器未安装 xelatex / pdflatex / tectonic，不能直接把 .tex 编译成高保真页面。请先安装 TeX Live/MiKTeX，或导入 Overleaf 编译后的 PDF。")

    last_output = ""
    for _ in range(2):
        proc = subprocess.run(
            compiler,
            cwd=str(tex_path.parent),
            stdout=subprocess.PIPE,
            stderr=subprocess.STDOUT,
            text=True,
            encoding="utf-8",
            errors="replace",
            timeout=120,
            check=False,
        )
        last_output = proc.stdout or ""
        if proc.returncode != 0:
            tail = "\n".join(last_output.splitlines()[-24:])
            raise RuntimeError("LaTeX 编译失败：\n" + tail)
    pdf_path = tex_path.with_suffix(".pdf")
    if not pdf_path.exists():
        raise RuntimeError("LaTeX 编译未生成 PDF 文件")
    return pdf_path.read_bytes()


def _count_generated_slides(latex: str) -> int:
    try:
        parsed = parse_latex_to_slides(latex or "")
        slides = parsed.get("slides") if isinstance(parsed, dict) else []
        if isinstance(slides, list):
            return len(slides)
    except Exception:
        logger.exception("Failed to count generated slides from LaTeX")
    return len(re.findall(r"\\begin\{frame\}", latex or ""))


def _extract_latex_frames(latex: str) -> list[str]:
    return re.findall(r"\\begin\{frame\}(?:\[[^\]]*\])?(?:\{[^{}]*\})?[\s\S]*?\\end\{frame\}", latex or "")


def _frame_title(frame: str) -> str:
    parsed_title = _frame_begin_title(frame)
    if parsed_title:
        return re.sub(r"\s+", " ", parsed_title).strip()
    match = re.search(r"\\frametitle\{([^{}]*)\}", frame or "")
    return re.sub(r"\s+", " ", match.group(1)).strip() if match else ""


def _balanced_group_span(text: str, open_pos: int) -> tuple[int, int] | None:
    if open_pos < 0 or open_pos >= len(text) or text[open_pos] != "{":
        return None
    depth = 0
    escaped = False
    for idx in range(open_pos, len(text)):
        char = text[idx]
        if escaped:
            escaped = False
            continue
        if char == "\\":
            escaped = True
            continue
        if char == "{":
            depth += 1
        elif char == "}":
            depth -= 1
            if depth == 0:
                return open_pos, idx + 1
    return None


def _frame_begin_title_span(frame: str) -> tuple[int, int, int, int] | None:
    text = frame or ""
    match = re.search(r"\\begin\{frame\}", text)
    if not match:
        return None
    pos = match.end()
    while pos < len(text) and text[pos].isspace():
        pos += 1
    if pos < len(text) and text[pos] == "[":
        close = text.find("]", pos + 1)
        if close == -1:
            return None
        pos = close + 1
        while pos < len(text) and text[pos].isspace():
            pos += 1
    if pos >= len(text) or text[pos] != "{":
        return None
    span = _balanced_group_span(text, pos)
    if not span:
        return None
    return span[0], span[1], span[0] + 1, span[1] - 1


def _frame_begin_title(frame: str) -> str:
    span = _frame_begin_title_span(frame)
    if not span:
        return ""
    return (frame or "")[span[2]:span[3]]


def _normalize_frame_body(frame: str) -> str:
    text = re.sub(r"\\begin\{frame\}(?:\[[^\]]*\])?(?:\{[^{}]*\})?", " ", frame or "")
    text = re.sub(r"\\frametitle\{[^{}]*\}", " ", text)
    text = re.sub(r"\\end\{frame\}", " ", text)
    text = re.sub(r"%.*", " ", text)
    text = re.sub(r"\\(?:begin|end)\{[^{}]*\}", " ", text)
    text = re.sub(r"\\[a-zA-Z]+\*?(?:\[[^\]]*\])?(?:\{([^{}]*)\})?", r" \1 ", text)
    text = re.sub(r"[{}\\$#&_^\[\](),.;:，。；：、（）]+", " ", text)
    return re.sub(r"\s+", " ", text).strip().lower()


def _frame_body_key(frame: str) -> str:
    body = _normalize_frame_body(frame)
    words = re.findall(r"[\u4e00-\u9fff]|[a-z0-9]+", body)
    return " ".join(words[:160])


def _frame_content_score(frame: str) -> int:
    body = _normalize_frame_body(frame)
    cjk_chars = len(re.findall(r"[\u4e00-\u9fff]", body))
    latin_words = len(re.findall(r"[a-z0-9]+", body))
    return cjk_chars + latin_words * 2


def _clean_content_anchor(value: str) -> str:
    text = re.sub(r"\s+", " ", str(value or "")).strip()
    text = re.sub(r"^#{1,6}\s*", "", text)
    text = re.sub(r"^[-*+]\s+", "", text)
    text = re.sub(r"^\d+[.)、]\s*", "", text)
    text = text.strip(" -:：;；,，。")
    if len(text) > 170:
        text = text[:170].rsplit(" ", 1)[0] or text[:170]
    return text.strip()


def _anchor_key(value: str) -> str:
    text = _clean_content_anchor(value).lower()
    text = re.sub(r"[`*_~#>\[\](){}，。；：、,.!?;:]+", " ", text)
    return re.sub(r"\s+", " ", text).strip()


def _extract_content_anchors(content: str, limit: int = 140) -> list[str]:
    anchors: list[str] = []
    seen: set[str] = set()

    def add(raw: str) -> None:
        clean = _clean_content_anchor(raw)
        if not clean:
            return
        score = len(re.findall(r"[\u4e00-\u9fff]", clean)) + len(re.findall(r"[A-Za-z0-9]+", clean)) * 2
        if score < 12:
            return
        key = _anchor_key(clean)
        if not key or key in seen:
            return
        seen.add(key)
        anchors.append(clean)

    for raw_line in (content or "").splitlines():
        line = raw_line.strip()
        if not line:
            continue
        if re.match(r"^#{1,6}\s+", line):
            add(line)
            continue
        if re.match(r"^[-*+]\s+", line) or re.match(r"^\d+[.)、]\s+", line):
            add(line)
            continue
        if re.search(r"(Figure|Fig\.?|图|表|Table)\s*[\d.]*", line, re.IGNORECASE):
            add(line)
            continue
        if re.search(r"(\$[^$]+\$|\\\(|\\\[|=|\\frac|\\sum|\\int|\\varphi|\\sigma)", line):
            add(line)
            continue
        if len(line) >= 24 and len(anchors) < limit // 2:
            add(line)
        if len(anchors) >= limit:
            break
    return anchors[:limit]


BIMSA_FIXED_SECTIONS = [
    "Optimal Selection Intensities For Maximizing Long-Term Response",
    "Effects Of Population Structure On Long-Term Response",
    "Asymptotic Response Due To Mutational Input",
]


def _frame_has_anchor(frame: str, anchors: list[str]) -> bool:
    if not anchors:
        return True
    body = _anchor_key(_normalize_frame_body(frame))
    for anchor in anchors:
        key = _anchor_key(anchor)
        if key and (key[:45] in body or body[:80] in key):
            return True
    return False


def _insert_frames_before_end_document(latex: str, frames: list[str]) -> str:
    if not frames:
        return latex
    end_tag = r"\end{document}"
    insert_at = latex.rfind(end_tag)
    addition = "\n\n" + "\n\n".join(frame.strip() for frame in frames if frame.strip()) + "\n\n"
    if insert_at == -1:
        return latex.rstrip() + addition + end_tag + "\n"
    return latex[:insert_at].rstrip() + addition + latex[insert_at:]


def _bimsa_latex_preamble() -> str:
    return r"""\documentclass[10pt, aspectratio=169]{ctexbeamer}
\usetheme{Madrid}

\usepackage{amsmath, amssymb, amsthm}
\usepackage{graphicx}
\usepackage{booktabs}
\usepackage{multirow}
\usepackage{caption}
\usepackage{hyperref}
\usepackage{tikz}
\usetikzlibrary{shapes.callouts, tikzmark}
\usetikzlibrary{shapes, positioning}

\newcommand{\E}{\mathrm{E}}
\newcommand{\bbE}{\mathbb{E}}
\newcommand{\bbR}{\mathbb{R}}
\newcommand{\cA}{\mathcal{A}}
\newcommand{\cB}{\mathcal{B}}
\newcommand{\cP}{\mathcal{P}}
\newcommand{\cF}{\mathcal{F}}

\definecolor{myline}{RGB}{0,116,112}
\definecolor{myblue}{RGB}{40,100,180}

\providecommand{\safelogoimage}[1]{\fbox{\parbox[c][30pt][c]{65pt}{\centering\tiny Logo}}}
\providecommand{\safecontentimage}[1]{\IfFileExists{#1}{\includegraphics[width=0.7\textwidth]{#1}}{\fbox{\parbox[c][0.34\textheight][c]{0.7\textwidth}{\centering\scriptsize Missing image\\\texttt{\detokenize{#1}}}}}}
\providecommand{\safeverticalimage}[1]{\IfFileExists{#1}{\includegraphics[width=\textwidth]{#1}}{\fbox{\parbox[c][0.34\textheight][c]{\textwidth}{\centering\scriptsize Missing image\\\texttt{\detokenize{#1}}}}}}

\setbeamertemplate{frametitle}{%
  \vspace*{0.2cm}%
  \begin{beamercolorbox}[wd=\paperwidth, leftskip=0.5cm, rightskip=0.5cm, ht=0.3cm, dp=0pt]{whitebg}%
    \usebeamerfont{frametitle}\textcolor{black}{\insertframetitle}%
  \end{beamercolorbox}%
  \vspace{0pt}%
  \begin{tikzpicture}[remember picture, overlay]
    \draw[myline, line width=1.5pt]
      ([yshift=-1.3cm] current page.north west) -- ([yshift=-1.3cm] current page.north east);
  \end{tikzpicture}%
  \vspace{0.1cm}%
}

\title[]{ Evolutionary Theory on\\ Polygenic Trait}
\subtitle{XII - Long-term Response: 2. Finite Population Size and Mutation (2)}
\author{Qi WU(吴琦)}
\date{2026-5-26}

\setbeamertemplate{title page}{%
  \begin{tikzpicture}[remember picture, overlay]
    \draw[line width=1.5pt, color=myline]
      ([yshift=-40pt] current page.north west) -- ([yshift=-40pt] current page.north east);
    \node[anchor=north west, inner sep=0, minimum width=0.25\paperwidth,
          minimum height=39pt, fill=gray!30, text=black, align=center]
          at (current page.north west) {Public course in BIMSA in 2026 spring semester};
    \node[anchor=north east, inner sep=0] (f3) at ([xshift=-3pt] current page.north east)
          {\safelogoimage{logo-right}};
    \node[anchor=north east, inner sep=0] at (f3.north west)
          {\safelogoimage{logo-left}};
  \end{tikzpicture}%
  \vspace*{36pt}
  \begin{center}
    \begin{tikzpicture}
      \node[draw=none, inner sep=8pt, fill=white, text=black,
            align=center, font=\Huge\bfseries] (titlebox) {\inserttitle};
      \node[draw=none, rounded corners=2pt,
            inner sep=8pt, fill=white, text=black,
            align=center, font=\large,
            below=5pt of titlebox] (subtitlebox) {\insertsubtitle};
      \node[below=5pt of subtitlebox.south east, anchor=north east, align=center, text=black] {
        \insertauthor \\[3pt]
        \insertdate
      };
    \end{tikzpicture}
  \end{center}
}
"""


def _bimsa_title_frame() -> str:
    return r"""\setbeamertemplate{footline}{}

{
\usebackgroundtemplate{%
  \begin{tikzpicture}[remember picture, overlay]
    \fill[gray!12] (current page.south west) rectangle (current page.north east);
  \end{tikzpicture}%
}
\begin{frame}
    \titlepage
\end{frame}
}
"""


def _safe_title_background_template() -> str:
    return (
        "\\usebackgroundtemplate{%\n"
        "  \\begin{tikzpicture}[remember picture, overlay]\n"
        "    \\fill[gray!12] (current page.south west) rectangle (current page.north east);\n"
        "  \\end{tikzpicture}%\n"
        "}"
    )


def _sanitize_unsafe_title_assets(latex: str) -> str:
    text = latex or ""
    text = text.replace("\t" + "itle", r"\title")
    text = text.replace("\b" + "egin{", r"\begin{")
    text = re.sub(
        r"(?m)^([ \t]*)itle(\s*(?:\[[^\]]*\])?\s*\{)",
        r"\1\\title\2",
        text,
    )
    text = re.sub(
        r"(?m)^([ \t]*)(title|subtitle|author|date)(\s*(?:\[[^\]]*\])?\s*\{)",
        r"\1\\\2\3",
        text,
    )
    text = re.sub(
        r"(?m)^([ \t]*)(begin|end)(\{(?:document|frame|center|itemize|enumerate|columns|column|tikzpicture|table|tabular|equation|align\*?|minipage)\})",
        r"\1\\\2\3",
        text,
    )
    text = re.sub(
        r"\\providecommand\{\\safelogoimage\}\[1\]\{[\s\S]*?\}\s*(?=\\providecommand\{\\safecontentimage\})",
        lambda _match: "\\providecommand{\\safelogoimage}[1]{\\fbox{\\parbox[c][30pt][c]{65pt}{\\centering\\tiny Logo}}}\n",
        text,
        count=1,
    )
    text = re.sub(
        r"\\usebackgroundtemplate\{\s*\\includegraphics\[[^\]]*\]\{fig/图片3\.png\}\s*\}",
        lambda _match: _safe_title_background_template(),
        text,
    )
    text = re.sub(
        r"\\setbeamertemplate\{footline\}\{[\s\S]*?fig/图片3\.png[\s\S]*?\n\}",
        lambda _match: _safe_title_background_template(),
        text,
    )
    return text


def _ensure_complete_bimsa_latex(latex: str) -> str:
    text = _strip_markdown_code_fence(latex or "")
    text = _sanitize_unsafe_title_assets(text)
    frames = _extract_latex_frames(text)
    has_document = r"\begin{document}" in text and r"\end{document}" in text
    if not has_document:
        body = "\n\n".join(frame.strip() for frame in frames)
        if not body:
            body = text.strip()
        return _bimsa_latex_preamble().rstrip() + "\n\n\\begin{document}\n\n" + _bimsa_title_frame().rstrip() + "\n\n" + body + "\n\n\\end{document}\n"

    if r"\documentclass" not in text:
        text = _bimsa_latex_preamble().rstrip() + "\n\n" + text
    if r"\titlepage" not in text:
        text = text.replace(r"\begin{document}", r"\begin{document}" + "\n\n" + _bimsa_title_frame().rstrip(), 1)
    return text


def _dedupe_latex_frames(latex: str) -> str:
    frames = _extract_latex_frames(latex)
    if not frames:
        return latex
    seen_titles: set[str] = set()
    seen_bodies: set[str] = set()
    replacements: dict[str, str] = {}
    for frame in frames:
        if r"\titlepage" in frame:
            continue
        title_key = _frame_title(frame).lower()
        body_key = _frame_body_key(frame)
        if (title_key and title_key in seen_titles) or (body_key and body_key in seen_bodies):
            replacements[frame] = ""
            continue
        if title_key:
            seen_titles.add(title_key)
        if body_key:
            seen_bodies.add(body_key)
    result = latex
    for frame, replacement in replacements.items():
        result = result.replace(frame, replacement, 1)
    return result


def _extract_display_formula_blocks(frame: str) -> list[str]:
    blocks: list[str] = []
    for match in re.finditer(r"\\\[[\s\S]*?\\\]", frame or ""):
        blocks.append(match.group(0).strip())
    for match in re.finditer(r"\\begin\{(?:equation|align|gather|multline)\*?\}[\s\S]*?\\end\{(?:equation|align|gather|multline)\*?\}", frame or ""):
        blocks.append(match.group(0).strip())
    return blocks


def _extract_item_texts(frame: str, limit: int = 3) -> list[str]:
    texts: list[str] = []
    for match in re.finditer(r"\\item(?:<[^>]*>)?(?:\[[\s\S]*?\])?\s*([\s\S]*?)(?=\\item|\\end\{itemize\}|\\end\{enumerate\})", frame or ""):
        item = re.sub(r"\s+", " ", match.group(1)).strip()
        if item:
            texts.append(item)
        if len(texts) >= limit:
            break
    return texts


def _short_formula_note(frame: str) -> str:
    items = _extract_item_texts(frame, 2)
    if not items:
        return ""
    note = " ".join(items)
    note = re.sub(r"\\begin\{[^{}]*\}|\\end\{[^{}]*\}", " ", note)
    note = re.sub(r"\s+", " ", note).strip()
    return note[:180]


def _formula_annotation_targets(formula_text: str) -> list[tuple[str, str, str]]:
    source = formula_text or ""
    labels: list[tuple[str, str, str]] = []
    if re.search(r"\\pi\s*_\s*\{?j\}?", source):
        labels.append(("pij", r"\pi", "突变 j 被固定的概率"))
    if re.search(r"\bp_b\b|p_\{b\}", source):
        labels.append(("pb", "p_b", "新突变有益的概率"))
    if re.search(r"\\Phi\b", source):
        labels.append(("phi", r"\Phi", "标准正态分布的累积函数"))
    if re.search(r"\\sum\b", source):
        labels.append(("sum", r"\sum", "把所有候选突变的效应加总"))
    if re.search(r"s\s*_\s*\{?j\}?", source):
        labels.append(("sj", "s_j", "突变 j 的选择优势"))
    if re.search(r"s\s*_\s*\{?\\ell\}?", source):
        labels.append(("sell", r"s_{\ell}", "第 $\\ell$ 个突变的选择优势"))
    if re.search(r"\\sqrt\{?n\}?|\bn\b", source):
        labels.append(("n", "n", "参与适应的性状数量"))
    if re.search(r"\br\b", source):
        labels.append(("r", "r", "新突变带来的表型改变大小"))
    if re.search(r"\bd\b", source):
        labels.append(("d", "d", r"当前表型到最优值 $\theta$ 的距离"))
    if re.search(r"\bk\b", source):
        labels.append(("k", "k", "可供选择的有益突变数"))
    if re.search(r"\bi\b", source):
        labels.append(("i", "i", "当前排序或行走步数"))
    if re.search(r"平均步数|average", source, re.IGNORECASE):
        labels.append(("avgsteps", r"\text{平均步数}", "预期需要走多少步"))
    if re.search(r"\\nu\b", source):
        labels.append(("nu", r"\nu", "最优值移动的速度"))
    if re.search(r"\bN\b|N\\mu", source):
        labels.append(("N", "N", "群体大小"))
    if re.search(r"\\mu\b", source):
        labels.append(("mu", r"\mu", "突变率"))
    return labels[:5]


def _formula_anchor_name(mark: str, page_index: int) -> str:
    suffix = re.sub(r"[^A-Za-z0-9]+", "", str(mark or "")) or "mark"
    return f"kgp{max(1, int(page_index or 1))}_{suffix}"


def _mark_formula_annotation_targets(formula: str, targets: list[tuple[str, str, str]], page_index: int = 1) -> str:
    result = formula or ""
    for mark, token, _label in targets:
        safe_mark = _formula_anchor_name(mark, page_index)
        if f"\\tikzmark{{{safe_mark}}}" in result:
            continue
        if token == "p_b":
            result = re.sub(r"p_b", r"\\tikzmark{" + safe_mark + r"}p_b", result, count=1)
        elif token == r"\Phi":
            result = re.sub(r"\\Phi", r"\\tikzmark{" + safe_mark + r"}\\Phi", result, count=1)
        elif token == r"\sum":
            result = re.sub(r"\\sum", r"\\tikzmark{" + safe_mark + r"}\\sum", result, count=1)
        elif token == r"\nu":
            result = re.sub(r"\\nu", r"\\tikzmark{" + safe_mark + r"}\\nu", result, count=1)
        elif token == r"\mu":
            result = re.sub(r"\\mu", r"\\tikzmark{" + safe_mark + r"}\\mu", result, count=1)
        elif token == r"\pi":
            result = re.sub(r"\\pi", r"\\tikzmark{" + safe_mark + r"}\\pi", result, count=1)
        elif token == r"\text{平均步数}":
            result = re.sub(r"\\text\{平均步数\}", r"\\tikzmark{" + safe_mark + r"}\\text{平均步数}", result, count=1)
        elif token == "s_j":
            result = re.sub(r"s\s*_\s*(?:\{j\}|j)", r"\\tikzmark{" + safe_mark + r"}s_{j}", result, count=1)
        elif token == r"s_{\ell}":
            result = re.sub(r"s\s*_\s*(?:\{\\ell\}|\\ell)", r"\\tikzmark{" + safe_mark + r"}s_{\\ell}", result, count=1)
        elif token in {"r", "n", "d", "N", "k", "i"}:
            result = re.sub(r"(?<![A-Za-z])" + re.escape(token) + r"(?![A-Za-z])", r"\\tikzmark{" + safe_mark + r"}" + token, result, count=1)
    return result


def _formula_callout_latex(formula_text: str, page_index: int = 1) -> str:
    labels = _formula_annotation_targets(formula_text)
    if not labels:
        return ""
    nodes: list[str] = ["  \\begin{tikzpicture}[remember picture, overlay]"]
    offsets = [
        ("-2.5cm", "-0.2cm", "2.9cm"),
        ("0.5cm", "1.0cm", "2.4cm"),
        ("0.5cm", "-0.8cm", "2.7cm"),
        ("2.5cm", "-0.5cm", "3.3cm"),
        ("-1.8cm", "1.1cm", "3.6cm"),
    ]
    for idx, (mark, _token, label) in enumerate(labels[:5]):
        x, y, width = offsets[idx]
        safe_mark = _formula_anchor_name(mark, page_index)
        nodes.append(
            f"    \\onslide<{idx + 4}->{{\n"
            "    \\node[rectangle callout, callout absolute pointer={(pic cs:" + safe_mark + ")}, "
            "draw=blue, fill=white, rounded corners, "
            f"text width={width}, align=center, font=\\footnotesize] "
            f"at ([shift={{({x},{y})}}] pic cs:{safe_mark}) "
            "{" + label + "};\n"
            "    }"
        )
    nodes.append("  \\end{tikzpicture}")
    return "\n".join(nodes)


def _enforce_formula_callout_anchor_positioning(latex: str) -> str:
    offsets = [
        ("2.0cm", "1.15cm"),
        ("3.2cm", "0.0cm"),
        ("2.0cm", "-1.15cm"),
        ("-2.2cm", "0.95cm"),
        ("-2.2cm", "-0.95cm"),
    ]
    index = 0

    pattern = re.compile(
        r"(\\node\[(?=[^\]]*rectangle callout)(?=[^\]]*callout absolute pointer=\{\(pic cs:(?P<mark>[^)}]+)\)\})[^\]]*\]\s*)"
        r"at\s*\(\s*\[xshift=[^\]]+\]\s*current page\.north west\s*\)",
        re.MULTILINE,
    )

    def replace(match: re.Match[str]) -> str:
        nonlocal index
        mark = match.group("mark").strip()
        x, y = offsets[index % len(offsets)]
        index += 1
        return f"{match.group(1)}at ([shift={{({x},{y})}}] pic cs:{mark})"

    return pattern.sub(replace, latex or "")


def _enforce_formula_callout_sequence_and_spacing(latex: str) -> str:
    text = latex or ""
    frames = _extract_latex_frames(text)
    if not frames:
        return text
    offsets = [
        ("2.0cm", "1.15cm"),
        ("3.2cm", "0.0cm"),
        ("2.0cm", "-1.15cm"),
        ("-2.2cm", "0.95cm"),
        ("-2.2cm", "-0.95cm"),
    ]
    result = text
    callout_pattern = re.compile(
        r"\\onslide<\d+->\{\s*"
        r"\\node\[(?P<options>[^\]]*rectangle callout[^\]]*callout absolute pointer=\{\(pic cs:(?P<ptr>[^)}]+)\)\}[^\]]*)\]\s*"
        r"at\s*\(\[shift=\{\([^)]*\)\}\]\s*pic cs:(?P=ptr)\)",
        re.DOTALL,
    )

    def replace_frame(frame: str) -> str:
        index = 0

        def replace_callout(match: re.Match[str]) -> str:
            nonlocal index
            mark = match.group("ptr").strip()
            x, y = offsets[index % len(offsets)]
            index += 1
            return (
                f"\\onslide<{index}->{{\n"
                f"      \\node[{match.group('options')}] at ([shift={{({x},{y})}}] pic cs:{mark})"
            )

        return callout_pattern.sub(replace_callout, frame)

    for frame in frames:
        if "rectangle callout" not in frame or "pic cs:" not in frame:
            continue
        replacement = replace_frame(frame)
        if replacement != frame:
            result = result.replace(frame, replacement, 1)
    return result


def _anchor_already_page_scoped(mark: str) -> bool:
    return bool(re.match(r"^(?:kg)?p\d+_", str(mark or "")))


def _scope_frame_tikzmarks(frame: str, page_index: int) -> str:
    labels: set[str] = set()
    for match in re.finditer(r"\\tikzmark\{([^{}]+)\}", frame or ""):
        label = match.group(1).strip()
        if label and not _anchor_already_page_scoped(label):
            labels.add(label)
    for match in re.finditer(r"pic cs:([A-Za-z0-9_.:-]+)", frame or ""):
        label = match.group(1).strip()
        if label and not _anchor_already_page_scoped(label):
            labels.add(label)
    if not labels:
        return frame

    scoped = frame
    for label in sorted(labels, key=len, reverse=True):
        safe_label = re.sub(r"[^A-Za-z0-9_.:-]+", "", label)
        if not safe_label:
            continue
        replacement = _formula_anchor_name(safe_label, page_index)
        scoped = scoped.replace(f"\\tikzmark{{{label}}}", f"\\tikzmark{{{replacement}}}")
        scoped = scoped.replace(f"pic cs:{label}", f"pic cs:{replacement}")
    return scoped


def _enforce_unique_tikzmark_names(latex: str) -> str:
    text = latex or ""
    frames = _extract_latex_frames(text)
    if not frames:
        return text
    result = text
    for frame_index, frame in enumerate(frames, 1):
        replacement = _scope_frame_tikzmarks(frame, frame_index)
        if replacement != frame:
            result = result.replace(frame, replacement, 1)
    return result


def _build_formula_breakdown_frame(title: str, formulas: list[str], note: str = "", page_index: int = 1) -> str:
    out = [f"\\begin{{frame}}{{{title}}}", "  \\vspace*{0.35cm}", "  \\begin{center}"]
    all_targets = _formula_annotation_targets("\n".join(formulas[:2]))
    for idx, formula in enumerate(formulas[:2]):
        body = formula
        if body.startswith(r"\[") and body.endswith(r"\]"):
            body = body[2:-2].strip()
        body = _mark_formula_annotation_targets(body, all_targets, page_index)
        out.append("    {\\Large")
        out.append("    \\[")
        out.append("      " + body)
        out.append("    \\]")
        out.append("    }")
        if idx == 0 and len(formulas[:2]) > 1:
            out.append("    \\vspace{0.35cm}")
    out.append("  \\end{center}")
    if note:
        out.append("  \\vspace{0.15cm}")
        out.append("  {\\scriptsize " + note + "}")
    callouts = _formula_callout_latex("\n".join(formulas[:2]), page_index)
    if callouts:
        out.append(callouts)
    out.append("\\end{frame}")
    return "\n".join(out)


def _enforce_formula_breakdown_layout(latex: str) -> str:
    text = latex or ""
    frames = _extract_latex_frames(text)
    if not frames:
        return text
    result = text
    for frame_index, frame in enumerate(frames, 1):
        if r"\safecontentimage" in frame or r"\safeverticalimage" in frame or r"\includegraphics" in frame:
            continue
        formulas = _extract_display_formula_blocks(frame)
        if not formulas:
            continue
        has_list = r"\begin{itemize}" in frame or r"\begin{enumerate}" in frame
        has_callout = (
            "rectangle callout" in frame
            and "callout absolute pointer={(pic cs:" in frame
            and r"\tikzmark{" in frame
        )
        if not has_list and has_callout:
            continue
        title = _frame_title(frame) or "Formula Breakdown"
        note = _short_formula_note(frame)
        replacement = _build_formula_breakdown_frame(title, formulas, note, frame_index)
        result = result.replace(frame, replacement, 1)
    return result


def _frame_inner_body(frame: str) -> str:
    text = frame or ""
    begin_match = re.search(r"\\begin\{frame\}", text)
    if begin_match:
        span = _frame_begin_title_span(text)
        if span:
            text = text[span[1]:]
        else:
            text = text[begin_match.end():]
            text = re.sub(r"^\s*\[[^\]]*\]", "", text, count=1)
    text = re.sub(r"\\frametitle\{[^{}]*\}", "", text)
    text = re.sub(r"\\end\{frame\}\s*$", "", text).strip()
    return text


def _clean_latex_heading_text(value: str) -> str:
    text = value or ""
    text = text.replace("\\\\", " ")
    text = re.sub(r"\\(?:textbf|textit|emph|textcolor)\*?(?:\[[^\]]*\])?\{([^{}]*)\}", r"\1", text)
    text = re.sub(r"\\[a-zA-Z]+\*?(?:\[[^\]]*\])?(?:\{([^{}]*)\})?", r" \1 ", text)
    text = re.sub(r"[{}$]", " ", text)
    return re.sub(r"\s+", " ", text).strip()


def _is_generic_outline_title(title: str) -> bool:
    clean = _clean_latex_heading_text(title).lower()
    clean = re.sub(r"[^a-z0-9 ]+", " ", clean)
    clean = re.sub(r"\s+", " ", clean).strip()
    return clean in {
        "outline",
        "chapter outline",
        "outline of this chapter",
        "contents",
        "table of contents",
    }


def _is_empty_outline_frame(frame: str) -> bool:
    title = _frame_title(frame)
    if not _is_generic_outline_title(title):
        return False
    body = _frame_inner_body(frame)
    body = re.sub(r"%.*", "", body).strip()
    return not body


def _outline_items_from_latex(latex: str) -> list[str]:
    sections = []
    for match in re.finditer(r"\\section\*?\{([^{}]+)\}", latex or ""):
        title = re.sub(r"\s+", " ", match.group(1)).strip()
        if title and title not in sections:
            sections.append(title)
    if sections:
        return sections[:8]
    titles = []
    for frame in _extract_latex_frames(latex):
        title = _frame_title(frame)
        if not title:
            continue
        low = _clean_latex_heading_text(title).lower()
        if _is_generic_outline_title(title) or low == "review of last course":
            continue
        if title not in titles:
            titles.append(title)
        if len(titles) >= 6:
            break
    return titles


def _document_chapter_title(latex: str) -> str:
    text = latex or ""
    match = re.search(r"\\title(?:\[[^\]]*\])?", text)
    if match:
        pos = match.end()
        while pos < len(text) and text[pos].isspace():
            pos += 1
        span = _balanced_group_span(text, pos) if pos < len(text) and text[pos] == "{" else None
        if span:
            title = _clean_latex_heading_text(text[span[0] + 1:span[1] - 1])
            if title:
                return title
    return "Evolutionary Theory on Polygenic Trait"


def _replace_frame_begin_title(frame: str, title: str) -> str:
    span = _frame_begin_title_span(frame)
    if span:
        return frame[:span[0]] + "{" + title + "}" + frame[span[1]:]
    begin_match = re.search(r"\\begin\{frame\}(?:\[[^\]]*\])?", frame or "")
    if not begin_match:
        return frame
    return frame[:begin_match.end()] + "{" + title + "}" + frame[begin_match.end():]


def _enforce_outline_frame_titles(latex: str, preferred_title: str = "") -> str:
    text = latex or ""
    chapter_title = preferred_title.strip() or _document_chapter_title(text)
    if not chapter_title:
        return text
    result = text
    changed_first_outline = False
    for frame in _extract_latex_frames(text):
        title = _frame_title(frame)
        if not title:
            continue
        body = _frame_inner_body(frame)
        if r"\begin{itemize}" not in body and r"\tableofcontents" not in body:
            continue
        clean_title = _clean_latex_heading_text(title).lower()
        is_template_outline_title = clean_title in {
            "evolutionary theory on polygenic trait",
            "long term response",
        }
        should_replace = (
            (bool(preferred_title.strip()) and not changed_first_outline)
            or _is_generic_outline_title(title)
            or is_template_outline_title
        )
        if not should_replace:
            continue
        replacement = _replace_frame_begin_title(frame, chapter_title)
        result = result.replace(frame, replacement, 1)
        changed_first_outline = True
    return result


def _build_filled_outline_frame(title: str, items: list[str], has_sections: bool) -> str:
    if not items:
        items = ["Key Concepts", "Main Results", "Summary"]
    out = [
        f"\\begin{{frame}}{{{title}}}",
        "  \\vfill",
        "  \\begin{center}",
        "    \\begin{minipage}{0.7\\textwidth}",
        "      \\begin{itemize}",
        "        \\setlength{\\itemsep}{0.3\\baselineskip}",
    ]
    for item in items:
        out.append(f"        \\item \\textcolor{{black}}{{{item}}}")
    out.extend([
        "      \\end{itemize}",
        "    \\end{minipage}",
        "  \\end{center}",
        "  \\vfill",
        "\\end{frame}",
    ])
    return "\n".join(out)


def _fill_empty_outline_frames(latex: str, preferred_title: str = "") -> str:
    text = latex or ""
    frames = _extract_latex_frames(text)
    if not frames:
        return _enforce_outline_frame_titles(text, preferred_title)
    items = _outline_items_from_latex(text)
    has_sections = bool(re.search(r"\\section\*?\{[^{}]+\}", text))
    chapter_title = preferred_title.strip() or _document_chapter_title(text)
    result = text
    for frame in frames:
        if not _is_empty_outline_frame(frame):
            continue
        replacement = _build_filled_outline_frame(chapter_title, items, has_sections)
        result = result.replace(frame, replacement, 1)
    return _enforce_outline_frame_titles(result, preferred_title)


def _is_empty_frame(frame: str) -> bool:
    body = _frame_inner_body(frame)
    body = re.sub(r"%.*", "", body).strip()
    return not body


def _content_snippet_for_empty_title(title: str, source_text: str) -> list[str]:
    title_words = [w.lower() for w in re.findall(r"[A-Za-z][A-Za-z0-9-]+", title or "") if len(w) > 3]
    lines = [re.sub(r"\s+", " ", line).strip() for line in (source_text or "").splitlines()]
    candidates: list[str] = []
    for line in lines:
        if not line or line.startswith("\\"):
            continue
        low = line.lower()
        if title_words and not any(word in low for word in title_words[:4]):
            continue
        plain = re.sub(r"^[#>*\-\d.\s]+", "", line).strip()
        if plain and plain not in candidates:
            candidates.append(plain[:180])
        if len(candidates) >= 4:
            break
    return candidates


def _build_nonempty_placeholder_frame(title: str, items: list[str]) -> str:
    if not items:
        items = [
            "该主题需要补充定义、核心假设和适用条件。",
            "请结合导入的补充章节或公式章节完善本页内容。",
        ]
    out = [f"\\begin{{frame}}{{{title or 'Topic'}}}", "  \\begin{itemize}"]
    for item in items[:4]:
        out.append(f"    \\item {_escape_latex_text_preserving_inline_math(item)}")
    out.extend(["  \\end{itemize}", "\\end{frame}"])
    return "\n".join(out)


def _fill_empty_content_frames(latex: str, source_text: str = "") -> str:
    text = latex or ""
    result = text
    for frame in _extract_latex_frames(text):
        if not _is_empty_frame(frame) or _is_empty_outline_frame(frame):
            continue
        title = _frame_title(frame) or "Topic"
        items = _content_snippet_for_empty_title(title, source_text)
        result = result.replace(frame, _build_nonempty_placeholder_frame(title, items), 1)
    return result


def _remove_numbered_outline_item_labels(latex: str) -> str:
    text = latex or ""
    text = re.sub(
        r"\\item\s*\[\s*\\textcolor\{(black|gray)\}\{\\textbf\{\d+\.\}\}\s*\]\s*\\textcolor\{\1\}\{",
        r"\\item \\textcolor{\1}{",
        text,
    )
    text = re.sub(r"\\item\s*\[\s*(?:\\textbf\{)?\d+\.(?:\})?\s*\]\s*", r"\\item ", text)
    text = re.sub(r"(?<=\})\s*\[\d+\.\]", "", text)
    text = re.sub(r"(?<=[A-Za-z\u4e00-\u9fff])\s*\[\d+\.\](?=\s*(?:\\\\|$|[}\n]))", "", text)
    return text


def _figure_caption_map_from_markdown(content: str) -> dict[str, str]:
    captions: dict[str, str] = {}
    pattern = re.compile(r"^\s*>?\s*Figure\s+(\d+(?:\.\d+)*)\s+(.+)$", re.IGNORECASE)
    for line in (content or "").splitlines():
        match = pattern.match(line.strip())
        if not match:
            continue
        number = match.group(1)
        caption = re.sub(r"\s+", " ", match.group(2)).strip()
        caption = re.sub(r"^(?:[·路\-:：]\s*)+", "", caption).strip()
        if caption and number not in captions:
            captions[number] = caption
    return captions


def _matched_image_caption(frame_title: str, image_path: str, content: str) -> str:
    caption_map = _figure_caption_map_from_markdown(content)
    candidates = [frame_title or "", image_path or "", Path(str(image_path or "")).stem]
    for value in candidates:
        match = re.search(r"(?:Figure|Fig\.?|图)?\s*(\d+(?:[._]\d+)*)", value, re.IGNORECASE)
        if not match:
            continue
        number = match.group(1).replace("_", ".")
        caption = caption_map.get(number)
        if caption:
            return _escape_latex_text_preserving_inline_math(f"Figure {number} {caption}")
    return ""


def _enforce_matched_image_captions(
    latex: str,
    content: str,
    figure_assets: dict[str, str] | None = None,
) -> str:
    if not content:
        return latex or ""
    frame_pattern = re.compile(r"\\begin\{frame\}(?:\[[^\]]*\])?(?:\{[^{}]*\})?[\s\S]*?\\end\{frame\}")
    image_pattern = re.compile(
        r"\\(?P<cmd>includegraphics|safecontentimage|safeverticalimage)(?:\[[^\]]*\])?\{(?P<path>[^{}]+)\}"
    )

    def replace_frame(match: re.Match[str]) -> str:
        frame = match.group(0)
        image_match = image_pattern.search(frame)
        if not image_match:
            return frame
        image_path = image_match.group("path").strip()
        frame_title = _figure_frame_title_for_path(image_path, figure_assets or {})
        body = re.sub(r"^\\begin\{frame\}(?:\[[^\]]*\])?(?:\{[^{}]*\})?", "", frame)
        body = re.sub(r"\\end\{frame\}\s*$", "", body)
        caption = _caption_text_from_image_frame_body(body, image_match.group(0))
        if not _is_uninformative_image_caption(caption, frame_title, image_path):
            return frame
        matched_caption = _matched_image_caption(frame_title, image_path, content)
        if not matched_caption:
            return frame
        orientation = "vertical" if image_match.group("cmd") == "safeverticalimage" else _image_orientation_for_path(image_path, figure_assets or {})
        if orientation == "vertical":
            return (
                f"\\begin{{frame}}{{{frame_title}}}\n"
                "  \\begin{columns}[T]\n"
                "    \\begin{column}{0.45\\textwidth}\n"
                f"      \\scriptsize {matched_caption}\n"
                "    \\end{column}\n"
                "    \\begin{column}{0.45\\textwidth}\n"
                "      \\centering\n"
                f"      \\safeverticalimage{{{image_path}}}\n"
                "    \\end{column}\n"
                "  \\end{columns}\n"
                "\\end{frame}"
            )
        return (
            f"\\begin{{frame}}{{{frame_title}}}\n"
            "  \\centering\n"
            f"  \\safecontentimage{{{image_path}}}\n"
            "  \\vspace{0.3cm}\n"
            "  \\begin{center}\n"
            f"    \\parbox{{0.95\\textwidth}}{{\\scriptsize {matched_caption}}}\n"
            "  \\end{center}\n"
            "\\end{frame}"
        )

    return frame_pattern.sub(replace_frame, latex or "")


def _escape_latex_text(value: str) -> str:
    text = str(value or "")
    replacements = {
        "\\": r"\textbackslash{}",
        "&": r"\&",
        "%": r"\%",
        "$": r"\$",
        "#": r"\#",
        "_": r"\_",
        "{": r"\{",
        "}": r"\}",
    }
    for src, dst in replacements.items():
        text = text.replace(src, dst)
    return text


def _escape_latex_text_preserving_inline_math(value: str) -> str:
    text = str(value or "")
    if not text:
        return ""
    tokens: list[str] = []

    def store(match: re.Match[str]) -> str:
        tokens.append(match.group(0))
        return f"@@PPTMATH{len(tokens) - 1}@@"

    text = re.sub(r"\\\([\s\S]*?\\\)", store, text)
    text = re.sub(r"\$[^$\n]+\$", store, text)
    text = _escape_latex_text(text)
    for idx, token in enumerate(tokens):
        text = text.replace(f"@@PPTMATH{idx}@@", token)
    return text


def _ppt_formula_like(value: str) -> bool:
    text = str(value or "").strip()
    if not text:
        return False
    if re.search(r"^(\$\$[\s\S]*\$\$|\\\[[\s\S]*\\\]|\\\([\s\S]*\\\))$", text):
        return True
    if re.search(r"\\(?:frac|sqrt|sum|prod|int|bar|overline|hat|vec|theta|alpha|beta|gamma|delta|lambda|mu|sigma|phi|omega|partial|nabla|mathrm|mathbb|mathcal)\b", text):
        return True
    return bool(re.search(r"[_^]\s*\{?|\b[xyzndr]\s*[=<>]", text) and not re.search(r"[\u4e00-\u9fff]{4,}", text))


def _ppt_formula_source(value: str) -> str:
    text = str(value or "").strip()
    text = re.sub(r"^\$\$\s*|\s*\$\$$", "", text)
    text = re.sub(r"^\\\[\s*|\s*\\\]$", "", text)
    text = re.sub(r"^\\\(\s*|\s*\\\)$", "", text)
    text = re.sub(r"^\$\s*|\s*\$$", "", text)
    return text.strip()


def _ppt_figure_numbers_from_text(value: str) -> list[str]:
    text = str(value or "")
    numbers: list[str] = []
    for match in re.finditer(r"(?:Figure|Fig\.?|图)\s*([0-9]+(?:[._][0-9]+)*)", text, re.IGNORECASE):
        number = match.group(1).replace("_", ".")
        if number and number not in numbers:
            numbers.append(number)
    return numbers


def _ppt_normalize_symbol_text(value: str) -> str:
    text = str(value or "")
    replacements = {
        "θ": r"$\theta$",
        "μ": r"$\mu$",
        "Φ": r"$\Phi$",
        "σ": r"$\sigma$",
        "¦È": r"$\theta$",
        "¦Ì": r"$\mu$",
        "¦Ð": r"$\Phi$",
        "¦Ó": r"$\sigma$",
    }
    for source, target in replacements.items():
        text = text.replace(source, target)
    return text


def _ppt_text_lines(shape: object) -> list[str]:
    if not getattr(shape, "has_text_frame", False):
        return []
    lines: list[str] = []
    for paragraph in getattr(shape.text_frame, "paragraphs", []) or []:
        parts = [str(getattr(run, "text", "") or "") for run in getattr(paragraph, "runs", []) or []]
        text = "".join(parts).strip() if parts else str(getattr(paragraph, "text", "") or "").strip()
        if text:
            lines.append(_ppt_normalize_symbol_text(text))
    if not lines:
        text = str(getattr(shape, "text", "") or "").strip()
        if text:
            lines.extend(_ppt_normalize_symbol_text(line.strip()) for line in text.splitlines() if line.strip())
    return lines


def _ppt_is_title_shape(shape: object) -> bool:
    if not getattr(shape, "has_text_frame", False):
        return False
    try:
        ph_type = shape.placeholder_format.type
        if ph_type in {1, 3, 4}:
            return True
    except Exception:
        pass
    try:
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                if run.font.size and int(run.font.size) >= 24 * 12700:
                    return True
                if run.font.bold:
                    return True
    except Exception:
        pass
    return False


def _ppt_table_latex(shape: object) -> str:
    if not getattr(shape, "has_table", False):
        return ""
    rows = []
    try:
        for row in shape.table.rows:
            rows.append([_escape_latex_text_preserving_inline_math(str(cell.text or "").strip()) for cell in row.cells])
    except Exception:
        return ""
    if not rows or not rows[0]:
        return ""
    col_count = max(len(row) for row in rows)
    spec = "l" * col_count
    out = [
        "  \\begin{table}",
        "    \\centering",
        f"    \\begin{{tabular}}{{{spec}}}",
        "      \\toprule",
    ]
    for idx, row in enumerate(rows):
        padded = row + [""] * (col_count - len(row))
        out.append("      " + " & ".join(padded) + r" \\")
        if idx == 0:
            out.append("      \\midrule")
    out.extend(["      \\bottomrule", "    \\end{tabular}", "  \\end{table}"])
    return "\n".join(out)


def _ppt_shape_bounds(shape: object) -> dict[str, int]:
    return {
        "left": int(getattr(shape, "left", 0) or 0),
        "top": int(getattr(shape, "top", 0) or 0),
        "width": int(getattr(shape, "width", 0) or 0),
        "height": int(getattr(shape, "height", 0) or 0),
    }


def _ppt_image_extension(shape: object) -> str:
    try:
        ext = str(shape.image.ext or "").strip().lower()
        if ext:
            return "." + ext.lstrip(".")
    except Exception:
        pass
    content_type = ""
    try:
        content_type = str(shape.image.content_type or "").lower()
    except Exception:
        pass
    if "jpeg" in content_type or "jpg" in content_type:
        return ".jpg"
    if "webp" in content_type:
        return ".webp"
    if "gif" in content_type:
        return ".gif"
    return ".png"


def _ppt_shape_image_bytes(shape: object) -> bytes:
    try:
        return bytes(shape.image.blob)
    except Exception:
        return b""


def _ppt_image_latex_path(slide_index: int, image_index: int, ext: str, preferred_number: str = "") -> tuple[str, str]:
    number = str(preferred_number or "").strip().replace("_", ".")
    if not number:
        number = f"{slide_index}.{image_index}"
    return f"Figure {number}", f"fig/{number}{ext}"


def _ppt_box_cm(bounds: dict[str, int], slide_width: int, slide_height: int) -> tuple[float, float, float, float]:
    width = max(1, int(slide_width or 12192000))
    height = max(1, int(slide_height or 6858000))
    return (
        bounds["left"] / width * 16.0,
        bounds["top"] / height * 9.0,
        max(0.2, bounds["width"] / width * 16.0),
        max(0.2, bounds["height"] / height * 9.0),
    )


def _ppt_tikz_node_at(bounds: dict[str, int], slide_width: int, slide_height: int, options: str, content: str) -> str:
    x, y, _w, _h = _ppt_box_cm(bounds, slide_width, slide_height)
    return (
        f"    \\node[{options}] at ([xshift={x:.3f}cm,yshift=-{y:.3f}cm] current page.north west) "
        f"{{{content}}};"
    )


def _ppt_box_center_cm(bounds: dict[str, int], slide_width: int, slide_height: int) -> tuple[float, float]:
    x, y, w, h = _ppt_box_cm(bounds, slide_width, slide_height)
    return x + w / 2.0, y + h / 2.0


def _ppt_nearest_pointer_coordinate(
    note_bounds: dict[str, int],
    target_bounds: list[dict[str, int]],
    slide_width: int,
    slide_height: int,
) -> str:
    nx, ny = _ppt_box_center_cm(note_bounds, slide_width, slide_height)
    if not target_bounds:
        return f"([xshift={nx:.3f}cm,yshift=-{ny:.3f}cm] current page.north west)"
    target = min(
        target_bounds,
        key=lambda bounds: (
            (_ppt_box_center_cm(bounds, slide_width, slide_height)[0] - nx) ** 2
            + (_ppt_box_center_cm(bounds, slide_width, slide_height)[1] - ny) ** 2
        ),
    )
    tx, ty = _ppt_box_center_cm(target, slide_width, slide_height)
    return f"([xshift={tx:.3f}cm,yshift=-{ty:.3f}cm] current page.north west)"


def _ppt_lines_to_latex(lines: list[str]) -> str:
    return r" \\ ".join(_escape_latex_text_preserving_inline_math(line) for line in lines if str(line or "").strip())


def _ppt_shape_is_blue_note(shape: object) -> bool:
    name = str(getattr(shape, "name", "") or "").lower()
    if "标注" in name or "callout" in name:
        return True
    try:
        rgb = shape.line.color.rgb
        if rgb and int(rgb[2]) > int(rgb[0]) and int(rgb[2]) > int(rgb[1]):
            return True
    except Exception:
        pass
    return False


def _ppt_picture_is_formula(bounds: dict[str, int], slide_width: int, slide_height: int, has_figure_number: bool) -> bool:
    if has_figure_number:
        return False
    width = max(1, bounds.get("width", 0))
    height = max(1, bounds.get("height", 0))
    if width / height >= 2.6:
        return True
    return height <= max(1, int(slide_height or 6858000)) * 0.24 and width >= max(1, int(slide_width or 12192000)) * 0.20


def _safe_upload_filename(filename: str, fallback: str) -> str:
    name = Path(str(filename or "").replace("\\", "/")).name.strip() or fallback
    name = re.sub(r'[<>:"/\\|?*\x00-\x1f]+', "_", name).strip(" .")
    return name or fallback


def _ppt_slide_to_latex_frame(
    slide: object,
    slide_index: int,
    package_dir: Path,
    asset_urls: dict[str, str],
    slide_width: int = 12192000,
    slide_height: int = 6858000,
) -> tuple[str, list[dict[str, str]]]:
    title = ""
    all_text: list[str] = []
    for shape in slide.shapes:
        lines = _ppt_text_lines(shape)
        if lines:
            all_text.extend(lines)
            bounds = _ppt_shape_bounds(shape)
            is_title_candidate = (
                _ppt_is_title_shape(shape)
                and not _ppt_shape_is_blue_note(shape)
                and bounds["top"] <= int(slide_height * 0.22)
            )
            if is_title_candidate and not title:
                title = " ".join(lines)

    preferred_numbers = _ppt_figure_numbers_from_text("\n".join(all_text))
    pointer_targets = [
        _ppt_shape_bounds(shape)
        for shape in slide.shapes
        if MSO_SHAPE_TYPE is not None and getattr(shape, "shape_type", None) == MSO_SHAPE_TYPE.PICTURE
    ]
    image_records: list[dict[str, str]] = []
    image_index = 0
    formula_index = 0
    nodes: list[str] = []

    for shape in slide.shapes:
        bounds = _ppt_shape_bounds(shape)
        x, y, w, h = _ppt_box_cm(bounds, slide_width, slide_height)
        shape_type = getattr(shape, "shape_type", None)

        if MSO_SHAPE_TYPE is not None and shape_type == MSO_SHAPE_TYPE.PICTURE:
            blob = _ppt_shape_image_bytes(shape)
            if not blob:
                continue
            ext = _ppt_image_extension(shape)
            has_figure_number = image_index < len(preferred_numbers)
            if _ppt_picture_is_formula(bounds, slide_width, slide_height, has_figure_number):
                formula_index += 1
                output_name = f"formula-{slide_index}.{formula_index}{ext}"
                latex_path = f"fig/{output_name}"
                label = f"Formula {slide_index}.{formula_index}"
                content = (
                    f"\\parbox[c][{h:.3f}cm][c]{{{w:.3f}cm}}"
                    f"{{\\centering\\includegraphics[width=\\textwidth,keepaspectratio]{{{latex_path}}}}}"
                )
                node_options = (
                    f"anchor=north west, inner sep=0pt, "
                    f"minimum width={w:.3f}cm, minimum height={h:.3f}cm"
                )
            else:
                image_index += 1
                preferred_number = preferred_numbers[image_index - 1] if image_index - 1 < len(preferred_numbers) else ""
                label, latex_path = _ppt_image_latex_path(slide_index, image_index, ext, preferred_number)
                content = (
                    f"\\parbox[c][{h:.3f}cm][c]{{{w:.3f}cm}}"
                    f"{{\\centering\\includegraphics[width=\\textwidth,keepaspectratio]{{{latex_path}}}}}"
                )
                node_options = (
                    f"anchor=north west, draw=black, inner sep=0pt, "
                    f"minimum width={w:.3f}cm, minimum height={h:.3f}cm, "
                    f"text width={w:.3f}cm, align=center, font=\\scriptsize"
                )
            output_name = Path(latex_path).name
            output_path = package_dir / output_name
            output_path.write_bytes(blob)
            url = f"/beamer-generator/uploads/{package_dir.name}/{output_name}"
            asset_urls[latex_path] = url
            asset_urls[output_name] = url
            image_records.append({"label": label, "path": latex_path, "url": url})
            nodes.append(_ppt_tikz_node_at(bounds, slide_width, slide_height, node_options, content))
            continue

        if getattr(shape, "has_table", False):
            table_latex = _ppt_table_latex(shape)
            if table_latex:
                content = f"\\parbox[t][{h:.3f}cm][t]{{{w:.3f}cm}}{{\\scriptsize {table_latex}}}"
                nodes.append(_ppt_tikz_node_at(bounds, slide_width, slide_height, "anchor=north west, inner sep=0pt", content))
            continue

        lines = _ppt_text_lines(shape)
        if not lines:
            continue
        formulas = [_ppt_formula_source(line) for line in lines if _ppt_formula_like(line)]
        if formulas and len(formulas) == len(lines):
            content = "\\parbox[t][" + f"{h:.3f}cm" + "][t]{" + f"{w:.3f}cm" + "}{" + " ".join(f"\\[{formula}\\]" for formula in formulas) + "}"
            nodes.append(_ppt_tikz_node_at(bounds, slide_width, slide_height, "anchor=north west, inner sep=0pt", content))
            continue
        text_content = _ppt_lines_to_latex(lines)
        if not text_content:
            continue
        if _ppt_shape_is_blue_note(shape):
            pointer = _ppt_nearest_pointer_coordinate(bounds, pointer_targets, slide_width, slide_height)
            content = f"\\parbox[c][{h:.3f}cm][c]{{{w:.3f}cm}}{{\\centering {text_content}}}"
            options = (
                f"rectangle callout, callout absolute pointer={{{pointer}}}, "
                f"anchor=north west, draw=blue, fill=white, rounded corners, "
                f"inner sep=0pt, text width={w:.3f}cm, minimum width={w:.3f}cm, minimum height={h:.3f}cm, "
                "align=center, font=\\footnotesize"
            )
        else:
            content = f"\\parbox[t][{h:.3f}cm][t]{{{w:.3f}cm}}{{{text_content}}}"
            options = f"anchor=north west, inner sep=0pt, text width={w:.3f}cm, align=left, font=\\scriptsize"
        nodes.append(_ppt_tikz_node_at(bounds, slide_width, slide_height, options, content))

    if not nodes:
        nodes.append("    \\node[anchor=center, font=\\small] at (current page.center) {暂无可提取内容};")
    out = [
        "\\begin{frame}[plain]",
        "  \\begin{tikzpicture}[remember picture, overlay]",
        *nodes,
        "  \\end{tikzpicture}",
        "\\end{frame}",
    ]
    return "\n".join(out), image_records


def _pptx_to_beamer_latex(file_bytes: bytes, filename: str) -> dict[str, object]:
    if Presentation is None:
        raise RuntimeError("服务器未安装 python-pptx，无法导入 PPT")
    prs = Presentation(BytesIO(file_bytes))
    package_id = uuid.uuid4().hex
    package_dir = UPLOAD_DIR / package_id
    package_dir.mkdir(parents=True, exist_ok=True)
    asset_urls: dict[str, str] = {}
    figure_assets: list[dict[str, str]] = []
    frames: list[str] = []
    slide_width = int(getattr(prs, "slide_width", 12192000) or 12192000)
    slide_height = int(getattr(prs, "slide_height", 6858000) or 6858000)
    for slide_index, slide in enumerate(prs.slides, 1):
        frame, images = _ppt_slide_to_latex_frame(slide, slide_index, package_dir, asset_urls, slide_width, slide_height)
        frames.append(frame)
        figure_assets.extend(images)
    title = Path(filename or "PPT").stem or "Imported PPT"
    latex = _bimsa_latex_preamble().rstrip()
    title_tex = _escape_latex_text_preserving_inline_math(title)
    latex = re.sub(
        r"\\title\[\]\{[^{}]*\}",
        lambda _match: f"\\title[]{{{title_tex}}}",
        latex,
        count=1,
    )
    latex += "\n\n\\begin{document}\n\n" + "\n\n".join(frames) + "\n\n\\end{document}\n"
    return {
        "success": True,
        "filename": filename,
        "latex": latex,
        "slide_count": len(frames),
        "asset_urls": asset_urls,
        "figure_assets": figure_assets,
        "package_id": package_id,
        "message": "PPT 已转换为 LaTeX",
    }


def _convert_legacy_ppt_to_pptx_bytes(file_bytes: bytes, filename: str) -> bytes:
    soffice = shutil.which("soffice") or shutil.which("libreoffice")
    if not soffice:
        raise RuntimeError("当前环境未安装 LibreOffice，无法直接转换 .ppt；请先另存为 .pptx 后导入。")
    with tempfile.TemporaryDirectory(prefix="kg-ppt-convert-") as temp_name:
        temp_dir = Path(temp_name)
        source_name = _safe_upload_filename(filename or "presentation.ppt", "presentation.ppt")
        source_path = temp_dir / source_name
        source_path.write_bytes(file_bytes)
        result = subprocess.run(
            [
                soffice,
                "--headless",
                "--convert-to",
                "pptx",
                "--outdir",
                str(temp_dir),
                str(source_path),
            ],
            stdout=subprocess.PIPE,
            stderr=subprocess.STDOUT,
            text=True,
            timeout=120,
        )
        pptx_files = list(temp_dir.glob("*.pptx"))
        if result.returncode != 0 or not pptx_files:
            tail = (result.stdout or "").strip()[-1200:]
            raise RuntimeError("PPT 转 PPTX 失败：" + (tail or "LibreOffice 未生成 pptx 文件"))
        return pptx_files[0].read_bytes()


def _english_frame_topic_title(value: str) -> str:
    text = str(value or "").strip()
    if not text:
        return "Core Topic"
    text = text.replace(r"\textbackslash{}", "\\")
    text = text.replace(r"\_", "_")
    text = re.sub(r"\\bar\s*\{\s*\\?(?:imath|i)\s*\}", "Average Selection Intensity", text, flags=re.IGNORECASE)
    text = re.sub(r"\bN\s*_\s*\{?\s*e\s*\}?", "Effective Population Size", text)
    text = re.sub(r"\$([^$]+)\$", r"\1", text)
    text = re.sub(r"\\\((.*?)\\\)", r"\1", text)
    text = re.sub(r"\\textbf\{([^{}]*)\}", r"\1", text)
    text = re.sub(r"\\[a-zA-Z]+\{([^{}]*)\}", r"\1", text)
    text = re.sub(r"\\[a-zA-Z]+", " ", text)
    text = text.replace("{", " ").replace("}", " ").replace("\\", " ")
    text = re.sub(r"[_^]", " ", text)
    text = re.sub(r"[\u4e00-\u9fff]+", " ", text)
    text = re.sub(r"\s+", " ", text).strip(" -:;,.")

    low = text.lower()
    if "tradeoff between" in low and "selection intensity" in low and "effective population size" in low:
        return "Tradeoff Between Selection Intensity and Effective Population Size"
    if not re.search(r"[A-Za-z]", text):
        return "Core Topic"
    return text[:90].strip(" -:;,.") or "Core Topic"


def _split_support_points(value: str) -> list[str]:
    raw = re.sub(r"\s+", " ", str(value or "")).strip()
    if not raw:
        return []
    parts = re.split(r"[。；;.!?]\s*|(?:\s+-\s+)", raw)
    seen: set[str] = set()
    points: list[str] = []
    for part in parts:
        clean = part.strip(" -,:：;；，。")
        if not clean or clean in seen:
            continue
        seen.add(clean)
        points.append(clean[:120])
    return points


def _topic_items_from_slide(slide: dict) -> list[str]:
    items: list[str] = []
    for item in slide.get("items") or []:
        clean = re.sub(r"\s+", " ", str(item or "")).strip()
        if clean:
            items.append(clean[:120])
    subtitle = re.sub(r"\s+", " ", str(slide.get("subtitle") or "")).strip()
    if subtitle:
        items.append(subtitle[:120])
    for part in _split_support_points(slide.get("notes") or ""):
        items.append(part[:120])
    seen: set[str] = set()
    result: list[str] = []
    for item in items:
        if item in seen:
            continue
        seen.add(item)
        result.append(item)
    return result


def _build_auto_expansion_frames(latex: str, minimum_slide_count: int, content: str = "") -> str:
    try:
        parsed = parse_latex_to_slides(latex or "")
    except Exception:
        logger.exception("Failed to parse LaTeX for automatic slide expansion")
        return latex

    slides = parsed.get("slides") if isinstance(parsed, dict) else []
    if not isinstance(slides, list) or len(slides) >= minimum_slide_count:
        return latex

    anchors = _extract_content_anchors(content)
    topics: list[dict[str, object]] = []
    for anchor in anchors:
        topics.append({"title": anchor, "items": [anchor]})
    for slide in slides:
        if not isinstance(slide, dict):
            continue
        slide_type = str(slide.get("type") or "").strip().lower()
        if slide_type in {"title", "toc"}:
            continue
        title = _english_frame_topic_title(slide.get("title") or "Core Topic")
        topic_items = _topic_items_from_slide(slide)
        topics.append({"title": title, "items": topic_items or [title]})

    if not topics:
        topics = [{
            "title": "Chapter Topic",
            "items": [
                "梳理本章核心概念与基本定义",
                "解释关键公式、变量含义和使用条件",
                "结合例子说明知识点之间的联系",
                "总结常见误区与课堂复习重点",
            ],
        }]

    templates = [
        ("Background and Motivation", [
            "说明 {topic} 的问题背景，以及它为什么是本章需要优先理解的核心主题。",
            "从研究目标、应用场景和课堂复习角度，梳理 {topic} 与前后知识点的连接。",
            "指出学习 {topic} 时最容易混淆的概念，并给出本页的辨析重点。",
            "结合知识图谱中的相关节点，建立后续公式、图表或例题讲解的上下文。",
        ]),
        ("Definition and Notation", [
            "给出 {topic} 中关键对象、变量和符号的中文解释，保留必要英文术语。",
            "说明每个变量的含义、单位或适用范围，避免只记公式而不理解条件。",
            "把知识图谱中的相关关系转化为定义之间的依赖链条。",
            "总结本页定义在后续推导、图表解读或课堂提问中的作用。",
        ]),
        ("Formula Explanation", [
            "围绕 {topic} 提取核心公式，并解释公式左端与右端分别代表的含义。",
            "逐项说明参数、变量和假设条件，强调哪些量可观测、哪些量由模型给出。",
            "讨论公式成立的前提，以及在什么情况下不能直接套用。",
            "用一句中文教学语言概括该公式传达的主要结论。",
        ]),
        ("Derivation Step", [
            "把 {topic} 的推导拆成前提、变换、结论三个层次，而不是直接给最终结果。",
            "说明每一步变换依赖的定义、近似或假设，帮助学生追踪逻辑来源。",
            "标出推导中容易跳步的位置，并解释该步骤为什么合理。",
            "总结该推导如何服务于本章的主线问题。",
        ]),
        ("Example Walkthrough", [
            "构造一个与 {topic} 相关的小例子，用具体变量展示概念如何落地。",
            "按输入条件、计算过程、结果解释的顺序组织讲解。",
            "指出例子中哪些结论可以推广，哪些只是当前条件下成立。",
            "给出课堂上可追问的问题，帮助检查学生是否真正理解。",
        ]),
        ("Common Mistakes", [
            "列出学习 {topic} 时常见误区，例如混淆变量、忽略条件或过度解释结果。",
            "用对比方式说明正确理解和错误理解之间的差别。",
            "解释这些误区为什么会影响公式使用、图表解读或结论判断。",
            "给出复习时可以自查的判断标准。",
        ]),
        ("Figure and Table Reading", [
            "说明 {topic} 相关图表应先看坐标、变量、趋势和异常点。",
            "把图表信息转化为中文结论，避免只描述视觉现象。",
            "联系知识图谱中的公式或定义，解释图表趋势背后的机制。",
            "总结该图表对本章核心问题提供了什么证据。",
        ]),
        ("Review and Concept Check", [
            "用问题形式回顾 {topic} 的关键定义、公式和适用条件。",
            "设计一个判断题或简短问答，检查学生是否能区分相近概念。",
            "要求学生说明结论背后的原因，而不是只复述术语。",
            "把本页复习点连接到下一页或下一节的学习目标。",
        ]),
    ]
    needed = minimum_slide_count - len(slides)
    new_frames: list[str] = []
    seen_body_keys = {_frame_body_key(frame) for frame in _extract_latex_frames(latex)}
    seen_titles = {_frame_title(frame).lower() for frame in _extract_latex_frames(latex) if _frame_title(frame)}

    idx = 0
    max_attempts = max(needed * 6, 24)
    while len(new_frames) < needed and idx < max_attempts:
        topic = topics[idx % len(topics)]
        section_title = BIMSA_FIXED_SECTIONS[idx % len(BIMSA_FIXED_SECTIONS)]
        topic_title = _english_frame_topic_title(str(topic.get("title") or section_title))
        template_name, template_bullets = templates[idx % len(templates)]
        frame_title_raw = f"{section_title} - {template_name}"
        if frame_title_raw.lower() in seen_titles:
            frame_title_raw = f"{topic_title} - {template_name} {idx + 1}"
        seen_titles.add(frame_title_raw.lower())
        frame_title = _escape_latex_text(frame_title_raw)
        topic_items = [str(item).strip() for item in list(topic.get("items") or []) if str(item).strip()]
        anchor = anchors[idx % len(anchors)] if anchors else (topic_items[idx % len(topic_items)] if topic_items else topic_title)

        bullets = [line.format(topic=topic_title) for line in template_bullets]
        bullets[0] = f"结合原始知识图谱要点“{anchor}”，说明它在固定章节“{section_title}”中的具体作用。"
        if topic_items and topic_items[0] != anchor:
            bullets[1] = f"补充关联要点“{topic_items[idx % len(topic_items)]}”，说明它与本页主题的依赖关系。"

        frame_lines = [f"\\begin{{frame}}{{{frame_title}}}", "  \\begin{itemize}"]
        for bullet in bullets:
            frame_lines.append(f"    \\item {_escape_latex_text(bullet)}")
        frame_lines.extend(["  \\end{itemize}", "\\end{frame}"])
        frame = "\n".join(frame_lines)
        body_key = _frame_body_key(frame)
        if body_key and body_key in seen_body_keys:
            idx += 1
            continue
        if _frame_content_score(frame) < 90:
            idx += 1
            continue
        seen_body_keys.add(body_key)
        new_frames.append(frame)
        idx += 1

    return _insert_frames_before_end_document(latex, new_frames)


def _strip_markdown_code_fence(text: str) -> str:
    cleaned = (text or "").strip()
    cleaned = re.sub(r"^```(?:latex|tex)?\s*", "", cleaned, flags=re.IGNORECASE)
    cleaned = re.sub(r"\s*```$", "", cleaned)
    return cleaned.strip()


async def _repair_latex_with_deepseek(
    *,
    client: DeepSeekClient,
    content: str,
    latex: str,
    target_count: int,
    current_count: int,
    model: str,
) -> str:
    missing = max(0, target_count - current_count)
    if missing <= 0:
        return latex
    batch_missing = min(missing, 12)
    frames = _extract_latex_frames(latex)
    existing_titles = [_frame_title(frame) or f"Slide {idx + 1}" for idx, frame in enumerate(frames)]
    title_list = "\n".join(f"- {title}" for title in existing_titles[:80])
    content_excerpt = content[:18000]
    anchor_list = "\n".join(f"- {anchor}" for anchor in _extract_content_anchors(content, limit=80))
    repair_system_prompt = (
        "你是 LaTeX Beamer 教学课件补页助手。只输出 LaTeX frame 代码，不输出解释、Markdown 代码块或完整文档。"
        "必须遵守：大标题/专业词汇英文，小标题和正文中文；不得重复已有页面标题或正文。"
    )
    repair_user_prompt = f"""
当前 Beamer 目标页数是 {target_count} 页，但已有 {current_count} 页，还缺 {missing} 页。

已有页面标题如下，禁止重复或只改后缀：
{title_list}

请基于下面知识图谱内容，本轮只补充 exactly {batch_missing} 个新的 \\begin{{frame}}...\\end{{frame}}。
要求：
- 只输出 frame 片段，不要输出 \\documentclass、\\begin{{document}} 或 \\end{{document}}。
- 每个 frame 标题必须不同，标题使用英文。
- 每页正文 3-6 个 bullet，正文解释使用中文，专业术语和变量保留英文。
- 每页必须承担不同教学功能：背景、定义、公式解释、推导、例子、图表解读、误区、概念检查、总结等。
- 禁止复制已有页正文；禁止出现正文要点完全相同但标题不同的页面。
- 优先覆盖知识图谱中尚未出现在已有标题中的节点、关系、公式、图表和例子。
- 每页正文必须引用至少一个下面“可用知识锚点”中的具体信息，不能只写“进一步展开”“复习重点”等空泛句子。
- 补充页必须服务于固定三章节之一：{BIMSA_FIXED_SECTIONS[0]}；{BIMSA_FIXED_SECTIONS[1]}；{BIMSA_FIXED_SECTIONS[2]}。
- 不要新增 section，不要输出 Test Title A2、Test title B2 或空白幻灯。

可用知识锚点：
{anchor_list}

知识图谱内容：
{content_excerpt}
""".strip()

    generated_parts: list[str] = []
    async for chunk in client.stream_generate(
        system_prompt=repair_system_prompt,
        user_prompt=repair_user_prompt,
        model=model,
        max_tokens=min(config.MAX_TOKENS, max(4096, batch_missing * 1100)),
        temperature=max(0.2, min(config.TEMPERATURE, 0.6)),
    ):
        generated_parts.append(chunk)
    generated = _strip_markdown_code_fence("".join(generated_parts))
    repair_frames = _extract_latex_frames(generated)
    if not repair_frames:
        return latex

    anchors = _extract_content_anchors(content)
    seen_titles = {_frame_title(frame).lower() for frame in frames if _frame_title(frame)}
    seen_bodies = {_frame_body_key(frame) for frame in frames if _frame_body_key(frame)}
    accepted: list[str] = []
    for frame in repair_frames:
        title = _frame_title(frame).lower()
        body_key = _frame_body_key(frame)
        if title and title in seen_titles:
            continue
        if body_key and body_key in seen_bodies:
            continue
        if _frame_content_score(frame) < 90:
            continue
        if not _frame_has_anchor(frame, anchors):
            logger.info("Rejecting unanchored repair frame: %s", title or "<untitled>")
            continue
        if title:
            seen_titles.add(title)
        if body_key:
            seen_bodies.add(body_key)
        accepted.append(frame)
        if len(accepted) >= batch_missing:
            break
    return _insert_frames_before_end_document(latex, accepted)


@router.post("/api/generate")
async def generate(payload: dict | None = Body(default=None)) -> StreamingResponse:
    payload = payload or {}
    content = _payload_text(payload, "content").strip()
    if not content:
        return _stream_error("请输入文案内容")

    style = _payload_text(payload, "style", "academic").strip() or "academic"
    custom_requirements = _payload_text(payload, "custom_requirements").strip()
    language = _payload_text(payload, "language", "title_terms_en_content_zh").strip() or "title_terms_en_content_zh"
    slide_count = max(1, min(_payload_int(payload, "slide_count", 7), 80))
    figure_assets = _payload_map(payload, "figure_assets")
    selected_sections = _payload_selected_sections(payload)
    outline_payload = payload.get("outline") if isinstance(payload.get("outline"), dict) else None
    outline = _validate_outline(outline_payload) if outline_payload else None
    selected_sections_requirement = _selected_sections_requirement(selected_sections)
    chapter_directory_title = _infer_chapter_directory_title(content)
    if selected_sections_requirement:
        custom_requirements = (
            custom_requirements + "\n" + selected_sections_requirement
            if custom_requirements
            else selected_sections_requirement
        )

    system_prompt = prompt_engine.build_system_prompt(
        style=style,
        custom_requirements=custom_requirements,
        slide_count=slide_count,
        language=language,
        figure_assets=figure_assets,
    )
    user_prompt = prompt_engine.build_user_prompt(
        content,
        custom_requirements,
        slide_count=slide_count,
        figure_assets=figure_assets,
        outline=outline,
    )
    api_key = _payload_text(payload, "api_key").strip() or _default_llm_api_key()
    if not api_key:
        return _stream_error("请输入本次生成使用的 GPT API Key")

    base_url = (
        _payload_text(payload, "base_url").strip()
        or _default_llm_base_url()
    )
    model = (
        _payload_text(payload, "model").strip()
        or _default_llm_model("pro")
    )
    logger.info("Generating Beamer: base_url=%s, model=%s, content_len=%s", base_url, model, len(content))
    client = DeepSeekClient(api_key=api_key, base_url=base_url)

    async def event_stream():
        queue: asyncio.Queue[dict[str, str]] = asyncio.Queue()

        async def send_event(event_type: str, content: str = "") -> None:
            await queue.put({"type": event_type, "content": content})

        async def worker() -> None:
            try:
                chunk_count = 0
                generated_parts: list[str] = []
                await send_event("heartbeat", f"已连接后端，正在请求模型 {model} 生成 LaTeX...")
                async for chunk_text in client.stream_generate(
                    system_prompt=system_prompt,
                    user_prompt=user_prompt,
                    model=model,
                    max_tokens=config.MAX_TOKENS,
                    temperature=config.TEMPERATURE,
                ):
                    chunk_count += 1
                    generated_parts.append(chunk_text)
                    if chunk_count % 8 == 0:
                        await send_event("heartbeat", f"模型正在生成 LaTeX，已收到 {chunk_count} 个片段...")
                await send_event("heartbeat", "模型已返回内容，正在检查页数和格式...")
                latex = _ensure_complete_bimsa_latex("".join(generated_parts))
                latex = _dedupe_latex_frames(latex)
                actual_slide_count = _count_generated_slides(latex)
                repair_attempts = 0
                while actual_slide_count < slide_count and repair_attempts < 4:
                    previous_count = actual_slide_count
                    repair_attempts += 1
                    await send_event("heartbeat", f"当前 {actual_slide_count}/{slide_count} 页，正在补齐第 {repair_attempts} 轮...")
                    logger.info(
                        "Generated %s/%s slides; requesting LLM repair pass %s",
                        actual_slide_count,
                        slide_count,
                        repair_attempts,
                    )
                    repaired = await _repair_latex_with_deepseek(
                        client=client,
                        content=content,
                        latex=latex,
                        target_count=slide_count,
                        current_count=actual_slide_count,
                        model=model,
                    )
                    latex = _dedupe_latex_frames(repaired)
                    actual_slide_count = _count_generated_slides(latex)
                    if actual_slide_count <= previous_count:
                        logger.info("LLM repair pass made no progress at %s slides", actual_slide_count)
                        break
                if actual_slide_count < slide_count:
                    await send_event("heartbeat", f"正在本地补足页面：{actual_slide_count}/{slide_count}...")
                    logger.info("LLM repair still produced %s/%s slides; using local dedup fallback", actual_slide_count, slide_count)
                    latex = _build_auto_expansion_frames(latex, slide_count, content=content)
                    latex = _dedupe_latex_frames(latex)
                    actual_slide_count = _count_generated_slides(latex)
                if actual_slide_count < slide_count:
                    logger.info("Local fallback still produced %s/%s slides; applying strict minimum pass", actual_slide_count, slide_count)
                    latex = _build_auto_expansion_frames(latex, slide_count, content=content)
                    actual_slide_count = _count_generated_slides(latex)
                await send_event("heartbeat", "正在应用图片、公式和目录格式约束...")
                latex = _prepare_generated_image_paths(latex, figure_assets)
                latex = _enforce_top_image_bottom_text_layout(latex, figure_assets)
                latex = _enforce_matched_image_captions(latex, content, figure_assets)
                latex = _merge_duplicate_image_frames(latex)
                latex = _enforce_formula_breakdown_layout(latex)
                latex = _enforce_formula_callout_anchor_positioning(latex)
                latex = _enforce_formula_callout_sequence_and_spacing(latex)
                latex = _enforce_unique_tikzmark_names(latex)
                latex = _fill_empty_outline_frames(latex, chapter_directory_title)
                latex = _fill_empty_content_frames(latex, content)
                latex = _enforce_formula_breakdown_layout(latex)
                latex = _enforce_formula_callout_anchor_positioning(latex)
                latex = _enforce_formula_callout_sequence_and_spacing(latex)
                latex = _enforce_unique_tikzmark_names(latex)
                latex = _remove_numbered_outline_item_labels(latex)
                latex = _ensure_safe_image_macros(latex)
                latex, missing_equations, resolved_equations = _apply_equation_reference_policy(latex)
                if missing_equations:
                    logger.info("Marked %s missing equation references", len(missing_equations))
                logger.info("Beamer generation done: %s chunks, %s slides", chunk_count, actual_slide_count)
                await send_event("chunk", latex)
                await send_event("done", "")
            except Exception as exc:
                logger.error("Beamer generation error: %s: %s", type(exc).__name__, exc)
                error_msg = str(exc)
                low = error_msg.lower()
                if "connect" in low or "timeout" in low:
                    error_msg = f"无法连接 API ({base_url})，请检查网络或 Base URL"
                elif "401" in error_msg or "auth" in low:
                    error_msg = "API Key 无效，请检查 Key 是否正确"
                elif "model" in low:
                    error_msg = f"模型 '{model}' 不可用，请检查模型名称"
                await send_event("error", error_msg)

        task = asyncio.create_task(worker())
        try:
            while True:
                try:
                    event = await asyncio.wait_for(queue.get(), timeout=20)
                except asyncio.TimeoutError:
                    event = {"type": "heartbeat", "content": "后端仍在生成 LaTeX，请继续等待..."}
                yield f"data: {json.dumps(event, ensure_ascii=False)}\n\n"
                if event.get("type") in {"done", "error"}:
                    break
        finally:
            if not task.done():
                task.cancel()

    return StreamingResponse(
        event_stream(),
        media_type="text/event-stream",
        headers={
            "Cache-Control": "no-cache",
            "Connection": "keep-alive",
            "X-Accel-Buffering": "no",
        },
    )


@router.post("/api/parse-slides")
async def parse_slides(req: ParseRequest) -> JSONResponse:
    try:
        latex = _sanitize_unsafe_title_assets(req.latex or "")
        latex = _enforce_top_image_bottom_text_layout(latex, {})
        latex = _enforce_formula_breakdown_layout(latex)
        latex = _enforce_formula_callout_anchor_positioning(latex)
        latex = _enforce_formula_callout_sequence_and_spacing(latex)
        latex = _enforce_unique_tikzmark_names(latex)
        latex = _fill_empty_outline_frames(latex)
        latex = _fill_empty_content_frames(latex)
        latex = _enforce_formula_breakdown_layout(latex)
        latex = _enforce_formula_callout_anchor_positioning(latex)
        latex = _enforce_formula_callout_sequence_and_spacing(latex)
        latex = _enforce_unique_tikzmark_names(latex)
        latex = _remove_numbered_outline_item_labels(latex)
        latex, missing_equations, resolved_equations = _apply_equation_reference_policy(latex)
        parsed = parse_latex_to_slides(latex)
        parsed = _attach_missing_equations_to_slides(parsed, latex, missing_equations)
        parsed["latex"] = latex
        parsed["missing_equations"] = missing_equations
        parsed["resolved_equations"] = resolved_equations
        return JSONResponse(content=parsed)
    except Exception as exc:
        logger.error("Beamer parse error: %s", exc)
        return JSONResponse(content={"error": str(exc)}, status_code=500)


@router.post("/api/render-latex-pages")
async def render_latex_pages(req: RenderLatexRequest) -> JSONResponse:
    try:
        latex = _sanitize_unsafe_title_assets(req.latex or "")
        latex = _enforce_top_image_bottom_text_layout(latex, {})
        latex = _enforce_formula_breakdown_layout(latex)
        latex = _enforce_formula_callout_anchor_positioning(latex)
        latex = _enforce_formula_callout_sequence_and_spacing(latex)
        latex = _enforce_unique_tikzmark_names(latex)
        latex = _fill_empty_outline_frames(latex)
        latex = _fill_empty_content_frames(latex)
        latex = _enforce_formula_breakdown_layout(latex)
        latex = _enforce_formula_callout_anchor_positioning(latex)
        latex = _enforce_formula_callout_sequence_and_spacing(latex)
        latex = _enforce_unique_tikzmark_names(latex)
        latex = _remove_numbered_outline_item_labels(latex)
        latex, missing_equations, resolved_equations = _apply_equation_reference_policy(latex)
        try:
            parsed = parse_latex_to_slides(latex)
            parsed = _attach_missing_equations_to_slides(parsed, latex, missing_equations)
        except Exception:
            parsed = {"title": Path(req.filename or "presentation.tex").stem, "slides": []}
        pdf_bytes = _compile_latex_to_pdf_bytes(latex, req.asset_urls)
        return JSONResponse(content={
            **parsed,
            "latex": latex,
            "missing_equations": missing_equations,
            "resolved_equations": resolved_equations,
            "rendered_pages": _render_pdf_bytes_to_pages(pdf_bytes),
            "render_source": "latex",
            "filename": req.filename,
        })
    except Exception as exc:
        logger.error("Beamer LaTeX render error: %s", exc)
        return JSONResponse(content={"error": str(exc)}, status_code=500)


@router.post("/api/overleaf-package")
@router.post("/api/overleaf-package/")
async def overleaf_package(req: OverleafPackageRequest) -> JSONResponse:
    try:
        zip_bytes, filename = _build_overleaf_package(req)
        snip_uri = "data:application/zip;base64," + base64.b64encode(zip_bytes).decode("ascii")
        return JSONResponse(content={
            "success": True,
            "filename": filename,
            "snip_uri": snip_uri,
            "size": len(zip_bytes),
        })
    except Exception as exc:
        logger.error("Overleaf package error: %s", exc)
        return JSONResponse(content={"success": False, "error": str(exc)}, status_code=500)


@router.post("/api/render-pdf-pages")
async def render_pdf_pages(file: UploadFile = File(...)) -> JSONResponse:
    try:
        if Path(file.filename or "").suffix.lower() != ".pdf":
            return JSONResponse(content={"error": "请选择 PDF 文件"}, status_code=400)
        content = await file.read()
        if not content:
            return JSONResponse(content={"error": "PDF 文件为空"}, status_code=400)
        return JSONResponse(content={
            "success": True,
            "filename": file.filename or "presentation.pdf",
            "rendered_pages": _render_pdf_bytes_to_pages(content),
            "render_source": "pdf",
        })
    except Exception as exc:
        logger.error("Beamer PDF render error: %s", exc)
        return JSONResponse(content={"error": str(exc)}, status_code=500)


def _select_uploaded_latex_tex(files: list[UploadFile]) -> UploadFile | None:
    candidates: list[tuple[bool, int, str, UploadFile]] = []
    for upload in files:
        rel_name = (upload.filename or "").replace("\\", "/")
        if Path(rel_name).suffix.lower() not in {".tex", ".latex"}:
            continue
        name = Path(rel_name).name.lower()
        candidates.append((name != "main.tex", rel_name.count("/"), rel_name.lower(), upload))
    candidates.sort(key=lambda item: item[:3])
    return candidates[0][3] if candidates else None


@router.post("/api/render-latex-project")
async def render_latex_project(files: List[UploadFile] = File(...)) -> JSONResponse:
    try:
        if not files:
            return JSONResponse(content={"error": "未收到 LaTeX 项目文件"}, status_code=400)

        tex_upload = _select_uploaded_latex_tex(files)
        if not tex_upload:
            return JSONResponse(content={"error": "项目目录中未找到 .tex 文件"}, status_code=400)

        file_names = [file.filename or "" for file in files]
        root_name = _detect_common_root(file_names)
        if root_name and Path(root_name).suffix:
            root_name = ""

        with tempfile.TemporaryDirectory(prefix="kg-latex-project-") as temp_name:
            temp_dir = Path(temp_name)
            tex_rel_name = ""
            asset_urls: dict[str, str] = {}
            package_id = uuid.uuid4().hex
            package_dir = UPLOAD_DIR / package_id
            package_dir.mkdir(parents=True, exist_ok=True)

            for upload in files:
                raw_name = upload.filename or ""
                rel_name = _normalize_package_rel_path(raw_name, root_name=root_name)
                if not rel_name:
                    continue
                rel_path = Path(rel_name)
                if any(part in ("..", "") for part in rel_path.parts):
                    continue

                data = await upload.read()
                target = temp_dir / rel_path
                target.parent.mkdir(parents=True, exist_ok=True)
                target.write_bytes(data)

                if upload is tex_upload:
                    tex_rel_name = rel_path.as_posix()

                if rel_path.suffix.lower() in OVERLEAF_ASSET_EXTENSIONS:
                    upload_target = package_dir / rel_path
                    upload_target.parent.mkdir(parents=True, exist_ok=True)
                    upload_target.write_bytes(data)
                    asset_urls[rel_path.as_posix()] = f"/beamer-generator/uploads/{package_id}/{rel_path.as_posix()}"

            if not tex_rel_name:
                tex_rel_name = _normalize_package_rel_path(tex_upload.filename or "", root_name=root_name)
            tex_path = temp_dir / tex_rel_name
            if not tex_path.exists():
                return JSONResponse(content={"error": "主 .tex 文件无法读取"}, status_code=400)

            latex = _decode_text_bytes(tex_path.read_bytes())
            try:
                parsed = parse_latex_to_slides(latex)
            except Exception:
                parsed = {"title": Path(tex_rel_name).stem, "slides": []}
            pdf_bytes = _compile_latex_file_to_pdf_bytes(tex_path)
            return JSONResponse(content={
                **parsed,
                "success": True,
                "filename": Path(tex_rel_name).name,
                "tex_filename": Path(tex_rel_name).name,
                "latex": latex,
                "asset_urls": asset_urls,
                "asset_base_url": f"/beamer-generator/uploads/{package_id}/",
                "rendered_pages": _render_pdf_bytes_to_pages(pdf_bytes),
                "render_source": "latex_project",
            })
    except Exception as exc:
        logger.error("Beamer LaTeX project render error: %s", exc)
        return JSONResponse(content={"error": str(exc)}, status_code=500)


def _decode_text_bytes(data: bytes) -> str:
    for encoding in ("utf-8-sig", "utf-8", "gb18030", "gbk"):
        try:
            return data.decode(encoding)
        except UnicodeDecodeError:
            continue
    return data.decode("utf-8", errors="ignore")


def _clean_markdown_source(text: str) -> str:
    lines = []
    in_front_matter = False
    for index, line in enumerate((text or "").splitlines()):
        if index == 0 and line.strip() == "---":
            in_front_matter = True
            continue
        if in_front_matter:
            if line.strip() == "---":
                in_front_matter = False
            continue
        lines.append(line.rstrip())
    cleaned = "\n".join(lines).strip()
    while "\n\n\n" in cleaned:
        cleaned = cleaned.replace("\n\n\n", "\n\n")
    return cleaned


def _normalize_package_rel_path(name: str, root_name: str = "") -> str:
    path = Path((name or "").replace("\\", "/"))
    parts = [part for part in path.parts if part not in ("", ".", "/")]
    if root_name and parts and parts[0] == root_name:
        parts = parts[1:]
    while parts and parts[0] in (".", ".."):
        parts = parts[1:]
    return "/".join(parts)


def _detect_common_root(file_names: list[str]) -> str:
    roots = []
    for name in file_names:
        parts = [part for part in Path(name.replace("\\", "/")).parts if part not in ("", ".", "/")]
        if parts:
            roots.append(parts[0])
    if not roots:
        return ""
    first = roots[0]
    if all(root == first for root in roots):
        return first
    return ""


OVERLEAF_ASSET_EXTENSIONS = {".png", ".jpg", ".jpeg", ".gif", ".webp", ".svg", ".bmp", ".pdf"}


def _select_overleaf_tex_member(infos: list[zipfile.ZipInfo], root_name: str = "") -> zipfile.ZipInfo | None:
    tex_infos = []
    for info in infos:
        rel_name = _normalize_package_rel_path(info.filename, root_name=root_name)
        if not rel_name or Path(rel_name).suffix.lower() not in {".tex", ".latex"}:
            continue
        name = Path(rel_name).name.lower()
        tex_infos.append((name != "main.tex", rel_name.count("/"), rel_name.lower(), info))
    tex_infos.sort(key=lambda item: item[:3])
    return tex_infos[0][3] if tex_infos else None


def _asset_url_for_graphic_target(target: str, asset_urls: dict[str, str]) -> str:
    raw = str(target or "").strip().strip('"').strip("'")
    if not raw or raw.startswith(("http://", "https://", "data:", "/")):
        return raw
    normalized = raw.replace("\\", "/").lstrip("./")
    normalized = normalized.split("?", 1)[0].split("#", 1)[0]
    key_map = {key.lower(): value for key, value in asset_urls.items()}
    direct = key_map.get(normalized.lower())
    if direct:
        return direct
    if Path(normalized).suffix:
        return raw
    for ext in OVERLEAF_ASSET_EXTENSIONS:
        match = key_map.get((normalized + ext).lower())
        if match:
            return match
    return raw


def _rewrite_latex_graphic_paths_to_upload_urls(latex: str, asset_urls: dict[str, str]) -> str:
    pattern = re.compile(r"(\\includegraphics\s*(?:\[[^\]]*\])?\s*)\{([^{}]+)\}", re.DOTALL)

    def replace(match: re.Match[str]) -> str:
        prefix = match.group(1)
        target = match.group(2)
        rewritten = _asset_url_for_graphic_target(target, asset_urls)
        return f"{prefix}{{{rewritten}}}"

    return pattern.sub(replace, latex or "")


def _import_overleaf_zip_bytes(data: bytes, filename: str) -> dict:
    package_id = uuid.uuid4().hex
    package_dir = UPLOAD_DIR / package_id
    package_dir.mkdir(parents=True, exist_ok=True)
    asset_urls: dict[str, str] = {}

    with zipfile.ZipFile(BytesIO(data)) as zf, tempfile.TemporaryDirectory(prefix="kg-overleaf-import-") as temp_name:
        infos = [
            info for info in zf.infolist()
            if not info.is_dir()
            and "__MACOSX/" not in info.filename.replace("\\", "/")
            and not Path(info.filename).name.startswith("._")
        ]
        root_name = _detect_common_root([info.filename for info in infos])
        if root_name and Path(root_name).suffix:
            root_name = ""
        tex_info = _select_overleaf_tex_member(infos, root_name=root_name)
        if not tex_info:
            raise ValueError("Overleaf ZIP 中未找到 .tex 文件")

        temp_dir = Path(temp_name)
        tex_rel_name = ""
        for info in infos:
            rel_name = _normalize_package_rel_path(info.filename, root_name=root_name)
            if not rel_name:
                continue
            rel_path = Path(rel_name)
            if any(part in ("..", "") for part in rel_path.parts):
                continue
            raw = zf.read(info)
            temp_target = temp_dir / rel_path
            temp_target.parent.mkdir(parents=True, exist_ok=True)
            temp_target.write_bytes(raw)

            if info.filename == tex_info.filename:
                tex_rel_name = rel_path.as_posix()

            if rel_path.suffix.lower() in OVERLEAF_ASSET_EXTENSIONS:
                upload_target = package_dir / rel_path
                upload_target.parent.mkdir(parents=True, exist_ok=True)
                upload_target.write_bytes(raw)
                asset_urls[rel_path.as_posix()] = f"/beamer-generator/uploads/{package_id}/{rel_path.as_posix()}"

        if not tex_rel_name:
            tex_rel_name = _normalize_package_rel_path(tex_info.filename, root_name=root_name)
        tex_path = temp_dir / tex_rel_name
        if not tex_path.exists():
            raise ValueError("Overleaf ZIP 中的主 .tex 文件无法读取")

        latex = _decode_text_bytes(tex_path.read_bytes())
        rewritten_latex = _rewrite_latex_graphic_paths_to_upload_urls(latex, asset_urls)
        parsed = parse_latex_to_slides(rewritten_latex)
        result = {
            **parsed,
            "success": True,
            "filename": filename,
            "tex_filename": Path(tex_rel_name).name,
            "latex": rewritten_latex,
            "asset_urls": asset_urls,
            "asset_base_url": f"/beamer-generator/uploads/{package_id}/",
            "render_source": "overleaf_zip",
        }
        try:
            pdf_bytes = _compile_latex_file_to_pdf_bytes(tex_path)
            result["rendered_pages"] = _render_pdf_bytes_to_pages(pdf_bytes)
        except Exception as exc:
            result["render_error"] = str(exc)
        return result


@router.post("/api/import-overleaf-package")
async def import_overleaf_package(file: UploadFile = File(...)) -> JSONResponse:
    try:
        filename = file.filename or "overleaf.zip"
        if Path(filename).suffix.lower() != ".zip":
            return JSONResponse(content={"error": "请选择从 Overleaf 下载的 .zip 源码包"}, status_code=400)
        content = await file.read()
        if not content:
            return JSONResponse(content={"error": "上传文件为空"}, status_code=400)
        return JSONResponse(content=_import_overleaf_zip_bytes(content, filename))
    except zipfile.BadZipFile:
        return JSONResponse(content={"error": "ZIP 文件无法读取，请确认文件未损坏"}, status_code=400)
    except Exception as exc:
        logger.error("Overleaf ZIP import error: %s", exc)
        return JSONResponse(content={"error": str(exc)}, status_code=500)


@router.post("/api/resolve-missing-equations")
async def resolve_missing_equations(payload: dict | None = Body(default=None)) -> JSONResponse:
    try:
        payload = payload or {}
        latex = _payload_text(payload, "latex")
        source = _payload_text(payload, "source")
        filename = _payload_text(payload, "filename") or "supplement"
        if not latex.strip():
            return JSONResponse(content={"error": "当前没有可补全的 LaTeX 内容"}, status_code=400)
        if not source.strip():
            return JSONResponse(content={"error": "补充章节内容为空"}, status_code=400)
        marked_latex, still_missing, resolved = _apply_equation_reference_policy(
            latex,
            source=source,
            source_id=filename,
            source_title=filename,
        )
        return JSONResponse(content={
            "success": True,
            "latex": marked_latex,
            "resolved": resolved,
            "missing": [item.get("label") or item.get("key") for item in still_missing],
            "missing_equations": still_missing,
        })
    except Exception as exc:
        logger.error("Equation resolve error: %s", exc)
        return JSONResponse(content={"error": str(exc)}, status_code=500)


def _rewrite_markdown_image_links(text: str, asset_urls: dict[str, str]) -> str:
    pattern = re.compile(r"(!\[[^\]]*?\]\()([^)]+)(\))")

    def replace(match: re.Match[str]) -> str:
        prefix, target, suffix = match.group(1), match.group(2), match.group(3)
        raw_target = target.strip().strip('"').strip("'")
        if raw_target.startswith(("http://", "https://", "data:", "/")):
            return match.group(0)
        normalized = raw_target.lstrip("./")
        normalized = normalized.replace("\\", "/")
        normalized = normalized.split("?", 1)[0].split("#", 1)[0]
        url = asset_urls.get(normalized)
        if not url:
            return match.group(0)
        return f"{prefix}{url}{suffix}"

    return pattern.sub(replace, text or "")


def _store_uploaded_package_files(
    files: list[UploadFile],
    package_id: str,
    require_markdown: bool = True,
) -> tuple[str, dict[str, str], str, str]:
    package_dir = UPLOAD_DIR / package_id
    package_dir.mkdir(parents=True, exist_ok=True)
    file_names = [file.filename or "" for file in files]
    root_name = _detect_common_root(file_names)
    asset_urls: dict[str, str] = {}
    markdown_bytes: bytes | None = None
    markdown_name = ""

    for upload in files:
        raw_name = upload.filename or ""
        rel_name = _normalize_package_rel_path(raw_name, root_name=root_name)
        if not rel_name:
            continue
        rel_path = Path(rel_name)
        if any(part in ("..", "") for part in rel_path.parts):
            continue

        target = package_dir / rel_path
        target.parent.mkdir(parents=True, exist_ok=True)
        data = upload.file.read()
        target.write_bytes(data)
        asset_urls[rel_path.as_posix()] = f"/beamer-generator/uploads/{package_id}/{rel_path.as_posix()}"

        if markdown_bytes is None and rel_path.suffix.lower() in {".md", ".markdown"}:
            markdown_bytes = data
            markdown_name = rel_path.name

    if markdown_bytes is None:
        if not require_markdown:
            return "", asset_urls, "", root_name
        raise HTTPException(status_code=400, detail="未找到 md / markdown 知识图谱文件")

    markdown_text = _clean_markdown_source(_decode_text_bytes(markdown_bytes))
    rewritten = _rewrite_markdown_image_links(markdown_text, asset_urls)
    return rewritten, asset_urls, markdown_name, root_name


def _extract_markdown_from_zip(data: bytes) -> tuple[str, list[str]]:
    contents: list[tuple[str, str]] = []
    with zipfile.ZipFile(BytesIO(data)) as zf:
        infos = [
            info for info in zf.infolist()
            if not info.is_dir()
            and not Path(info.filename).name.startswith("._")
            and Path(info.filename).suffix.lower() in {".md", ".markdown", ".txt"}
        ]
        infos.sort(key=lambda item: item.filename.lower())
        for info in infos:
            raw = zf.read(info)
            text = _clean_markdown_source(_decode_text_bytes(raw))
            if not text:
                continue
            title = Path(info.filename).stem
            contents.append((info.filename, f"# {title}\n\n{text}"))

    merged = "\n\n---\n\n".join(text for _name, text in contents)
    return merged, [name for name, _text in contents]


@router.post("/api/import-markdown-source")
async def import_markdown_source(file: UploadFile = File(...)):
    try:
        filename = file.filename or "source.md"
        ext = Path(filename).suffix.lower()
        content = await file.read()
        if not content:
            return JSONResponse(content={"error": "上传文件为空"}, status_code=400)

        if ext not in {".md", ".markdown"}:
            return JSONResponse(content={"error": "仅支持 .md/.markdown 知识图谱文件"}, status_code=400)

        if ext == ".zip":
            text, files = _extract_markdown_from_zip(content)
            if not text:
                return JSONResponse(content={"error": "压缩包中未找到 .md/.markdown/.txt 文件"}, status_code=400)
            return {
                "filename": filename,
                "files": files,
                "content": text,
                "char_count": len(text),
            }

        if ext not in {".md", ".markdown"}:
            return JSONResponse(content={"error": "仅支持 .md/.markdown/.txt/.zip 文件"}, status_code=400)

        text = _clean_markdown_source(_decode_text_bytes(content))
        return {
            "filename": filename,
            "files": [filename],
            "content": text,
            "char_count": len(text),
        }
    except zipfile.BadZipFile:
        return JSONResponse(content={"error": "ZIP 文件无法读取，请确认文件未损坏"}, status_code=400)
    except Exception as exc:
        logger.error("Beamer markdown import error: %s", exc)
        return JSONResponse(content={"error": str(exc)}, status_code=500)


@router.post("/api/import-knowledge-package")
async def import_knowledge_package(files: List[UploadFile] = File(...)):
    try:
        if not files:
            return JSONResponse(content={"error": "未收到知识图谱文件包"}, status_code=400)

        package_id = uuid.uuid4().hex
        markdown_text, asset_urls, markdown_name, root_name = _store_uploaded_package_files(files, package_id)

        title = ""
        for line in markdown_text.splitlines():
            stripped = line.strip()
            if stripped.startswith("#"):
                title = stripped.lstrip("#").strip()
                break
        if not title:
            title = Path(markdown_name or root_name or "knowledge-package").stem

        return {
            "success": True,
            "file_name": markdown_name or root_name,
            "graph_type": "knowledge_package",
            "chapter_hint": {
                "title": title,
                "content": markdown_text,
            },
            "markdown_content": markdown_text,
            "parsed": {
                "nodes": 0,
                "relations": 0,
                "files": len(files),
                "assets": max(0, len(files) - 1),
            },
            "result": {
                "package_id": package_id,
                "package_name": root_name or package_id,
                "markdown_file": markdown_name,
                "asset_base_url": f"/beamer-generator/uploads/{package_id}/",
                "asset_urls": asset_urls,
            },
            "message": "知识图谱文件包导入成功",
            "imported_at": datetime.now().isoformat(),
        }
    except HTTPException:
        raise
    except Exception as exc:
        logger.error("Beamer knowledge package import error: %s", exc)
        return JSONResponse(content={"error": str(exc)}, status_code=500)


@router.post("/api/import-image-package")
async def import_image_package(files: List[UploadFile] = File(...)):
    try:
        if not files:
            return JSONResponse(content={"error": "未收到图片包文件"}, status_code=400)

        image_exts = {".png", ".jpg", ".jpeg", ".gif", ".webp", ".svg", ".bmp"}
        image_files = [file for file in files if Path(file.filename or "").suffix.lower() in image_exts]
        if not image_files:
            return JSONResponse(content={"error": "图片包中未找到可用图片文件"}, status_code=400)

        package_id = uuid.uuid4().hex
        _, asset_urls, _, root_name = _store_uploaded_package_files(image_files, package_id, require_markdown=False)

        return {
            "success": True,
            "graph_type": "image_package",
            "package_id": package_id,
            "package_name": root_name or package_id,
            "parsed": {
                "files": len(image_files),
                "assets": len(asset_urls),
            },
            "result": {
                "package_id": package_id,
                "package_name": root_name or package_id,
                "asset_base_url": f"/beamer-generator/uploads/{package_id}/",
                "asset_urls": asset_urls,
            },
            "message": "图片包导入成功",
            "imported_at": datetime.now().isoformat(),
        }
    except HTTPException:
        raise
    except Exception as exc:
        logger.error("Beamer image package import error: %s", exc)
        return JSONResponse(content={"error": str(exc)}, status_code=500)


@router.post("/api/import-ppt-latex")
async def import_ppt_latex(file: UploadFile = File(...)):
    try:
        filename = file.filename or "presentation.pptx"
        ext = Path(filename).suffix.lower()
        if ext != ".pptx":
            return JSONResponse(content={"error": "仅支持 .pptx 文件"}, status_code=400)
        content = await file.read()
        if not content:
            return JSONResponse(content={"error": "上传文件为空"}, status_code=400)
        if len(content) > 80 * 1024 * 1024:
            return JSONResponse(content={"error": "PPT 文件超过 80MB 限制"}, status_code=400)
        return _pptx_to_beamer_latex(content, filename)
    except Exception as exc:
        logger.error("Beamer PPT to LaTeX import error: %s", exc)
        return JSONResponse(content={"error": str(exc)}, status_code=500)


@router.options("/api/export-pptx")
@router.options("/api/export-pptx/")
async def export_pptx_options() -> Response:
    return Response(status_code=204)


@router.post("/api/export-pptx")
@router.post("/api/export-pptx/")
async def export_pptx(req: ExportRequest) -> Response:
    try:
        pptx_bytes = generate_pptx(req.model_dump(), upload_dir=str(UPLOAD_DIR))
        return Response(
            content=pptx_bytes,
            media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            headers={"Content-Disposition": "attachment; filename=presentation.pptx"},
        )
    except Exception as exc:
        logger.error("Beamer export error: %s", exc)
        return JSONResponse(content={"error": str(exc)}, status_code=500)


@router.options("/api/save-project")
@router.options("/api/save-project/")
async def save_project_options() -> Response:
    return Response(status_code=204)


@router.post("/api/save-project")
@router.post("/api/save-project/")
async def save_project(req: SaveProjectRequest) -> Response:
    try:
        data = req.model_dump()
        latex = data.pop("latex", "") or ""
        chapter_id = _safe_saved_project_id(data.pop("chapter_id", "") or data.get("title") or "presentation")
        chapter_title = data.pop("chapter_title", "") or data.get("title") or chapter_id
        if latex.strip():
            latex, missing_equations, resolved_equations = _apply_equation_reference_policy(
                latex,
                source_id=chapter_id,
                source_title=chapter_title,
            )
            _update_equation_index_from_source(latex, chapter_id, chapter_title)
            data["missing_equations"] = missing_equations
            data["resolved_equations"] = resolved_equations

        saved_payload = {
            **data,
            "chapter_id": chapter_id,
            "chapter_title": chapter_title,
            "latex": latex,
            "updated_at": datetime.now().isoformat(),
        }
        _saved_project_path(chapter_id).write_text(
            json.dumps(saved_payload, ensure_ascii=False, indent=2),
            encoding="utf-8",
        )

        return JSONResponse(content={
            "success": True,
            "chapter_id": chapter_id,
            "chapter_title": chapter_title,
            "title": saved_payload.get("title") or chapter_title,
            "slide_count": len(saved_payload.get("slides") or []),
            "slides": [
                {
                    "page_index": idx,
                    "title": (slide or {}).get("title") or f"Slide {idx + 1}",
                    "type": (slide or {}).get("type") or "content",
                }
                for idx, slide in enumerate(saved_payload.get("slides") or [])
            ],
            "missing_equations": saved_payload.get("missing_equations") or [],
            "missing_equation_count": len(saved_payload.get("missing_equations") or []),
            "resolved_equations": saved_payload.get("resolved_equations") or [],
            "updated_at": saved_payload["updated_at"],
        })
    except Exception as exc:
        logger.error("Beamer save project error: %s", exc)
        return JSONResponse(content={"error": str(exc)}, status_code=500)


@router.get("/api/saved-projects")
async def list_saved_projects() -> JSONResponse:
    projects = []
    for path in sorted(SAVED_PROJECT_DIR.glob("*.json"), key=lambda p: p.stat().st_mtime, reverse=True):
        summary = _saved_project_summary(path)
        if summary:
            projects.append(summary)
    return JSONResponse(content={"projects": projects})


@router.get("/api/saved-projects/{chapter_id}")
async def get_saved_project(chapter_id: str) -> JSONResponse:
    path = _saved_project_path(chapter_id)
    if not path.exists():
        raise HTTPException(status_code=404, detail="Saved PPT project not found")
    return JSONResponse(content=json.loads(path.read_text(encoding="utf-8")))


@router.options("/api/saved-projects/{chapter_id}")
async def delete_saved_project_options(chapter_id: str) -> Response:
    return Response(status_code=204)


@router.delete("/api/saved-projects/{chapter_id}")
async def delete_saved_project(chapter_id: str) -> JSONResponse:
    path = _saved_project_path(chapter_id)
    if not path.exists():
        raise HTTPException(status_code=404, detail="Saved PPT project not found")
    path.unlink()
    return JSONResponse(content={
        "success": True,
        "chapter_id": _safe_saved_project_id(chapter_id),
    })


@router.get("/api/saved-projects/{chapter_id}/slides/{page_index}")
async def get_saved_project_slide(chapter_id: str, page_index: int) -> JSONResponse:
    path = _saved_project_path(chapter_id)
    if not path.exists():
        raise HTTPException(status_code=404, detail="Saved PPT project not found")
    data = json.loads(path.read_text(encoding="utf-8"))
    slides = data.get("slides") if isinstance(data.get("slides"), list) else []
    if page_index < 0 or page_index >= len(slides):
        raise HTTPException(status_code=404, detail="Saved PPT slide not found")
    return JSONResponse(content={
        "chapter_id": data.get("chapter_id") or chapter_id,
        "chapter_title": data.get("chapter_title") or data.get("title") or chapter_id,
        "page_index": page_index,
        "slide": slides[page_index],
        "latex": data.get("latex") or "",
    })


@router.post("/api/upload-image")
async def upload_image(file: UploadFile = File(...)):
    try:
        ext = Path(file.filename or "img.png").suffix or ".png"
        filename = f"{uuid.uuid4().hex}{ext}"
        filepath = UPLOAD_DIR / filename
        content = await file.read()
        filepath.write_bytes(content)
        return {
            "url": f"/beamer-generator/uploads/{filename}",
            "filename": filename,
            "size": len(content),
        }
    except Exception as exc:
        logger.error("Beamer image upload error: %s", exc)
        return JSONResponse(content={"error": str(exc)}, status_code=500)

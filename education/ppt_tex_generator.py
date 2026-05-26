"""Generate PPT/TeX teaching artifacts from graph-backed lesson content."""

from __future__ import annotations

import hashlib
import json
import re
from datetime import datetime
from pathlib import Path
from typing import Any, Dict, Iterable, List, Optional

from pptx import Presentation
from pptx.util import Inches, Pt

from KGTS.core.bridge import RUNTIME_DIR


ARTIFACT_DIR = RUNTIME_DIR / "artifacts" / "ppt_tex"


def normalize_generated_slides(raw: Any, *, fallback_title: str, fallback_content: str, max_slides: int = 12) -> List[Dict[str, Any]]:
    slides = _coerce_slides(raw)
    normalized: List[Dict[str, Any]] = []
    for index, slide in enumerate(slides[: max(1, max_slides)], start=1):
        if not isinstance(slide, dict):
            continue
        title = _clean_text(slide.get("title") or slide.get("heading") or f"{fallback_title} {index}")
        bullets = _normalize_bullets(slide.get("bullets") or slide.get("points") or slide.get("body"))
        notes = _clean_text(slide.get("notes") or slide.get("speaker_notes") or "")
        content = _clean_text(slide.get("content") or "\n".join(f"- {item}" for item in bullets))
        if not bullets and content:
            bullets = _content_to_bullets(content)
        if not content:
            content = "\n".join(f"- {item}" for item in bullets)
        if not title and not content:
            continue
        normalized.append(
            {
                "index": len(normalized) + 1,
                "title": title or f"第 {len(normalized) + 1} 页",
                "content": content,
                "notes": notes,
                "has_images": False,
                "image_count": 0,
                "images": [],
                "tables": [],
                "body_texts": bullets,
                "raw_text": _build_raw_text(title, bullets, notes),
            }
        )
    if normalized:
        return normalized
    return fallback_slides_from_context(fallback_title, fallback_content, max_slides=max_slides)


def fallback_slides_from_context(title: str, content: str, *, max_slides: int = 12) -> List[Dict[str, Any]]:
    clean_title = _clean_text(title) or "图谱生成课件"
    sections = _split_context_sections(content)
    if not sections:
        sections = [{"title": clean_title, "content": content}]
    slides: List[Dict[str, Any]] = [
        _slide_payload(
            1,
            clean_title,
            ["课程目标", "关键知识结构", "课堂讲解顺序"],
            "导入本节课主题，说明本课会按图谱结构展开。",
        )
    ]
    for section in sections[: max(1, max_slides - 1)]:
        bullets = _content_to_bullets(section.get("content", ""))
        if not bullets:
            bullets = [section.get("title", "核心知识点")]
        slides.append(
            _slide_payload(
                len(slides) + 1,
                section.get("title") or f"知识点 {len(slides)}",
                bullets[:5],
                "围绕本页要点解释概念、关系和必要公式。",
            )
        )
    return slides[:max_slides]


def build_tex_from_slides(title: str, slides: Iterable[Dict[str, Any]], *, style_reference: Optional[Dict[str, Any]] = None) -> str:
    style_profile = _style_profile(style_reference)
    lines = _tex_preamble(title, style_profile)
    for slide in slides:
        slide_title = _tex_escape(str(slide.get("title") or f"第 {slide.get('index', '')} 页"))
        lines.append(r"\begin{frame}{" + slide_title + "}")
        bullets = slide.get("body_texts") or _content_to_bullets(str(slide.get("content") or ""))
        if bullets:
            lines.append(r"\begin{itemize}")
            for bullet in bullets[:8]:
                lines.append(r"\item " + _tex_escape(str(bullet)))
            lines.append(r"\end{itemize}")
        else:
            lines.append(_tex_escape(str(slide.get("content") or "")))
        notes = str(slide.get("notes") or "").strip()
        if notes:
            lines.append(r"\vspace{0.5em}")
            lines.append(r"\footnotesize{" + _tex_escape(notes[:420]) + "}")
        lines.append(r"\end{frame}")
    lines.append(r"\end{document}")
    return "\n".join(lines)


def build_pptx_artifact(
    title: str,
    slides: List[Dict[str, Any]],
    *,
    source_node_ids: Optional[List[str]] = None,
    style_reference: Optional[Dict[str, Any]] = None,
) -> Dict[str, Any]:
    ARTIFACT_DIR.mkdir(parents=True, exist_ok=True)
    slug = hashlib.md5(f"{title}:{datetime.now().isoformat()}".encode("utf-8")).hexdigest()[:12]
    pptx_path = ARTIFACT_DIR / f"generated_{slug}.pptx"
    tex_path = ARTIFACT_DIR / f"generated_{slug}.tex"
    tex_content = build_tex_from_slides(title, slides, style_reference=style_reference)
    tex_path.write_text(tex_content, encoding="utf-8")

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    for slide in slides:
        layout = prs.slide_layouts[1]
        ppt_slide = prs.slides.add_slide(layout)
        ppt_slide.shapes.title.text = str(slide.get("title") or f"第 {slide.get('index', '')} 页")
        body = ppt_slide.placeholders[1]
        frame = body.text_frame
        frame.clear()
        bullets = slide.get("body_texts") or _content_to_bullets(str(slide.get("content") or ""))
        for bullet_index, bullet in enumerate(bullets[:7]):
            paragraph = frame.paragraphs[0] if bullet_index == 0 else frame.add_paragraph()
            paragraph.text = str(bullet)
            paragraph.level = 0
            paragraph.font.size = Pt(22)
    prs.save(pptx_path)

    return {
        "kind": "generated_ppt_tex",
        "pptx_path": str(pptx_path),
        "tex_path": str(tex_path),
        "pptx_url": f"/api/education/artifacts/{pptx_path.name}",
        "tex_url": f"/api/education/artifacts/{tex_path.name}",
        "tex_content": tex_content,
        "slide_count": len(slides),
        "source_node_ids": source_node_ids or [],
        "generated_at": datetime.now().isoformat(),
    }


def _style_profile(style_reference: Optional[Dict[str, Any]]) -> Dict[str, Any]:
    if not isinstance(style_reference, dict):
        return {}
    profile = style_reference.get("profile") if isinstance(style_reference.get("profile"), dict) else style_reference
    return profile if isinstance(profile, dict) else {}


def _tex_preamble(title: str, style_profile: Dict[str, Any]) -> List[str]:
    document_class = str(style_profile.get("document_class") or "beamer").strip()
    document_options = [
        str(item).strip()
        for item in (style_profile.get("document_options") or [])
        if str(item).strip()
    ]
    if "aspectratio=169" not in document_options:
        document_options.append("aspectratio=169")
    if document_class not in {"beamer", "ctexbeamer"}:
        document_class = "beamer"

    lines = [
        r"\documentclass[" + ",".join(document_options[:6]) + "]{" + document_class + "}",
    ]
    theme = _theme_value(style_profile, "usetheme") or "Madrid"
    lines.append(r"\usetheme{" + _tex_escape(theme) + "}")
    lines.extend(
        [
            r"\usepackage{amsmath,amssymb}",
            r"\usepackage{graphicx}",
            r"\usepackage{tikz}",
            r"\usetikzlibrary{positioning}",
        ]
    )
    if document_class != "ctexbeamer":
        lines.append(r"\usepackage{ctex}")

    color_names = set()
    for color in style_profile.get("colors") or []:
        if not isinstance(color, dict):
            continue
        name = _safe_tex_name(str(color.get("name") or ""))
        model = str(color.get("model") or "").strip()
        value = str(color.get("value") or "").strip()
        if not name or not re.match(r"^[A-Za-z]+$", model) or not re.match(r"^[0-9A-Za-z, .!+-]+$", value):
            continue
        lines.append(r"\definecolor{" + name + "}{" + model + "}{" + value + "}")
        color_names.add(name)
    if "myline" in color_names:
        line_color = "myline"
    else:
        lines.append(r"\definecolor{kgtsline}{RGB}{0,116,112}")
        line_color = "kgtsline"

    signals = {str(item) for item in (style_profile.get("style_signals") or [])}
    templates = {str(item) for item in (style_profile.get("beamertemplates") or [])}
    if "thin top rule under frame titles" in signals or "frametitle" in templates:
        lines.append(r"\setbeamercolor{whitebg}{bg=white,fg=black}")
        lines.extend(
            [
                r"\setbeamertemplate{frametitle}{%",
                r"  \vspace*{0.2cm}%",
                r"  \begin{beamercolorbox}[wd=\paperwidth,leftskip=0.5cm,rightskip=0.5cm,ht=0.3cm,dp=0pt]{whitebg}%",
                r"    \usebeamerfont{frametitle}\textcolor{black}{\insertframetitle}%",
                r"  \end{beamercolorbox}%",
                r"  \begin{tikzpicture}[remember picture,overlay]",
                rf"    \draw[{line_color},line width=1.5pt] ([yshift=-1.3cm] current page.north west) -- ([yshift=-1.3cm] current page.north east);",
                r"  \end{tikzpicture}%",
                r"  \vspace{0.1cm}%",
                r"}",
            ]
        )
    if "white title blocks with black text" in signals or "title page" in templates:
        lines.extend(
            [
                r"\setbeamertemplate{title page}{%",
                r"  \begin{tikzpicture}[remember picture,overlay]",
                rf"    \draw[line width=1.5pt,color={line_color}] ([yshift=-40pt] current page.north west) -- ([yshift=-40pt] current page.north east);",
                r"  \end{tikzpicture}%",
                r"  \vspace*{36pt}",
                r"  \begin{center}",
                r"    \begin{tikzpicture}",
                r"      \node[draw=none,inner sep=8pt,fill=white,text=black,align=center,font=\Huge\bfseries] (titlebox) {\inserttitle};",
                r"      \node[draw=none,rounded corners=2pt,inner sep=8pt,fill=white,text=black,align=center,font=\large,below=5pt of titlebox] {\insertdate};",
                r"    \end{tikzpicture}",
                r"  \end{center}",
                r"}",
            ]
        )

    lines.extend(
        [
            r"\title{" + _tex_escape(title or "图谱生成课件") + "}",
            r"\date{}",
            r"\begin{document}",
            r"\frame{\titlepage}",
        ]
    )
    return lines


def _theme_value(style_profile: Dict[str, Any], command: str) -> str:
    prefix = f"{command}:"
    for theme in style_profile.get("themes") or []:
        text = str(theme or "")
        if text.startswith(prefix):
            return text[len(prefix):].strip()
    return ""


def _safe_tex_name(value: str) -> str:
    text = re.sub(r"[^A-Za-z0-9@]+", "", value)
    return text[:40]


def _coerce_slides(raw: Any) -> List[Any]:
    if isinstance(raw, list):
        return raw
    if isinstance(raw, dict):
        slides = raw.get("slides") or raw.get("pages")
        return slides if isinstance(slides, list) else []
    if not isinstance(raw, str):
        return []
    text = raw.strip()
    candidates = [text]
    fenced = re.search(r"```(?:json)?\s*([\s\S]*?)```", text, flags=re.I)
    if fenced:
        candidates.insert(0, fenced.group(1).strip())
    object_match = re.search(r"(\{[\s\S]*\})", text)
    array_match = re.search(r"(\[[\s\S]*\])", text)
    if object_match:
        candidates.append(object_match.group(1))
    if array_match:
        candidates.append(array_match.group(1))
    for candidate in candidates:
        try:
            parsed = json.loads(candidate)
        except json.JSONDecodeError:
            continue
        return _coerce_slides(parsed)
    return []


def _split_context_sections(content: str) -> List[Dict[str, str]]:
    sections: List[Dict[str, str]] = []
    current_title = ""
    current_lines: List[str] = []
    for raw_line in (content or "").splitlines():
        line = raw_line.strip()
        if not line:
            continue
        heading = re.match(r"^#{2,4}\s+(.+)$", line)
        typed_heading = re.match(r"^###\s+\[[^\]]+\]\s+(.+)$", line)
        if heading or typed_heading:
            if current_title or current_lines:
                sections.append({"title": current_title, "content": "\n".join(current_lines)})
            current_title = (typed_heading or heading).group(1).strip()
            current_lines = []
        elif current_title:
            current_lines.append(line)
    if current_title or current_lines:
        sections.append({"title": current_title, "content": "\n".join(current_lines)})
    return [item for item in sections if item.get("title") or item.get("content")]


def _content_to_bullets(content: str) -> List[str]:
    text = re.sub(r"\s+", " ", content or "").strip()
    if not text:
        return []
    sentences = re.split(r"(?<=[。！？.!?])\s+|[;；]\s*", text)
    bullets = []
    for sentence in sentences:
        clean = _clean_text(sentence)
        if len(clean) > 120:
            clean = clean[:117].rstrip() + "..."
        if clean and clean not in bullets:
            bullets.append(clean)
        if len(bullets) >= 5:
            break
    if bullets:
        return bullets
    return [text[:120]]


def _normalize_bullets(value: Any) -> List[str]:
    if isinstance(value, list):
        return [_clean_text(item) for item in value if _clean_text(item)]
    if isinstance(value, str):
        lines = [re.sub(r"^\s*[-*•\d.、)]+\s*", "", line).strip() for line in value.splitlines()]
        return [line for line in lines if line][:8] or _content_to_bullets(value)
    return []


def _slide_payload(index: int, title: str, bullets: List[str], notes: str = "") -> Dict[str, Any]:
    return {
        "index": index,
        "title": _clean_text(title) or f"第 {index} 页",
        "content": "\n".join(f"- {item}" for item in bullets),
        "notes": _clean_text(notes),
        "has_images": False,
        "image_count": 0,
        "images": [],
        "tables": [],
        "body_texts": bullets,
        "raw_text": _build_raw_text(title, bullets, notes),
    }


def _build_raw_text(title: str, bullets: List[str], notes: str) -> str:
    parts = [f"## {_clean_text(title)}"]
    parts.extend(f"- {_clean_text(item)}" for item in bullets if _clean_text(item))
    if notes:
        parts.append(f"[备注] {_clean_text(notes)}")
    return "\n".join(part for part in parts if part.strip())


def _clean_text(value: Any) -> str:
    return re.sub(r"\s+", " ", str(value or "")).strip()


def _tex_escape(value: str) -> str:
    protected: List[str] = []

    def protect(match: re.Match[str]) -> str:
        protected.append(match.group(0))
        return f"@@MATH{len(protected) - 1}@@"

    text = re.sub(r"\$\$[\s\S]*?\$\$|\$[^$\n]+\$", protect, value or "")
    replacements = {
        "\\": r"\textbackslash{}",
        "&": r"\&",
        "%": r"\%",
        "$": r"\$",
        "#": r"\#",
        "_": r"\_",
        "{": r"\{",
        "}": r"\}",
    }
    escaped = "".join(replacements.get(char, char) for char in text)
    for index, math in enumerate(protected):
        escaped = escaped.replace(f"@@MATH{index}@@", math)
    return escaped
